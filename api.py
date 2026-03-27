"""
4uPDF API - FastAPI backend
PDF tools: merge, split, compress, convert, OCR
User authentication, subscriptions, and usage tracking
"""

import asyncio
import fitz
import numpy as np
import re
import os
import json
import time
import uuid
import threading
import zipfile
import io
import shutil
import atexit
import hashlib
import sqlite3
import secrets
from pathlib import Path
from typing import List, Optional, Dict, Any
from datetime import datetime, timedelta
from contextlib import contextmanager
from concurrent.futures import ThreadPoolExecutor, as_completed
import tempfile

from fastapi import FastAPI, UploadFile, File, Form, HTTPException, BackgroundTasks, Depends, Request, Header
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from rapidocr_onnxruntime import RapidOCR

import jwt
import bcrypt

from rate_limiter import check_rate_limit, get_rate_limit_headers, rate_limiter
from file_validator import validate_upload, sanitize_filename, ALLOWED_MIME_TYPES, MAX_FILE_SIZE

try:
    import stripe
    STRIPE_AVAILABLE = True
except ImportError:
    STRIPE_AVAILABLE = False
    stripe = None

# Find a Unicode-capable TrueType font for PDF text rendering (Romanian diacritics etc.)
UNICODE_FONT_PATH = None
for _fp in [
    "C:/Windows/Fonts/arial.ttf",
    "C:/Windows/Fonts/times.ttf",
    "C:/Windows/Fonts/calibri.ttf",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
    "/usr/share/fonts/TTF/DejaVuSans.ttf",
]:
    if os.path.exists(_fp):
        UNICODE_FONT_PATH = _fp
        break

app = FastAPI(title="4uPDF API")

allowed_origins = os.environ.get("ALLOWED_ORIGINS", "http://localhost:3098,https://4updf.com").split(",")
app.add_middleware(
    CORSMiddleware,
    allow_origins=allowed_origins,
    allow_methods=["GET", "POST", "PUT", "DELETE"],
    allow_headers=["*"],
    allow_credentials=True,
    max_age=3600
)

# ============================================================================
# Security & Auth Configuration
# ============================================================================

security = HTTPBearer(auto_error=False)
# bcrypt is used directly for password hashing (passlib incompatible with Python 3.14)

# JWT Configuration - These can be overridden via settings API
JWT_SECRET_KEY = os.environ.get("JWT_SECRET_KEY", secrets.token_hex(32))
JWT_ALGORITHM = "HS256"
JWT_EXPIRY_HOURS = 24 * 7  # 7 days

# SuperAdmin key from environment (deprecated - kept for backward compat only)
SUPER_ADMIN_KEY = os.environ.get("SUPER_ADMIN_KEY", "")

# Import and include superadmin auth router (email/password based)
from routes.superadmin import router as superadmin_router, init_superadmin_routes
app.include_router(superadmin_router)

# In-memory heartbeat tracking for real-time active users
_active_users: Dict[str, float] = {}  # ip -> last_heartbeat_timestamp

# ============================================================================
# Database Configuration
# ============================================================================

DATA_DIR = Path("data")
DATA_DIR.mkdir(exist_ok=True)
DB_PATH = DATA_DIR / "4updf.db"

def get_db():
    """Get database connection."""
    conn = sqlite3.connect(str(DB_PATH), check_same_thread=False)
    conn.row_factory = sqlite3.Row
    return conn

@contextmanager
def db_session():
    """Context manager for database sessions."""
    conn = get_db()
    try:
        yield conn
        conn.commit()
    except Exception:
        conn.rollback()
        raise
    finally:
        conn.close()

def init_database():
    """Initialize database tables."""
    with db_session() as conn:
        cursor = conn.cursor()

        # Settings table for admin configuration
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS settings (
                key TEXT PRIMARY KEY,
                value TEXT NOT NULL,
                updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)

        # Users table
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS users (
                id TEXT PRIMARY KEY,
                email TEXT UNIQUE NOT NULL,
                password_hash TEXT NOT NULL,
                plan TEXT DEFAULT 'free',
                stripe_customer_id TEXT,
                stripe_subscription_id TEXT,
                subscription_status TEXT DEFAULT 'inactive',
                subscription_end_date TIMESTAMP,
                base_plan TEXT DEFAULT 'free',
                custom_max_file_size_mb INTEGER,
                custom_pages_per_day INTEGER,
                custom_features TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)

        # Anonymous usage tracking
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS anonymous_usage (
                id TEXT PRIMARY KEY,
                ip_address TEXT NOT NULL,
                fingerprint TEXT,
                local_token TEXT,
                composite_hash TEXT UNIQUE,
                pages_used_today INTEGER DEFAULT 0,
                last_reset_date TEXT,
                total_pages_all_time INTEGER DEFAULT 0,
                first_seen TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                last_seen TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)

        # Vouchers table
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS vouchers (
                id TEXT PRIMARY KEY,
                code TEXT UNIQUE NOT NULL,
                target_plan TEXT NOT NULL,
                duration_days INTEGER NOT NULL,
                max_uses INTEGER DEFAULT 1,
                current_uses INTEGER DEFAULT 0,
                expires_at TIMESTAMP,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                is_active INTEGER DEFAULT 1
            )
        """)

        # Voucher redemptions tracking
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS voucher_redemptions (
                id TEXT PRIMARY KEY,
                voucher_id TEXT NOT NULL,
                user_id TEXT NOT NULL,
                redeemed_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                expires_at TIMESTAMP NOT NULL,
                original_plan TEXT,
                FOREIGN KEY (voucher_id) REFERENCES vouchers(id),
                FOREIGN KEY (user_id) REFERENCES users(id)
            )
        """)

        # Usage history for authenticated users
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS usage_history (
                id TEXT PRIMARY KEY,
                user_id TEXT,
                anonymous_id TEXT,
                operation TEXT NOT NULL,
                pages_processed INTEGER DEFAULT 0,
                file_size_bytes INTEGER DEFAULT 0,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)

        # Analytics events table for local analytics tracking
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS analytics_events (
                id TEXT PRIMARY KEY,
                event_type TEXT NOT NULL,
                page TEXT,
                user_id TEXT,
                ip TEXT,
                user_agent TEXT,
                referrer TEXT,
                country TEXT,
                timestamp TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)

        # Heartbeats table for real-time active user tracking
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS heartbeats (
                id TEXT PRIMARY KEY,
                user_id TEXT,
                ip TEXT NOT NULL,
                user_agent TEXT,
                current_page TEXT,
                tool_name TEXT,
                session_id TEXT,
                last_seen TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                first_seen TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)

        # Add last_active, is_banned, and role columns to users (safe ALTER)
        for col, col_def in [
            ("last_active", "TIMESTAMP"),
            ("is_banned", "INTEGER DEFAULT 0"),
            ("role", "TEXT DEFAULT 'user'"),
        ]:
            try:
                cursor.execute(f"ALTER TABLE users ADD COLUMN {col} {col_def}")
            except Exception:
                pass  # Column already exists

        # Add new columns to analytics_events (safe ALTER)
        for col, col_def in [
            ("tool_name", "TEXT"),
            ("action_type", "TEXT"),
            ("file_size", "INTEGER"),
            ("processing_duration", "INTEGER"),
        ]:
            try:
                cursor.execute(f"ALTER TABLE analytics_events ADD COLUMN {col} {col_def}")
            except Exception:
                pass  # Column already exists

        # Create indexes
        cursor.execute("CREATE INDEX IF NOT EXISTS idx_anonymous_composite ON anonymous_usage(composite_hash)")
        cursor.execute("CREATE INDEX IF NOT EXISTS idx_anonymous_ip ON anonymous_usage(ip_address)")
        cursor.execute("CREATE INDEX IF NOT EXISTS idx_users_email ON users(email)")
        cursor.execute("CREATE INDEX IF NOT EXISTS idx_vouchers_code ON vouchers(code)")
        cursor.execute("CREATE INDEX IF NOT EXISTS idx_usage_user ON usage_history(user_id)")
        cursor.execute("CREATE INDEX IF NOT EXISTS idx_usage_anon ON usage_history(anonymous_id)")
        cursor.execute("CREATE INDEX IF NOT EXISTS idx_analytics_timestamp ON analytics_events(timestamp)")
        cursor.execute("CREATE INDEX IF NOT EXISTS idx_analytics_event_type ON analytics_events(event_type)")
        cursor.execute("CREATE INDEX IF NOT EXISTS idx_analytics_page ON analytics_events(page)")
        cursor.execute("CREATE INDEX IF NOT EXISTS idx_analytics_tool ON analytics_events(tool_name)")
        cursor.execute("CREATE INDEX IF NOT EXISTS idx_analytics_action ON analytics_events(action_type)")
        cursor.execute("CREATE INDEX IF NOT EXISTS idx_heartbeats_ip ON heartbeats(ip)")
        cursor.execute("CREATE INDEX IF NOT EXISTS idx_heartbeats_session ON heartbeats(session_id)")
        cursor.execute("CREATE INDEX IF NOT EXISTS idx_heartbeats_last_seen ON heartbeats(last_seen)")

        # Insert default settings if not exist
        default_settings = {
            "stripe_secret_key": "",
            "stripe_publishable_key": "",
            "stripe_webhook_secret": "",
            "jwt_secret_key": JWT_SECRET_KEY,
            "bronze_monthly_price_id": "",
            "bronze_annual_price_id": "",
            "silver_monthly_price_id": "",
            "silver_annual_price_id": "",
            "gold_monthly_price_id": "",
            "gold_annual_price_id": "",
            "rate_limit": "100",
            "max_file_size_mb": "500",
            "maintenance_mode": "false",
        }

        for key, value in default_settings.items():
            cursor.execute(
                "INSERT OR IGNORE INTO settings (key, value) VALUES (?, ?)",
                (key, value)
            )

# ============================================================================
# Plan Definitions
# ============================================================================

PLAN_LIMITS = {
    "free": {
        "max_file_size_mb": 50,
        "pages_per_day": 200,
        "has_ads": True,
        "batch_processing": False,
        "smart_tools": False,
        "api_access": False,
        "price_monthly_eur": 0,
        "price_annual_eur": 0,
    },
    "bronze": {
        "max_file_size_mb": 150,
        "pages_per_day": 500,
        "has_ads": False,
        "batch_processing": True,
        "smart_tools": False,
        "api_access": False,
        "price_monthly_eur": 3.99,
        "price_annual_eur": 38.30,  # ~20% discount
    },
    "silver": {
        "max_file_size_mb": 300,
        "pages_per_day": 1000,
        "has_ads": False,
        "batch_processing": True,
        "smart_tools": True,
        "api_access": False,
        "price_monthly_eur": 7.99,
        "price_annual_eur": 76.70,  # ~20% discount
    },
    "gold": {
        "max_file_size_mb": 500,
        "pages_per_day": -1,  # unlimited
        "has_ads": False,
        "batch_processing": True,
        "smart_tools": True,
        "api_access": True,
        "price_monthly_eur": 17.99,
        "price_annual_eur": 172.70,  # ~20% discount
    },
    "custom": {
        "max_file_size_mb": 500,
        "pages_per_day": -1,
        "has_ads": False,
        "batch_processing": True,
        "smart_tools": True,
        "api_access": True,
        "price_monthly_eur": 0,
        "price_annual_eur": 0,
    }
}

# ============================================================================
# Helper Functions for Auth & Settings
# ============================================================================

def get_setting(key: str, default: str = "") -> str:
    """Get a setting value from database."""
    with db_session() as conn:
        cursor = conn.cursor()
        cursor.execute("SELECT value FROM settings WHERE key = ?", (key,))
        row = cursor.fetchone()
        return row["value"] if row else default

def set_setting(key: str, value: str) -> None:
    """Set a setting value in database."""
    with db_session() as conn:
        cursor = conn.cursor()
        cursor.execute(
            "INSERT OR REPLACE INTO settings (key, value, updated_at) VALUES (?, ?, ?)",
            (key, value, datetime.utcnow().isoformat())
        )

def get_stripe_client():
    """Get configured Stripe client."""
    if not STRIPE_AVAILABLE:
        return None
    secret_key = get_setting("stripe_secret_key")
    if secret_key:
        stripe.api_key = secret_key
        return stripe
    return None

def hash_password(password: str) -> str:
    """Hash a password."""
    return bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt()).decode('utf-8')

def verify_password(plain_password: str, hashed_password: str) -> bool:
    """Verify a password against its hash."""
    return bcrypt.checkpw(plain_password.encode('utf-8'), hashed_password.encode('utf-8'))

def create_jwt_token(user_id: str, email: str) -> str:
    """Create a JWT token for a user."""
    secret = get_setting("jwt_secret_key", JWT_SECRET_KEY)
    payload = {
        "sub": user_id,
        "email": email,
        "exp": datetime.utcnow() + timedelta(hours=JWT_EXPIRY_HOURS),
        "iat": datetime.utcnow()
    }
    return jwt.encode(payload, secret, algorithm=JWT_ALGORITHM)

def decode_jwt_token(token: str) -> Optional[Dict[str, Any]]:
    """Decode and verify a JWT token."""
    try:
        secret = get_setting("jwt_secret_key", JWT_SECRET_KEY)
        payload = jwt.decode(token, secret, algorithms=[JWT_ALGORITHM])
        return payload
    except jwt.ExpiredSignatureError:
        return None
    except jwt.InvalidTokenError:
        return None

def get_current_user(credentials: HTTPAuthorizationCredentials = Depends(security)) -> Optional[Dict[str, Any]]:
    """Get current authenticated user from JWT token."""
    if not credentials:
        return None

    payload = decode_jwt_token(credentials.credentials)
    if not payload:
        return None

    with db_session() as conn:
        cursor = conn.cursor()
        cursor.execute("SELECT * FROM users WHERE id = ?", (payload["sub"],))
        row = cursor.fetchone()
        if row:
            return dict(row)
    return None

def get_current_user_required(credentials: HTTPAuthorizationCredentials = Depends(security)) -> Dict[str, Any]:
    """Get current authenticated user, raise 401 if not authenticated."""
    user = get_current_user(credentials)
    if not user:
        raise HTTPException(status_code=401, detail="Authentication required")
    return user

def generate_composite_hash(ip: str, fingerprint: str = "", local_token: str = "") -> str:
    """Generate composite hash for anonymous user identification."""
    composite = f"{ip}:{fingerprint}:{local_token}"
    return hashlib.sha256(composite.encode()).hexdigest()

def get_client_ip(request: Request) -> str:
    """Extract client IP from request."""
    forwarded = request.headers.get("X-Forwarded-For")
    if forwarded:
        return forwarded.split(",")[0].strip()
    return request.client.host if request.client else "unknown"

def verify_superadmin(admin_key: str):
    """Verify superadmin key from env or DB setting."""
    if SUPER_ADMIN_KEY and admin_key == SUPER_ADMIN_KEY:
        return True
    stored_key = get_setting("admin_api_key", "")
    if stored_key and admin_key == stored_key:
        return True
    raise HTTPException(status_code=403, detail="Invalid admin key")


def get_user_plan_limits(user: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    """Get plan limits for a user (or free limits if no user)."""
    if not user:
        return PLAN_LIMITS["free"]

    plan = user.get("plan", "free")
    base_limits = PLAN_LIMITS.get(plan, PLAN_LIMITS["free"]).copy()

    # Apply custom overrides for custom plan
    if plan == "custom":
        if user.get("custom_max_file_size_mb"):
            base_limits["max_file_size_mb"] = user["custom_max_file_size_mb"]
        if user.get("custom_pages_per_day"):
            base_limits["pages_per_day"] = user["custom_pages_per_day"]
        if user.get("custom_features"):
            try:
                features = json.loads(user["custom_features"])
                base_limits.update(features)
            except json.JSONDecodeError:
                pass

    return base_limits


def record_user_usage(user_id: str, operation: str, pages_processed: int, file_size_bytes: int = 0) -> None:
    """Record usage for an authenticated user."""
    with db_session() as conn:
        cursor = conn.cursor()
        cursor.execute("""
            INSERT INTO usage_history (id, user_id, operation, pages_processed, file_size_bytes)
            VALUES (?, ?, ?, ?, ?)
        """, (uuid.uuid4().hex, user_id, operation, pages_processed, file_size_bytes))


def check_user_limits(
    user: Optional[Dict[str, Any]],
    pages_to_process: int,
    file_size_bytes: int,
    request: Request,
    fingerprint: str = "",
    local_token: str = ""
) -> Dict[str, Any]:
    """Check if user (or anonymous) can process the given pages/file size. Returns check result."""
    limits = get_user_plan_limits(user)

    # Check file size
    max_bytes = limits["max_file_size_mb"] * 1024 * 1024
    if file_size_bytes > max_bytes:
        return {
            "allowed": False,
            "reason": "file_too_large",
            "detail": f"File size exceeds {limits['max_file_size_mb']}MB limit for your plan",
            "limits": limits
        }

    # Check page limit (skip if unlimited)
    if limits["pages_per_day"] == -1:
        return {"allowed": True, "limits": limits}

    today = datetime.utcnow().strftime("%Y-%m-%d")

    if user:
        # Authenticated user
        with db_session() as conn:
            cursor = conn.cursor()
            cursor.execute("""
                SELECT COALESCE(SUM(pages_processed), 0) as pages_today
                FROM usage_history
                WHERE user_id = ? AND DATE(created_at) = ?
            """, (user["id"], today))
            usage = cursor.fetchone()
            pages_used = usage["pages_today"] if usage else 0
    else:
        # Anonymous user
        ip = get_client_ip(request)
        composite_hash = generate_composite_hash(ip, fingerprint, local_token)

        with db_session() as conn:
            cursor = conn.cursor()
            cursor.execute("""
                SELECT pages_used_today FROM anonymous_usage
                WHERE composite_hash = ? AND last_reset_date = ?
            """, (composite_hash, today))
            anon = cursor.fetchone()
            pages_used = anon["pages_used_today"] if anon else 0

    if pages_used + pages_to_process > limits["pages_per_day"]:
        return {
            "allowed": False,
            "reason": "limit_reached",
            "detail": f"Daily page limit ({limits['pages_per_day']} pages) reached",
            "pages_used": pages_used,
            "pages_limit": limits["pages_per_day"],
            "limits": limits
        }

    return {"allowed": True, "pages_used": pages_used, "limits": limits}

def check_and_update_voucher_expiry(user_id: str) -> None:
    """Check if user's voucher has expired and revert to original plan."""
    with db_session() as conn:
        cursor = conn.cursor()
        cursor.execute("""
            SELECT vr.*, u.base_plan
            FROM voucher_redemptions vr
            JOIN users u ON u.id = vr.user_id
            WHERE vr.user_id = ? AND vr.expires_at > datetime('now')
            ORDER BY vr.expires_at DESC
            LIMIT 1
        """, (user_id,))
        active_redemption = cursor.fetchone()

        if not active_redemption:
            # No active voucher, check if we need to revert
            cursor.execute("""
                SELECT base_plan FROM users WHERE id = ?
            """, (user_id,))
            user = cursor.fetchone()
            if user:
                cursor.execute("""
                    UPDATE users SET plan = base_plan WHERE id = ? AND plan != base_plan
                """, (user_id,))

# Directories
UPLOAD_DIR = Path("uploads")
OUTPUT_DIR = Path("output")
UPLOAD_DIR.mkdir(exist_ok=True)
OUTPUT_DIR.mkdir(exist_ok=True)

# File expiration (24 hours)
FILE_EXPIRY_SECONDS = 24 * 60 * 60

# Global state for jobs
jobs = {}
ocr_engine = None
batch_processing_jobs = {}
MAX_BATCH_WORKERS = 4

def get_ocr():
    global ocr_engine
    if ocr_engine is None:
        ocr_engine = RapidOCR()
    return ocr_engine


def extract_order_from_page(page, ocr, pattern, crop_left=0.5, crop_top=0.0, crop_right=1.0, crop_bottom=0.2, dpi=150):
    """OCR the specified region of a PDF page and extract order number."""
    pw, ph = page.rect.width, page.rect.height
    clip = fitz.Rect(pw * crop_left, ph * crop_top, pw * crop_right, ph * crop_bottom)
    mat = fitz.Matrix(dpi / 72, dpi / 72)
    pix = page.get_pixmap(matrix=mat, clip=clip)
    img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, 3)

    result, _ = ocr(img)
    texts = [r[1] for r in result] if result else []
    full_text = ' '.join(texts)

    match = re.search(pattern, full_text, re.IGNORECASE)
    order_num = match.group(1) if match else None

    return order_num, full_text


def process_pdf_job(job_id, input_path, output_dir, pattern, crop_config, dpi, filename_template="{order}"):
    """Background job to process a PDF."""
    job = jobs[job_id]
    try:
        ocr = get_ocr()
        doc = fitz.open(input_path)
        total_pages = len(doc)
        job["total_pages"] = total_pages
        job["status"] = "scanning"

        # Phase 1: Scan all pages
        page_orders = []
        for i in range(total_pages):
            page = doc[i]
            order_num, raw_text = extract_order_from_page(
                page, ocr, pattern,
                crop_config["left"], crop_config["top"],
                crop_config["right"], crop_config["bottom"],
                dpi
            )
            page_orders.append((i, order_num))
            job["current_page"] = i + 1
            job["log"].append({
                "page": i + 1,
                "order": order_num,
                "text": raw_text[:150]
            })

        # Phase 2: Group pages
        job["status"] = "grouping"
        groups = {}
        current_order = None
        orphan_pages = []

        for page_idx, order_num in page_orders:
            if order_num and order_num != current_order:
                current_order = order_num
                if current_order not in groups:
                    groups[current_order] = []
                groups[current_order].append(page_idx)
            elif current_order:
                groups[current_order].append(page_idx)
            else:
                orphan_pages.append(page_idx)

        job["orders_found"] = len(groups)
        job["groups"] = {k: [p + 1 for p in v] for k, v in groups.items()}

        # Phase 3: Create split PDFs
        job["status"] = "splitting"
        os.makedirs(output_dir, exist_ok=True)
        created = []

        for idx, (order, pages) in enumerate(groups.items(), 1):
            fname = filename_template.replace("{order}", order).replace("{pages}", str(len(pages))).replace("{index}", str(idx))
            output_path = os.path.join(output_dir, f"{fname}.pdf")
            counter = 1
            while os.path.exists(output_path):
                output_path = os.path.join(output_dir, f"{fname}_{counter}.pdf")
                counter += 1

            new_doc = fitz.open()
            for page_idx in pages:
                new_doc.insert_pdf(doc, from_page=page_idx, to_page=page_idx)
            new_doc.save(output_path)
            new_doc.close()

            file_info = {
                "name": os.path.basename(output_path),
                "path": output_path,
                "order": order,
                "pages": len(pages),
                "size_mb": round(os.path.getsize(output_path) / (1024 * 1024), 2)
            }
            created.append(file_info)

        if orphan_pages:
            orphan_path = os.path.join(output_dir, "_no_order_detected.pdf")
            new_doc = fitz.open()
            for page_idx in orphan_pages:
                new_doc.insert_pdf(doc, from_page=page_idx, to_page=page_idx)
            new_doc.save(orphan_path)
            new_doc.close()
            created.append({
                "name": "_no_order_detected.pdf",
                "path": orphan_path,
                "order": "N/A",
                "pages": len(orphan_pages),
                "size_mb": round(os.path.getsize(orphan_path) / (1024 * 1024), 2)
            })

        doc.close()
        job["status"] = "done"
        job["files"] = created
        job["finished_at"] = time.time()

    except Exception as e:
        job["status"] = "error"
        job["error"] = str(e)


@app.get("/api/health")
def health():
    return {"status": "ok"}


# ============================================================================
# Authentication API Endpoints
# ============================================================================

@app.post("/api/auth/register")
async def register(
    email: str = Form(...),
    password: str = Form(...),
    plan: str = Form("free")
):
    """Register a new user."""
    if not email or not password:
        raise HTTPException(status_code=400, detail="Email and password are required")

    if len(password) < 8:
        raise HTTPException(status_code=400, detail="Password must be at least 8 characters")

    email = email.lower().strip()

    with db_session() as conn:
        cursor = conn.cursor()

        # Check if email already exists
        cursor.execute("SELECT id FROM users WHERE email = ?", (email,))
        if cursor.fetchone():
            raise HTTPException(status_code=400, detail="Email already registered")

        # Create user
        user_id = uuid.uuid4().hex
        password_hash = hash_password(password)

        cursor.execute("""
            INSERT INTO users (id, email, password_hash, plan, base_plan)
            VALUES (?, ?, ?, ?, ?)
        """, (user_id, email, password_hash, "free", "free"))

        token = create_jwt_token(user_id, email)

        return {
            "success": True,
            "token": token,
            "user": {
                "id": user_id,
                "email": email,
                "plan": "free"
            }
        }


@app.post("/api/auth/login")
async def login(request: Request):
    """Login and get JWT token. Accepts JSON or Form data."""
    content_type = request.headers.get("content-type", "")
    if "application/json" in content_type:
        body = await request.json()
        email = body.get("email", "")
        password = body.get("password", "")
    else:
        form = await request.form()
        email = form.get("email", "")
        password = form.get("password", "")

    if not email or not password:
        raise HTTPException(status_code=400, detail="Email and password are required")

    email = email.lower().strip()

    with db_session() as conn:
        cursor = conn.cursor()
        cursor.execute("SELECT * FROM users WHERE email = ?", (email,))
        user = cursor.fetchone()

        if not user or not verify_password(password, user["password_hash"]):
            raise HTTPException(status_code=401, detail="Invalid email or password")

        # Check and update voucher expiry
        check_and_update_voucher_expiry(user["id"])

        # Refresh user data after voucher check
        cursor.execute("SELECT * FROM users WHERE id = ?", (user["id"],))
        user = cursor.fetchone()
        user_dict = dict(user)

        role = user_dict.get("role", "user") or "user"
        token = create_jwt_token(user_dict["id"], user_dict["email"])

        response_data = {
            "success": True,
            "token": token,
            "user": {
                "id": user_dict["id"],
                "email": user_dict["email"],
                "name": user_dict.get("name", ""),
                "role": role,
                "plan": user_dict["plan"],
                "subscription_status": user_dict["subscription_status"]
            },
            "message": "Login successful"
        }

        # If superadmin, also set httpOnly cookie for superadmin panel access
        if role == "superadmin":
            secret = get_setting("jwt_secret_key", JWT_SECRET_KEY)
            sa_token = jwt.encode(
                {
                    "sub": user_dict["id"],
                    "email": user_dict["email"],
                    "role": "superadmin",
                    "exp": datetime.utcnow() + timedelta(hours=24),
                    "iat": datetime.utcnow(),
                },
                secret,
                algorithm=JWT_ALGORITHM,
            )
            from fastapi.responses import JSONResponse
            response = JSONResponse(content=response_data)
            response.set_cookie(
                key="superadmin_jwt",
                value=sa_token,
                httponly=True,
                secure=True,
                samesite="strict",
                max_age=86400,
                path="/",
            )
            return response

        return response_data


@app.get("/api/auth/me")
async def get_me(user: Dict[str, Any] = Depends(get_current_user_required)):
    """Get current user profile."""
    check_and_update_voucher_expiry(user["id"])

    # Refresh user data
    with db_session() as conn:
        cursor = conn.cursor()
        cursor.execute("SELECT * FROM users WHERE id = ?", (user["id"],))
        user = dict(cursor.fetchone())

    plan_limits = get_user_plan_limits(user)

    # Get usage stats
    with db_session() as conn:
        cursor = conn.cursor()
        today = datetime.utcnow().strftime("%Y-%m-%d")
        cursor.execute("""
            SELECT COALESCE(SUM(pages_processed), 0) as pages_today
            FROM usage_history
            WHERE user_id = ? AND DATE(created_at) = ?
        """, (user["id"], today))
        usage = cursor.fetchone()

    return {
        "id": user["id"],
        "email": user["email"],
        "plan": user["plan"],
        "role": user.get("role", "user") or "user",
        "base_plan": user["base_plan"],
        "subscription_status": user["subscription_status"],
        "subscription_end_date": user["subscription_end_date"],
        "limits": plan_limits,
        "usage": {
            "pages_today": usage["pages_today"] if usage else 0,
            "pages_limit": plan_limits["pages_per_day"]
        }
    }


@app.post("/api/auth/change-password")
async def change_password(
    current_password: str = Form(...),
    new_password: str = Form(...),
    user: Dict[str, Any] = Depends(get_current_user_required)
):
    """Change user password."""
    if len(new_password) < 8:
        raise HTTPException(status_code=400, detail="Password must be at least 8 characters")

    if not verify_password(current_password, user["password_hash"]):
        raise HTTPException(status_code=401, detail="Current password is incorrect")

    with db_session() as conn:
        cursor = conn.cursor()
        cursor.execute(
            "UPDATE users SET password_hash = ?, updated_at = ? WHERE id = ?",
            (hash_password(new_password), datetime.utcnow().isoformat(), user["id"])
        )

    return {"success": True}


# ============================================================================
# Anonymous Usage Tracking API
# ============================================================================

@app.post("/api/track/anonymous")
async def track_anonymous_usage(
    request: Request,
    fingerprint: str = Form(""),
    local_token: str = Form(""),
    pages: int = Form(0)
):
    """Track anonymous user usage. Returns usage stats and limits."""
    ip = get_client_ip(request)
    composite_hash = generate_composite_hash(ip, fingerprint, local_token)
    today = datetime.utcnow().strftime("%Y-%m-%d")

    with db_session() as conn:
        cursor = conn.cursor()

        # Find or create anonymous user record
        cursor.execute("SELECT * FROM anonymous_usage WHERE composite_hash = ?", (composite_hash,))
        anon_user = cursor.fetchone()

        if anon_user:
            # Check if we need to reset daily counter
            if anon_user["last_reset_date"] != today:
                cursor.execute("""
                    UPDATE anonymous_usage
                    SET pages_used_today = ?, last_reset_date = ?, last_seen = ?,
                        total_pages_all_time = total_pages_all_time + ?
                    WHERE id = ?
                """, (pages, today, datetime.utcnow().isoformat(), pages, anon_user["id"]))
                pages_used = pages
            else:
                new_total = anon_user["pages_used_today"] + pages
                cursor.execute("""
                    UPDATE anonymous_usage
                    SET pages_used_today = ?, last_seen = ?, total_pages_all_time = total_pages_all_time + ?
                    WHERE id = ?
                """, (new_total, datetime.utcnow().isoformat(), pages, anon_user["id"]))
                pages_used = new_total
            anon_id = anon_user["id"]
        else:
            # Create new anonymous user
            anon_id = uuid.uuid4().hex
            cursor.execute("""
                INSERT INTO anonymous_usage
                (id, ip_address, fingerprint, local_token, composite_hash, pages_used_today, last_reset_date, total_pages_all_time)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?)
            """, (anon_id, ip, fingerprint, local_token, composite_hash, pages, today, pages))
            pages_used = pages

        # Record usage history
        if pages > 0:
            cursor.execute("""
                INSERT INTO usage_history (id, anonymous_id, operation, pages_processed)
                VALUES (?, ?, ?, ?)
            """, (uuid.uuid4().hex, anon_id, "pdf_operation", pages))

    limits = PLAN_LIMITS["free"]
    limit_reached = pages_used >= limits["pages_per_day"]

    return {
        "anonymous_id": anon_id,
        "pages_used_today": pages_used,
        "pages_limit": limits["pages_per_day"],
        "limit_reached": limit_reached,
        "plan": "free",
        "has_ads": limits["has_ads"],
        "upgrade_message": "Create an account to unlock more features!" if limit_reached else None
    }


@app.post("/api/check-limits")
async def check_processing_limits(
    request: Request,
    pages: int = Form(0),
    file_size_bytes: int = Form(0),
    fingerprint: str = Form(""),
    local_token: str = Form(""),
    credentials: HTTPAuthorizationCredentials = Depends(security)
):
    """Check if user can process the given pages/file before starting operation."""
    user = get_current_user(credentials)

    if user:
        check_and_update_voucher_expiry(user["id"])

    result = check_user_limits(user, pages, file_size_bytes, request, fingerprint, local_token)

    if not result["allowed"]:
        return JSONResponse(
            status_code=403,
            content={
                "allowed": False,
                "reason": result["reason"],
                "detail": result["detail"],
                "limits": result["limits"],
                "upgrade_url": "/pricing"
            }
        )

    return {
        "allowed": True,
        "pages_used": result.get("pages_used", 0),
        "limits": result["limits"]
    }


@app.get("/api/track/usage")
async def get_usage_status(
    request: Request,
    fingerprint: str = "",
    local_token: str = "",
    credentials: HTTPAuthorizationCredentials = Depends(security)
):
    """Get current usage status for user or anonymous visitor."""
    user = get_current_user(credentials)

    if user:
        check_and_update_voucher_expiry(user["id"])
        limits = get_user_plan_limits(user)

        with db_session() as conn:
            cursor = conn.cursor()
            today = datetime.utcnow().strftime("%Y-%m-%d")
            cursor.execute("""
                SELECT COALESCE(SUM(pages_processed), 0) as pages_today
                FROM usage_history
                WHERE user_id = ? AND DATE(created_at) = ?
            """, (user["id"], today))
            usage = cursor.fetchone()

        pages_used = usage["pages_today"] if usage else 0
        pages_limit = limits["pages_per_day"]

        return {
            "authenticated": True,
            "user_id": user["id"],
            "email": user["email"],
            "plan": user["plan"],
            "pages_used_today": pages_used,
            "pages_limit": pages_limit,
            "limit_reached": pages_limit != -1 and pages_used >= pages_limit,
            "has_ads": limits["has_ads"],
            "limits": limits
        }
    else:
        # Anonymous user
        ip = get_client_ip(request)
        composite_hash = generate_composite_hash(ip, fingerprint, local_token)

        with db_session() as conn:
            cursor = conn.cursor()
            today = datetime.utcnow().strftime("%Y-%m-%d")
            cursor.execute("""
                SELECT * FROM anonymous_usage
                WHERE composite_hash = ? AND last_reset_date = ?
            """, (composite_hash, today))
            anon_user = cursor.fetchone()

        pages_used = anon_user["pages_used_today"] if anon_user else 0
        limits = PLAN_LIMITS["free"]

        return {
            "authenticated": False,
            "plan": "free",
            "pages_used_today": pages_used,
            "pages_limit": limits["pages_per_day"],
            "limit_reached": pages_used >= limits["pages_per_day"],
            "has_ads": limits["has_ads"],
            "limits": limits
        }


# ============================================================================
# Voucher System API
# ============================================================================

@app.post("/api/voucher/redeem")
async def redeem_voucher(
    code: str = Form(...),
    user: Dict[str, Any] = Depends(get_current_user_required)
):
    """Redeem a voucher code."""
    code = code.strip().upper()

    with db_session() as conn:
        cursor = conn.cursor()

        # Find voucher
        cursor.execute("""
            SELECT * FROM vouchers
            WHERE code = ? AND is_active = 1
            AND (expires_at IS NULL OR expires_at > datetime('now'))
            AND (max_uses = 0 OR current_uses < max_uses)
        """, (code,))
        voucher = cursor.fetchone()

        if not voucher:
            raise HTTPException(status_code=400, detail="Invalid or expired voucher code")

        # Check if user already redeemed this voucher
        cursor.execute("""
            SELECT id FROM voucher_redemptions
            WHERE voucher_id = ? AND user_id = ?
        """, (voucher["id"], user["id"]))
        if cursor.fetchone():
            raise HTTPException(status_code=400, detail="You have already redeemed this voucher")

        # Calculate expiry date
        expires_at = datetime.utcnow() + timedelta(days=voucher["duration_days"])

        # Create redemption record
        cursor.execute("""
            INSERT INTO voucher_redemptions (id, voucher_id, user_id, expires_at, original_plan)
            VALUES (?, ?, ?, ?, ?)
        """, (uuid.uuid4().hex, voucher["id"], user["id"], expires_at.isoformat(), user["plan"]))

        # Update user's plan
        cursor.execute("""
            UPDATE users SET plan = ?, base_plan = ?, updated_at = ?
            WHERE id = ?
        """, (voucher["target_plan"], user["plan"], datetime.utcnow().isoformat(), user["id"]))

        # Increment voucher usage
        cursor.execute("""
            UPDATE vouchers SET current_uses = current_uses + 1 WHERE id = ?
        """, (voucher["id"],))

    return {
        "success": True,
        "plan": voucher["target_plan"],
        "expires_at": expires_at.isoformat(),
        "message": f"Successfully upgraded to {voucher['target_plan'].title()} plan for {voucher['duration_days']} days!"
    }


# Admin endpoint for creating vouchers
@app.post("/api/admin/vouchers")
async def create_voucher(
    request: Request,
    code: str = Form(...),
    target_plan: str = Form(...),
    duration_days: int = Form(...),
    max_uses: int = Form(1),
    expires_at: str = Form(None),
    admin_key: str = Form(None),
):
    """Create a new voucher (admin only)."""
    if admin_key:
        verify_superadmin(admin_key)
    else:
        _require_superadmin(request)

    if target_plan not in PLAN_LIMITS:
        raise HTTPException(status_code=400, detail=f"Invalid plan: {target_plan}")

    code = code.strip().upper()

    with db_session() as conn:
        cursor = conn.cursor()

        # Check if code already exists
        cursor.execute("SELECT id FROM vouchers WHERE code = ?", (code,))
        if cursor.fetchone():
            raise HTTPException(status_code=400, detail="Voucher code already exists")

        voucher_id = uuid.uuid4().hex
        cursor.execute("""
            INSERT INTO vouchers (id, code, target_plan, duration_days, max_uses, expires_at)
            VALUES (?, ?, ?, ?, ?, ?)
        """, (voucher_id, code, target_plan, duration_days, max_uses, expires_at))

    return {
        "success": True,
        "voucher_id": voucher_id,
        "code": code,
        "target_plan": target_plan,
        "duration_days": duration_days,
        "max_uses": max_uses
    }


@app.get("/api/admin/vouchers")
async def list_vouchers(request: Request, admin_key: str = None):
    """List all vouchers (admin only)."""
    if admin_key:
        verify_superadmin(admin_key)
    else:
        _require_superadmin(request)

    with db_session() as conn:
        cursor = conn.cursor()
        cursor.execute("""
            SELECT v.*,
                   (SELECT COUNT(*) FROM voucher_redemptions WHERE voucher_id = v.id) as redemption_count
            FROM vouchers v
            ORDER BY created_at DESC
        """)
        vouchers = [dict(row) for row in cursor.fetchall()]

    return {"vouchers": vouchers}


# ============================================================================
# Stripe Payment Integration
# ============================================================================

@app.get("/api/plans")
async def get_plans():
    """Get available subscription plans."""
    plans = []
    for plan_name, limits in PLAN_LIMITS.items():
        if plan_name == "custom":
            continue

        plan_info = {
            "name": plan_name,
            "display_name": plan_name.title(),
            "price_monthly_eur": limits["price_monthly_eur"],
            "price_annual_eur": limits["price_annual_eur"],
            "max_file_size_mb": limits["max_file_size_mb"],
            "pages_per_day": limits["pages_per_day"],
            "has_ads": limits["has_ads"],
            "batch_processing": limits["batch_processing"],
            "smart_tools": limits["smart_tools"],
            "api_access": limits["api_access"],
        }

        if plan_name != "free":
            plan_info["monthly_price_id"] = get_setting(f"{plan_name}_monthly_price_id")
            plan_info["annual_price_id"] = get_setting(f"{plan_name}_annual_price_id")

        plans.append(plan_info)

    return {
        "plans": plans,
        "stripe_publishable_key": get_setting("stripe_publishable_key"),
        "custom_plan_contact": {
            "whatsapp": "+40749591399",
            "email": "office@4updf.com"
        }
    }


@app.post("/api/stripe/create-checkout")
async def create_checkout_session(
    plan: str = Form(...),
    billing_period: str = Form("monthly"),  # monthly or annual
    user: Dict[str, Any] = Depends(get_current_user_required)
):
    """Create Stripe checkout session."""
    stripe_client = get_stripe_client()
    if not stripe_client:
        raise HTTPException(status_code=503, detail="Stripe is not configured")

    if plan not in ["bronze", "silver", "gold"]:
        raise HTTPException(status_code=400, detail="Invalid plan")

    price_id = get_setting(f"{plan}_{billing_period}_price_id")
    if not price_id:
        raise HTTPException(status_code=503, detail=f"Price ID not configured for {plan} {billing_period}")

    try:
        # Get or create Stripe customer
        customer_id = user.get("stripe_customer_id")
        if not customer_id:
            customer = stripe_client.Customer.create(
                email=user["email"],
                metadata={"user_id": user["id"]}
            )
            customer_id = customer.id

            with db_session() as conn:
                cursor = conn.cursor()
                cursor.execute(
                    "UPDATE users SET stripe_customer_id = ? WHERE id = ?",
                    (customer_id, user["id"])
                )

        # Create checkout session
        session = stripe_client.checkout.Session.create(
            customer=customer_id,
            payment_method_types=["card"],
            line_items=[{"price": price_id, "quantity": 1}],
            mode="subscription",
            success_url=f"https://4updf.com/dashboard?success=true&session_id={{CHECKOUT_SESSION_ID}}",
            cancel_url="https://4updf.com/pricing?canceled=true",
            metadata={
                "user_id": user["id"],
                "plan": plan,
                "billing_period": billing_period
            }
        )

        return {"checkout_url": session.url, "session_id": session.id}

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/stripe/webhook")
async def stripe_webhook(request: Request):
    """Handle Stripe webhooks."""
    stripe_client = get_stripe_client()
    if not stripe_client:
        raise HTTPException(status_code=503, detail="Stripe is not configured")

    payload = await request.body()
    sig_header = request.headers.get("stripe-signature")
    webhook_secret = get_setting("stripe_webhook_secret")

    try:
        event = stripe_client.Webhook.construct_event(
            payload, sig_header, webhook_secret
        )
    except ValueError:
        raise HTTPException(status_code=400, detail="Invalid payload")
    except stripe.error.SignatureVerificationError:
        raise HTTPException(status_code=400, detail="Invalid signature")

    # Handle different event types
    if event.type == "checkout.session.completed":
        session = event.data.object
        user_id = session.metadata.get("user_id")
        plan = session.metadata.get("plan")

        if user_id and plan:
            with db_session() as conn:
                cursor = conn.cursor()
                cursor.execute("""
                    UPDATE users
                    SET plan = ?, base_plan = ?, stripe_subscription_id = ?,
                        subscription_status = 'active', updated_at = ?
                    WHERE id = ?
                """, (plan, plan, session.subscription, datetime.utcnow().isoformat(), user_id))

    elif event.type == "customer.subscription.updated":
        subscription = event.data.object
        customer_id = subscription.customer

        with db_session() as conn:
            cursor = conn.cursor()
            cursor.execute("SELECT id FROM users WHERE stripe_customer_id = ?", (customer_id,))
            user = cursor.fetchone()
            if user:
                status = "active" if subscription.status == "active" else "inactive"
                end_date = datetime.fromtimestamp(subscription.current_period_end).isoformat()
                cursor.execute("""
                    UPDATE users
                    SET subscription_status = ?, subscription_end_date = ?, updated_at = ?
                    WHERE id = ?
                """, (status, end_date, datetime.utcnow().isoformat(), user["id"]))

    elif event.type == "customer.subscription.deleted":
        subscription = event.data.object
        customer_id = subscription.customer

        with db_session() as conn:
            cursor = conn.cursor()
            cursor.execute("SELECT id FROM users WHERE stripe_customer_id = ?", (customer_id,))
            user = cursor.fetchone()
            if user:
                cursor.execute("""
                    UPDATE users
                    SET plan = 'free', base_plan = 'free', subscription_status = 'canceled',
                        stripe_subscription_id = NULL, updated_at = ?
                    WHERE id = ?
                """, (datetime.utcnow().isoformat(), user["id"]))

    return {"received": True}


@app.post("/api/stripe/create-portal")
async def create_customer_portal(user: Dict[str, Any] = Depends(get_current_user_required)):
    """Create Stripe customer portal session for subscription management."""
    stripe_client = get_stripe_client()
    if not stripe_client:
        raise HTTPException(status_code=503, detail="Stripe is not configured")

    customer_id = user.get("stripe_customer_id")
    if not customer_id:
        raise HTTPException(status_code=400, detail="No subscription found")

    try:
        session = stripe_client.billing_portal.Session.create(
            customer=customer_id,
            return_url="https://4updf.com/dashboard"
        )
        return {"portal_url": session.url}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ============================================================================
# SuperAdmin Authentication
# ============================================================================

## Old key-based superadmin login/verify removed — now handled by routes/superadmin.py


def _require_superadmin(request: Request):
    """Require superadmin JWT from httpOnly cookie or Authorization header."""
    token = request.cookies.get("superadmin_jwt")
    if not token:
        auth = request.headers.get("Authorization", "")
        if auth.startswith("Bearer "):
            token = auth.split(" ", 1)[1]
        else:
            raise HTTPException(status_code=401, detail="Authentication required")
    try:
        secret = get_setting("jwt_secret_key", JWT_SECRET_KEY)
        payload = jwt.decode(token, secret, algorithms=[JWT_ALGORITHM])
        if payload.get("role") != "superadmin":
            raise HTTPException(status_code=403, detail="Not superadmin")
    except Exception:
        raise HTTPException(status_code=401, detail="Invalid token")


# ============================================================================
# Admin Settings API
# ============================================================================

@app.get("/api/admin/settings")
async def get_admin_settings(request: Request, admin_key: str = None):
    """Get admin settings (admin only)."""
    if admin_key:
        verify_superadmin(admin_key)
    else:
        _require_superadmin(request)

    settings_keys = [
        "stripe_secret_key", "stripe_publishable_key", "stripe_webhook_secret",
        "jwt_secret_key", "admin_api_key",
        "bronze_monthly_price_id", "bronze_annual_price_id",
        "silver_monthly_price_id", "silver_annual_price_id",
        "gold_monthly_price_id", "gold_annual_price_id",
        "rate_limit", "max_file_size_mb", "maintenance_mode",
    ]

    settings = {}
    for key in settings_keys:
        value = get_setting(key, "")
        if "secret" in key.lower() or "key" in key.lower():
            if value:
                settings[key] = value[:8] + "..." + value[-4:] if len(value) > 12 else "***"
            else:
                settings[key] = ""
        else:
            settings[key] = value

    return {"settings": settings}


@app.post("/api/admin/settings")
async def update_admin_settings(request: Request):
    """Update admin settings (admin only)."""
    _require_superadmin(request)

    form_data = await request.form()

    allowed_keys = [
        "stripe_secret_key", "stripe_publishable_key", "stripe_webhook_secret",
        "jwt_secret_key", "admin_api_key",
        "bronze_monthly_price_id", "bronze_annual_price_id",
        "silver_monthly_price_id", "silver_annual_price_id",
        "gold_monthly_price_id", "gold_annual_price_id",
        "rate_limit", "max_file_size_mb", "maintenance_mode",
    ]

    updated = []
    for key in allowed_keys:
        if key in form_data:
            value = form_data[key]
            if value is not None:
                set_setting(key, str(value))
                updated.append(key)

    return {"success": True, "updated": updated}


# ============================================================================
# Admin User Management API
# ============================================================================

@app.get("/api/admin/users")
async def list_users(
    request: Request,
    admin_key: str = None,
    limit: int = 50,
    offset: int = 0,
    search: str = None,
    plan_filter: str = None,
    sort_by: str = "created_at",
    sort_dir: str = "desc",
):
    """List users (admin only) with search, filter, sort."""
    if admin_key:
        verify_superadmin(admin_key)
    else:
        _require_superadmin(request)

    allowed_sort = {"created_at", "last_active", "email", "plan"}
    if sort_by not in allowed_sort:
        sort_by = "created_at"
    if sort_dir not in ("asc", "desc"):
        sort_dir = "desc"

    with db_session() as conn:
        cursor = conn.cursor()

        where_clauses = []
        params: list = []

        if search:
            where_clauses.append("email LIKE ?")
            params.append(f"%{search}%")
        if plan_filter:
            where_clauses.append("plan = ?")
            params.append(plan_filter)

        where_sql = f"WHERE {' AND '.join(where_clauses)}" if where_clauses else ""

        cursor.execute(f"""
            SELECT id, email, plan, base_plan, subscription_status,
                   created_at, updated_at, last_active, is_banned
            FROM users
            {where_sql}
            ORDER BY {sort_by} {sort_dir}
            LIMIT ? OFFSET ?
        """, params + [limit, offset])
        users = [dict(row) for row in cursor.fetchall()]

        cursor.execute(f"SELECT COUNT(*) as count FROM users {where_sql}", params)
        total = cursor.fetchone()["count"]

    return {"users": users, "total": total}


@app.post("/api/admin/users/{user_id}/ban")
async def toggle_ban_user(user_id: str, request: Request):
    """Toggle ban/unban for a user."""
    _require_superadmin(request)

    with db_session() as conn:
        cursor = conn.cursor()
        cursor.execute("SELECT id, is_banned FROM users WHERE id = ?", (user_id,))
        row = cursor.fetchone()
        if not row:
            raise HTTPException(status_code=404, detail="User not found")
        new_status = 0 if row["is_banned"] else 1
        cursor.execute("UPDATE users SET is_banned = ?, updated_at = ? WHERE id = ?",
                        (new_status, datetime.utcnow().isoformat(), user_id))

    return {"success": True, "is_banned": bool(new_status)}


@app.post("/api/admin/users/{user_id}/plan")
async def update_user_plan(
    user_id: str,
    request: Request,
    plan: str = Form(...),
    custom_max_file_size_mb: int = Form(None),
    custom_pages_per_day: int = Form(None),
    admin_key: str = Form(None),
):
    """Update user plan (admin only)."""
    if admin_key:
        verify_superadmin(admin_key)
    else:
        _require_superadmin(request)

    if plan not in PLAN_LIMITS:
        raise HTTPException(status_code=400, detail=f"Invalid plan: {plan}")

    with db_session() as conn:
        cursor = conn.cursor()
        cursor.execute("SELECT id FROM users WHERE id = ?", (user_id,))
        if not cursor.fetchone():
            raise HTTPException(status_code=404, detail="User not found")

        cursor.execute("""
            UPDATE users
            SET plan = ?, base_plan = ?, custom_max_file_size_mb = ?,
                custom_pages_per_day = ?, updated_at = ?
            WHERE id = ?
        """, (plan, plan, custom_max_file_size_mb, custom_pages_per_day,
              datetime.utcnow().isoformat(), user_id))

    return {"success": True}


# ============================================================================
# Admin Voucher Management (enhanced CRUD)
# ============================================================================

@app.put("/api/admin/vouchers/{voucher_id}")
async def update_voucher(voucher_id: str, request: Request):
    """Update a voucher (toggle active, change fields)."""
    _require_superadmin(request)
    form_data = await request.form()

    with db_session() as conn:
        cursor = conn.cursor()
        cursor.execute("SELECT id FROM vouchers WHERE id = ?", (voucher_id,))
        if not cursor.fetchone():
            raise HTTPException(status_code=404, detail="Voucher not found")

        updates = []
        params = []
        for field in ["is_active", "max_uses", "expires_at", "target_plan", "duration_days"]:
            if field in form_data:
                updates.append(f"{field} = ?")
                params.append(form_data[field])

        if updates:
            params.append(voucher_id)
            cursor.execute(f"UPDATE vouchers SET {', '.join(updates)} WHERE id = ?", params)

    return {"success": True}


@app.delete("/api/admin/vouchers/{voucher_id}")
async def delete_voucher(voucher_id: str, request: Request):
    """Delete a voucher."""
    _require_superadmin(request)

    with db_session() as conn:
        cursor = conn.cursor()
        cursor.execute("DELETE FROM vouchers WHERE id = ?", (voucher_id,))
        if cursor.rowcount == 0:
            raise HTTPException(status_code=404, detail="Voucher not found")

    return {"success": True}


# ============================================================================
# Admin Stats & Real-time Active Users (SSE)
# ============================================================================

@app.get("/api/admin/stats")
async def get_admin_stats(request: Request):
    """Get admin dashboard stats."""
    _require_superadmin(request)

    with db_session() as conn:
        cursor = conn.cursor()

        # Total users
        cursor.execute("SELECT COUNT(*) as c FROM users")
        total_users = cursor.fetchone()["c"]

        # Users by plan
        cursor.execute("SELECT plan, COUNT(*) as c FROM users GROUP BY plan")
        users_by_plan = {row["plan"]: row["c"] for row in cursor.fetchall()}

        # Total uploads/operations
        cursor.execute("SELECT COUNT(*) as c FROM usage_history")
        total_operations = cursor.fetchone()["c"]

        # Operations per tool
        cursor.execute("""
            SELECT operation, COUNT(*) as c
            FROM usage_history
            GROUP BY operation
            ORDER BY c DESC
        """)
        operations_by_tool = {row["operation"]: row["c"] for row in cursor.fetchall()}

        # Operations last 30 days (daily)
        cursor.execute("""
            SELECT DATE(created_at) as day, COUNT(*) as c
            FROM usage_history
            WHERE created_at >= DATE('now', '-30 days')
            GROUP BY DATE(created_at)
            ORDER BY day
        """)
        daily_operations = [{"date": row["day"], "count": row["c"]} for row in cursor.fetchall()]

        # New users last 30 days (daily)
        cursor.execute("""
            SELECT DATE(created_at) as day, COUNT(*) as c
            FROM users
            WHERE created_at >= DATE('now', '-30 days')
            GROUP BY DATE(created_at)
            ORDER BY day
        """)
        daily_new_users = [{"date": row["day"], "count": row["c"]} for row in cursor.fetchall()]

        # Active users (heartbeat within last 5 min)
        now = time.time()
        active_count = sum(1 for t in _active_users.values() if now - t < 300)

    return {
        "total_users": total_users,
        "users_by_plan": users_by_plan,
        "total_operations": total_operations,
        "operations_by_tool": operations_by_tool,
        "daily_operations": daily_operations,
        "daily_new_users": daily_new_users,
        "active_users_now": active_count,
    }


@app.post("/api/heartbeat")
async def heartbeat(request: Request):
    """Track active users via heartbeat with page/tool info."""
    ip = get_client_ip(request)
    _active_users[ip] = time.time()

    # Parse optional JSON body for enhanced tracking
    current_page = None
    tool_name = None
    session_id = None
    user_id = None
    try:
        body = await request.json()
        current_page = body.get("current_page")
        tool_name = body.get("tool_name")
        session_id = body.get("session_id")
        user_id = body.get("user_id")
    except Exception:
        pass

    # Also check JWT for user_id
    if not user_id:
        auth = request.headers.get("Authorization", "")
        if auth.startswith("Bearer "):
            try:
                payload = decode_jwt_token(auth.split(" ", 1)[1])
                if payload and payload.get("sub"):
                    user_id = payload["sub"]
            except Exception:
                pass

    user_agent = request.headers.get("user-agent", "")

    try:
        with db_session() as conn:
            cursor = conn.cursor()
            # Upsert heartbeat by IP + session_id
            hb_key = f"{ip}:{session_id or 'default'}"
            cursor.execute("SELECT id FROM heartbeats WHERE id = ?", (hb_key,))
            if cursor.fetchone():
                cursor.execute("""
                    UPDATE heartbeats SET last_seen = ?, current_page = ?, tool_name = ?,
                    user_id = ?, user_agent = ? WHERE id = ?
                """, (datetime.utcnow().isoformat(), current_page, tool_name, user_id, user_agent, hb_key))
            else:
                now = datetime.utcnow().isoformat()
                cursor.execute("""
                    INSERT INTO heartbeats (id, user_id, ip, user_agent, current_page, tool_name, session_id, last_seen, first_seen)
                    VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
                """, (hb_key, user_id, ip, user_agent, current_page, tool_name, session_id, now, now))
    except Exception:
        pass

    return {"ok": True}


@app.post("/api/analytics/heartbeat")
async def analytics_heartbeat(request: Request):
    """Alias for heartbeat - frontend sends every 30s with current page."""
    return await heartbeat(request)


@app.get("/api/admin/active-users")
async def sse_active_users(request: Request):
    """SSE endpoint for real-time active users with details."""
    _require_superadmin(request)

    async def event_generator():
        while True:
            now = time.time()
            # In-memory count
            count = sum(1 for t in _active_users.values() if now - t < 300)

            # Detailed active users from heartbeats table (last 5 min)
            active_list = []
            try:
                cutoff = (datetime.utcnow() - timedelta(minutes=5)).isoformat()
                with db_session() as conn:
                    cursor = conn.cursor()
                    cursor.execute("""
                        SELECT h.user_id, h.ip, h.user_agent, h.current_page, h.tool_name,
                               h.first_seen, h.last_seen,
                               u.email
                        FROM heartbeats h
                        LEFT JOIN users u ON h.user_id = u.id
                        WHERE h.last_seen >= ?
                        ORDER BY h.last_seen DESC
                    """, (cutoff,))
                    for row in cursor.fetchall():
                        r = dict(row)
                        # Calculate session duration
                        try:
                            first = datetime.fromisoformat(r["first_seen"])
                            last = datetime.fromisoformat(r["last_seen"])
                            duration_sec = int((last - first).total_seconds())
                        except Exception:
                            duration_sec = 0

                        # Parse browser from user agent
                        ua = r.get("user_agent") or ""
                        browser = "Unknown"
                        if "Firefox" in ua:
                            browser = "Firefox"
                        elif "Edg/" in ua:
                            browser = "Edge"
                        elif "Chrome" in ua:
                            browser = "Chrome"
                        elif "Safari" in ua:
                            browser = "Safari"

                        email = r.get("email") or (f"Anonymous#{hashlib.md5(r['ip'].encode()).hexdigest()[:6]}" if r.get("ip") else "Anonymous")

                        active_list.append({
                            "email": email,
                            "current_page": r.get("current_page") or "/",
                            "tool": r.get("tool_name") or "-",
                            "duration_sec": duration_sec,
                            "ip": r.get("ip") or "",
                            "browser": browser,
                        })
            except Exception:
                pass

            yield f"data: {json.dumps({'active_users': count, 'users': active_list, 'timestamp': datetime.utcnow().isoformat()})}\n\n"
            await asyncio.sleep(5)

    return StreamingResponse(event_generator(), media_type="text/event-stream",
                              headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"})


# ============================================================================
# Analytics Tracking & Dashboard API
# ============================================================================

@app.middleware("http")
async def analytics_middleware(request: Request, call_next):
    """Track page views and API tool calls for local analytics."""
    start_time = time.time()

    # Read body for POST requests before call_next consumes it
    content_length = 0
    try:
        cl = request.headers.get("content-length")
        if cl:
            content_length = int(cl)
    except Exception:
        pass

    response = await call_next(request)

    path = request.url.path
    ip = get_client_ip(request)
    user_agent = request.headers.get("user-agent", "")
    referrer = request.headers.get("referer", "")
    processing_duration = int((time.time() - start_time) * 1000)  # ms

    # Extract user_id from JWT if present
    user_id = None
    auth = request.headers.get("Authorization", "")
    if auth.startswith("Bearer "):
        try:
            token = auth.split(" ", 1)[1]
            payload = decode_jwt_token(token)
            if payload and payload.get("sub"):
                user_id = payload["sub"]
        except Exception:
            pass

    # Track pageviews (GET to pages)
    if (request.method == "GET"
        and not path.startswith("/api/")
        and not path.startswith("/_next/")
        and not path.startswith("/static/")
        and not path.endswith((".js", ".css", ".ico", ".png", ".jpg", ".svg", ".woff2"))):
        try:
            with db_session() as conn:
                cursor = conn.cursor()
                cursor.execute("""
                    INSERT INTO analytics_events (id, event_type, page, user_id, ip, user_agent, referrer, timestamp)
                    VALUES (?, 'pageview', ?, ?, ?, ?, ?, ?)
                """, (uuid.uuid4().hex, path, user_id, ip, user_agent, referrer,
                      datetime.utcnow().isoformat()))
        except Exception:
            pass

    # Track API tool calls (POST to tool endpoints)
    _tool_endpoints = {
        "/api/merge": "merge", "/api/split": "split", "/api/split-pages": "split-pages",
        "/api/split-ocr": "split-ocr", "/api/compress": "compress",
        "/api/pdf-to-jpg": "pdf-to-jpg", "/api/pdf-to-png": "pdf-to-png",
        "/api/pdf-to-word": "pdf-to-word", "/api/pdf-to-excel": "pdf-to-excel",
        "/api/pdf-to-powerpoint": "pdf-to-ppt", "/api/pdf-to-text": "pdf-to-text",
        "/api/word-to-pdf": "word-to-pdf", "/api/excel-to-pdf": "excel-to-pdf",
        "/api/html-to-pdf": "html-to-pdf", "/api/rotate": "rotate",
        "/api/extract-pages": "extract-pages", "/api/crop": "crop",
        "/api/watermark": "watermark", "/api/protect": "protect",
        "/api/unlock": "unlock", "/api/redact-pdf": "redact",
        "/api/organize-pdf": "organize", "/api/flatten": "flatten",
        "/api/repair": "repair", "/api/split-by-text": "split-by-text",
        "/api/split-invoices": "split-invoices", "/api/ocr-layer": "ocr-layer",
        "/api/extract-text-ocr": "extract-text-ocr", "/api/sign-pdf": "sign",
    }
    if request.method == "POST" and path in _tool_endpoints:
        tool_name = _tool_endpoints[path]
        action_type = "process"
        status_code = response.status_code
        try:
            with db_session() as conn:
                cursor = conn.cursor()
                cursor.execute("""
                    INSERT INTO analytics_events
                    (id, event_type, page, user_id, ip, user_agent, tool_name, action_type,
                     file_size, processing_duration, timestamp)
                    VALUES (?, 'tool_use', ?, ?, ?, ?, ?, ?, ?, ?, ?)
                """, (uuid.uuid4().hex, path, user_id, ip, user_agent, tool_name,
                      action_type, content_length, processing_duration,
                      datetime.utcnow().isoformat()))
        except Exception:
            pass

    # Update last_active for authenticated users
    if user_id:
        try:
            with db_session() as conn:
                cursor = conn.cursor()
                cursor.execute("UPDATE users SET last_active = ? WHERE id = ?",
                                (datetime.utcnow().isoformat(), user_id))
        except Exception:
            pass

    return response


@app.post("/api/analytics/track")
async def track_analytics_event(
    request: Request,
    event_type: str = Form(...),
    page: str = Form(None),
):
    """Track custom analytics events from frontend."""
    ip = get_client_ip(request)
    user_agent = request.headers.get("user-agent", "")

    with db_session() as conn:
        cursor = conn.cursor()
        cursor.execute("""
            INSERT INTO analytics_events (id, event_type, page, ip, user_agent, timestamp)
            VALUES (?, ?, ?, ?, ?, ?)
        """, (uuid.uuid4().hex, event_type, page, ip, user_agent,
              datetime.utcnow().isoformat()))

    return {"success": True}


@app.get("/api/admin/analytics")
async def get_admin_analytics(
    request: Request,
    period: str = "30d",
):
    """Get local analytics data (page views, unique visitors, top pages)."""
    _require_superadmin(request)

    period_map = {"7d": 7, "30d": 30, "90d": 90, "365d": 365}
    days = period_map.get(period, 30)

    with db_session() as conn:
        cursor = conn.cursor()

        # Total page views in period
        cursor.execute("""
            SELECT COUNT(*) as c FROM analytics_events
            WHERE event_type = 'pageview' AND timestamp >= DATE('now', ?)
        """, (f"-{days} days",))
        total_pageviews = cursor.fetchone()["c"]

        # Unique visitors (by IP) in period
        cursor.execute("""
            SELECT COUNT(DISTINCT ip) as c FROM analytics_events
            WHERE event_type = 'pageview' AND timestamp >= DATE('now', ?)
        """, (f"-{days} days",))
        unique_visitors = cursor.fetchone()["c"]

        # Top pages
        cursor.execute("""
            SELECT page, COUNT(*) as views, COUNT(DISTINCT ip) as visitors
            FROM analytics_events
            WHERE event_type = 'pageview' AND timestamp >= DATE('now', ?)
            GROUP BY page
            ORDER BY views DESC
            LIMIT 20
        """, (f"-{days} days",))
        top_pages = [dict(row) for row in cursor.fetchall()]

        # Daily pageviews
        cursor.execute("""
            SELECT DATE(timestamp) as day, COUNT(*) as views, COUNT(DISTINCT ip) as visitors
            FROM analytics_events
            WHERE event_type = 'pageview' AND timestamp >= DATE('now', ?)
            GROUP BY DATE(timestamp)
            ORDER BY day
        """, (f"-{days} days",))
        daily_stats = [dict(row) for row in cursor.fetchall()]

        # Top referrers
        cursor.execute("""
            SELECT referrer, COUNT(*) as c
            FROM analytics_events
            WHERE event_type = 'pageview' AND referrer != '' AND timestamp >= DATE('now', ?)
            GROUP BY referrer
            ORDER BY c DESC
            LIMIT 10
        """, (f"-{days} days",))
        top_referrers = [dict(row) for row in cursor.fetchall()]

    return {
        "total_pageviews": total_pageviews,
        "unique_visitors": unique_visitors,
        "top_pages": top_pages,
        "daily_stats": daily_stats,
        "top_referrers": top_referrers,
        "period": period,
    }


@app.get("/api/admin/sessions")
async def get_admin_sessions(
    request: Request,
    period: str = "24h",
):
    """Get user sessions for the given period."""
    _require_superadmin(request)

    period_map = {"24h": 1, "7d": 7, "30d": 30}
    days = period_map.get(period, 1)

    with db_session() as conn:
        cursor = conn.cursor()
        cutoff = (datetime.utcnow() - timedelta(days=days)).isoformat()

        # Group analytics events by IP into sessions (gap > 30 min = new session)
        cursor.execute("""
            SELECT ae.ip, ae.user_id, ae.page, ae.tool_name, ae.user_agent, ae.timestamp,
                   u.email
            FROM analytics_events ae
            LEFT JOIN users u ON ae.user_id = u.id
            WHERE ae.timestamp >= ?
            ORDER BY ae.ip, ae.timestamp
        """, (cutoff,))

        rows = [dict(r) for r in cursor.fetchall()]

    # Build sessions from sequential events grouped by IP
    sessions = []
    current_session = None
    for row in rows:
        ip = row["ip"] or ""
        ts = row["timestamp"] or ""
        try:
            ts_dt = datetime.fromisoformat(ts)
        except Exception:
            continue

        # Start new session if different IP or gap > 30 min
        if (current_session is None
            or current_session["ip"] != ip
            or (ts_dt - current_session["_last_ts"]).total_seconds() > 1800):
            if current_session:
                current_session.pop("_last_ts", None)
                current_session.pop("ip", None)
                sessions.append(current_session)
            email = row.get("email") or (f"Anonymous#{hashlib.md5(ip.encode()).hexdigest()[:6]}" if ip else "Anonymous")

            ua = row.get("user_agent") or ""
            country = "Unknown"
            # Parse browser for session info
            browser = "Unknown"
            if "Firefox" in ua:
                browser = "Firefox"
            elif "Edg/" in ua:
                browser = "Edge"
            elif "Chrome" in ua:
                browser = "Chrome"
            elif "Safari" in ua:
                browser = "Safari"

            current_session = {
                "user": email,
                "pages": [],
                "tools": [],
                "start_time": ts,
                "end_time": ts,
                "duration_sec": 0,
                "country": country,
                "browser": browser,
                "ip": ip,
                "_last_ts": ts_dt,
            }
        else:
            current_session["end_time"] = ts
            current_session["_last_ts"] = ts_dt
            current_session["duration_sec"] = int((ts_dt - datetime.fromisoformat(current_session["start_time"])).total_seconds())

        page = row.get("page") or ""
        if page and page not in current_session["pages"]:
            current_session["pages"].append(page)
        tool = row.get("tool_name") or ""
        if tool and tool not in current_session["tools"]:
            current_session["tools"].append(tool)

    if current_session:
        current_session.pop("_last_ts", None)
        current_session.pop("ip", None)
        sessions.append(current_session)

    # Sort by start_time descending, limit to 200
    sessions.sort(key=lambda s: s["start_time"], reverse=True)
    return {"sessions": sessions[:200], "period": period}


@app.get("/api/admin/tool-usage")
async def get_admin_tool_usage(
    request: Request,
    period: str = "30d",
):
    """Get tool usage breakdown for the given period."""
    _require_superadmin(request)

    period_map = {"7d": 7, "30d": 30, "90d": 90, "365d": 365}
    days = period_map.get(period, 30)

    with db_session() as conn:
        cursor = conn.cursor()

        # From analytics_events with tool_name
        cursor.execute("""
            SELECT tool_name, COUNT(*) as count
            FROM analytics_events
            WHERE tool_name IS NOT NULL AND tool_name != ''
            AND timestamp >= DATE('now', ?)
            GROUP BY tool_name
            ORDER BY count DESC
        """, (f"-{days} days",))
        tool_events = {row["tool_name"]: row["count"] for row in cursor.fetchall()}

        # Also count from usage_history operations
        cursor.execute("""
            SELECT operation, COUNT(*) as count
            FROM usage_history
            WHERE created_at >= DATE('now', ?)
            GROUP BY operation
            ORDER BY count DESC
        """, (f"-{days} days",))
        for row in cursor.fetchall():
            op = row["operation"]
            tool_events[op] = tool_events.get(op, 0) + row["count"]

    tools = [{"tool": k, "count": v} for k, v in sorted(tool_events.items(), key=lambda x: x[1], reverse=True)]
    return {"tools": tools, "period": period}


@app.get("/api/admin/activity-log")
async def get_admin_activity_log(
    request: Request,
    page: int = 1,
    per_page: int = 100,
):
    """Get paginated activity log."""
    _require_superadmin(request)

    offset = (page - 1) * per_page

    with db_session() as conn:
        cursor = conn.cursor()

        # Total count
        cursor.execute("SELECT COUNT(*) as c FROM analytics_events")
        total = cursor.fetchone()["c"]

        # Paginated results
        cursor.execute("""
            SELECT ae.id, ae.event_type, ae.page, ae.user_id, ae.ip, ae.tool_name,
                   ae.action_type, ae.file_size, ae.processing_duration,
                   ae.timestamp, u.email
            FROM analytics_events ae
            LEFT JOIN users u ON ae.user_id = u.id
            ORDER BY ae.timestamp DESC
            LIMIT ? OFFSET ?
        """, (per_page, offset))

        events = []
        for row in cursor.fetchall():
            r = dict(row)
            email = r.get("email") or (f"Anonymous#{hashlib.md5(r['ip'].encode()).hexdigest()[:6]}" if r.get("ip") else "Anonymous")
            events.append({
                "timestamp": r["timestamp"],
                "user": email,
                "action": r.get("action_type") or r.get("event_type") or "pageview",
                "tool": r.get("tool_name") or "-",
                "page": r.get("page") or "-",
                "file_size": r.get("file_size"),
                "duration_ms": r.get("processing_duration"),
                "status": "success",
            })

    return {
        "events": events,
        "total": total,
        "page": page,
        "per_page": per_page,
        "total_pages": (total + per_page - 1) // per_page,
    }


@app.post("/api/split-ocr")
@app.post("/api/split")
async def start_split(
    file: UploadFile = File(None),
    input_path: str = Form(None),
    output_dir: str = Form("output"),
    pattern: str = Form(r"(?:Order\s*(?:No\.?|#|:)?\s*)(\d{8})"),
    crop_left: float = Form(0.5),
    crop_top: float = Form(0.0),
    crop_right: float = Form(1.0),
    crop_bottom: float = Form(0.2),
    dpi: int = Form(150),
    filename_template: str = Form("{order}"),
):
    """Start a PDF split job."""
    # Handle file upload or path
    if file and file.filename:
        os.makedirs("uploads", exist_ok=True)
        save_path = os.path.join("uploads", file.filename)
        with open(save_path, "wb") as f:
            content = await file.read()
            f.write(content)
        input_path = save_path
    elif not input_path:
        return JSONResponse({"error": "No file or input_path provided"}, status_code=400)

    if not os.path.exists(input_path):
        return JSONResponse({"error": f"File not found: {input_path}"}, status_code=404)

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "id": job_id,
        "status": "starting",
        "input": input_path,
        "output_dir": output_dir,
        "total_pages": 0,
        "current_page": 0,
        "orders_found": 0,
        "groups": {},
        "files": [],
        "log": [],
        "started_at": time.time(),
        "finished_at": None,
        "error": None
    }

    crop_config = {
        "left": crop_left,
        "top": crop_top,
        "right": crop_right,
        "bottom": crop_bottom
    }

    thread = threading.Thread(
        target=process_pdf_job,
        args=(job_id, input_path, output_dir, pattern, crop_config, dpi, filename_template),
        daemon=True
    )
    thread.start()

    return {"job_id": job_id, "status": "started"}


@app.get("/api/status/{job_id}")
def get_status(job_id: str):
    """Get job status and progress."""
    if job_id not in jobs:
        return JSONResponse({"error": "Job not found"}, status_code=404)

    job = jobs[job_id]
    progress = 0
    if job["total_pages"] > 0:
        progress = round(job["current_page"] / job["total_pages"] * 100, 1)

    elapsed = time.time() - job["started_at"]
    eta = 0
    if job["current_page"] > 0 and job["status"] == "scanning":
        per_page = elapsed / job["current_page"]
        remaining = job["total_pages"] - job["current_page"]
        eta = round(per_page * remaining)

    return {
        "id": job["id"],
        "status": job["status"],
        "progress": progress,
        "current_page": job["current_page"],
        "total_pages": job["total_pages"],
        "orders_found": job["orders_found"],
        "elapsed_seconds": round(elapsed),
        "eta_seconds": eta,
        "files": job["files"],
        "groups": job["groups"],
        "error": job["error"],
        "recent_log": job["log"][-5:] if job["log"] else []
    }


@app.get("/api/jobs")
def list_jobs():
    """List all jobs."""
    return [
        {
            "id": j.get("id"),
            "status": j.get("status"),
            "input": j.get("input", ""),
            "total_pages": j.get("total_pages", 0),
            "orders_found": j.get("orders_found", 0),
            "started_at": j.get("started_at", 0)
        }
        for j in jobs.values()
    ]


@app.get("/api/download/{filename}")
def download_file(filename: str, output_dir: str = "output"):
    """Download a split PDF file."""
    path = os.path.join(output_dir, filename)
    if not os.path.exists(path):
        return JSONResponse({"error": "File not found"}, status_code=404)
    return FileResponse(path, filename=filename, media_type="application/pdf")


@app.get("/api/browse")
def browse_folder(path: str = "."):
    """Browse filesystem for input/output folder selection."""
    try:
        p = Path(path).resolve()
        if not p.exists():
            return JSONResponse({"error": "Path not found"}, status_code=404)

        items = []
        if p.parent != p:
            items.append({"name": "..", "type": "dir", "path": str(p.parent)})

        for item in sorted(p.iterdir()):
            try:
                if item.is_dir():
                    items.append({"name": item.name, "type": "dir", "path": str(item)})
                elif item.suffix.lower() == '.pdf':
                    size_mb = round(item.stat().st_size / (1024 * 1024), 2)
                    items.append({"name": item.name, "type": "pdf", "path": str(item), "size_mb": size_mb})
            except PermissionError:
                continue

        return {"current": str(p), "items": items}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/api/download-all")
def download_all(output_dir: str = "output"):
    """Download all split PDFs as a ZIP file."""
    p = Path(output_dir)
    if not p.exists():
        return JSONResponse({"error": "Output directory not found"}, status_code=404)

    pdfs = list(p.glob("*.pdf"))
    if not pdfs:
        return JSONResponse({"error": "No PDF files found"}, status_code=404)

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for pdf in pdfs:
            zf.write(pdf, pdf.name)
    buf.seek(0)

    return StreamingResponse(
        buf,
        media_type="application/zip",
        headers={"Content-Disposition": "attachment; filename=split-pdfs.zip"}
    )


@app.get("/api/defaults")
def get_defaults():
    """Return server-appropriate default paths."""
    base = Path(__file__).parent.resolve()
    input_dir = base / "input"
    output_dir = base / "output"
    input_dir.mkdir(exist_ok=True)
    output_dir.mkdir(exist_ok=True)
    return {
        "inputPath": str(input_dir),
        "outputDir": str(output_dir),
    }


# ============================================================================
# File Cleanup - Auto-delete files after 24 hours
# ============================================================================

def cleanup_old_files():
    """Delete files older than FILE_EXPIRY_SECONDS."""
    now = time.time()
    for directory in [UPLOAD_DIR, OUTPUT_DIR]:
        if not directory.exists():
            continue
        for file_path in directory.iterdir():
            if file_path.is_file():
                file_age = now - file_path.stat().st_mtime
                if file_age > FILE_EXPIRY_SECONDS:
                    try:
                        file_path.unlink()
                    except Exception:
                        pass


def start_cleanup_thread():
    """Background thread for periodic file cleanup."""
    def cleanup_loop():
        while True:
            cleanup_old_files()
            time.sleep(3600)  # Run every hour
    thread = threading.Thread(target=cleanup_loop, daemon=True)
    thread.start()


# Start cleanup on app startup
@app.on_event("startup")
async def startup_event():
    init_database()
    # Initialize superadmin routes with app dependencies
    init_superadmin_routes(JWT_SECRET_KEY, JWT_ALGORITHM, get_setting, db_session)
    start_cleanup_thread()
    cleanup_old_files()


# ============================================================================
# Helper Functions
# ============================================================================

def save_upload_file(file: UploadFile) -> Path:
    """Save uploaded file with unique name."""
    ext = Path(file.filename).suffix
    unique_name = f"{uuid.uuid4().hex}{ext}"
    file_path = UPLOAD_DIR / unique_name
    with open(file_path, "wb") as f:
        shutil.copyfileobj(file.file, f)
    return file_path


def create_output_path(prefix: str, ext: str = ".pdf") -> Path:
    """Create unique output file path."""
    unique_name = f"{prefix}_{uuid.uuid4().hex[:8]}{ext}"
    return OUTPUT_DIR / unique_name


# ============================================================================
# Merge PDF API
# ============================================================================

@app.post("/api/merge")
async def merge_pdfs(files: List[UploadFile] = File(...), order: str = Form(None)):
    """Merge multiple PDFs into one."""
    if len(files) < 2:
        raise HTTPException(status_code=400, detail="At least 2 PDF files required")

    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_files = []
    try:
        # Save uploaded files
        for f in files:
            saved_files.append(save_upload_file(f))

        # Parse order if provided
        if order:
            indices = [int(i) for i in order.split(",")]
            saved_files = [saved_files[i] for i in indices if i < len(saved_files)]

        # Merge PDFs
        output_path = create_output_path("merged")
        merged = fitz.open()
        total = len(saved_files)

        for i, pdf_path in enumerate(saved_files):
            doc = fitz.open(str(pdf_path))
            merged.insert_pdf(doc)
            doc.close()
            jobs[job_id]["progress"] = int((i + 1) / total * 100)

        merged.save(str(output_path))
        merged.close()

        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100
        jobs[job_id]["output_file"] = output_path.name

        # Get page count properly without memory leak
        temp_doc = fitz.open(str(output_path))
        page_count = temp_doc.page_count
        temp_doc.close()

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_path.name}",
            "filename": output_path.name,
            "pages": page_count
        }
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        for f in saved_files:
            try:
                f.unlink()
            except:
                pass


# ============================================================================
# Split PDF API (by page ranges)
# ============================================================================

@app.post("/api/split-pages")
async def split_pdf_pages(
    file: UploadFile = File(...),
    ranges: str = Form(...),  # e.g., "1-3,5,7-10" or "all" for single pages
    split_mode: str = Form("ranges")  # "ranges" or "every" or "single"
):
    """Split PDF by page ranges or extract specific pages."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = save_upload_file(file)
    try:
        doc = fitz.open(str(saved_file))
        total_pages = doc.page_count
        output_files = []

        if split_mode == "single":
            # Split into single pages
            for i in range(total_pages):
                out_path = create_output_path(f"page_{i+1}")
                new_doc = fitz.open()
                new_doc.insert_pdf(doc, from_page=i, to_page=i)
                new_doc.save(str(out_path))
                new_doc.close()
                output_files.append(out_path.name)
                jobs[job_id]["progress"] = int((i + 1) / total_pages * 100)

        elif split_mode == "every":
            # Split every N pages
            n = int(ranges)
            for start in range(0, total_pages, n):
                end = min(start + n - 1, total_pages - 1)
                out_path = create_output_path(f"pages_{start+1}-{end+1}")
                new_doc = fitz.open()
                new_doc.insert_pdf(doc, from_page=start, to_page=end)
                new_doc.save(str(out_path))
                new_doc.close()
                output_files.append(out_path.name)
                jobs[job_id]["progress"] = int((end + 1) / total_pages * 100)

        else:
            # Parse ranges like "1-3,5,7-10"
            page_sets = []
            for part in ranges.split(","):
                part = part.strip()
                if "-" in part:
                    start, end = map(int, part.split("-"))
                    page_sets.append((start - 1, end - 1))
                else:
                    p = int(part) - 1
                    page_sets.append((p, p))

            for idx, (start, end) in enumerate(page_sets):
                start = max(0, min(start, total_pages - 1))
                end = max(0, min(end, total_pages - 1))
                out_path = create_output_path(f"pages_{start+1}-{end+1}")
                new_doc = fitz.open()
                new_doc.insert_pdf(doc, from_page=start, to_page=end)
                new_doc.save(str(out_path))
                new_doc.close()
                output_files.append(out_path.name)
                jobs[job_id]["progress"] = int((idx + 1) / len(page_sets) * 100)

        doc.close()
        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100
        jobs[job_id]["output_files"] = output_files

        return {
            "job_id": job_id,
            "status": "done",
            "files": [{"filename": f, "download_url": f"/api/download/{f}"} for f in output_files],
            "total_files": len(output_files)
        }
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            saved_file.unlink()
        except:
            pass


# ============================================================================
# Compress PDF API
# ============================================================================

@app.post("/api/compress")
async def compress_pdf(
    file: UploadFile = File(...),
    quality: str = Form("medium")  # low, medium, high
):
    """Compress PDF with specified quality level."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = save_upload_file(file)
    try:
        doc = fitz.open(str(saved_file))
        original_size = saved_file.stat().st_size

        # Quality settings
        quality_settings = {
            "low": {"image_quality": 90, "garbage": 1},
            "medium": {"image_quality": 75, "garbage": 3},
            "high": {"image_quality": 50, "garbage": 4}
        }
        settings = quality_settings.get(quality, quality_settings["medium"])

        output_path = create_output_path("compressed")

        # Compress images in PDF
        total_pages = doc.page_count
        for i, page in enumerate(doc):
            image_list = page.get_images(full=True)
            for img_idx, img in enumerate(image_list):
                xref = img[0]
                try:
                    base_image = doc.extract_image(xref)
                    if base_image:
                        image_bytes = base_image["image"]
                        # Re-encode at lower quality using PIL if available
                        try:
                            from PIL import Image
                            img_io = io.BytesIO(image_bytes)
                            pil_img = Image.open(img_io)
                            if pil_img.mode in ("RGBA", "P"):
                                pil_img = pil_img.convert("RGB")
                            out_io = io.BytesIO()
                            pil_img.save(out_io, format="JPEG", quality=settings["image_quality"], optimize=True)
                            page.delete_image(xref)
                            # Note: Full image replacement requires more complex logic
                        except ImportError:
                            pass
                except:
                    pass
            jobs[job_id]["progress"] = int((i + 1) / total_pages * 50)

        # Save with compression options
        doc.save(
            str(output_path),
            garbage=settings["garbage"],
            deflate=True,
            clean=True
        )

        compressed_size = output_path.stat().st_size
        doc.close()

        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_path.name}",
            "filename": output_path.name,
            "original_size": original_size,
            "compressed_size": compressed_size,
            "reduction_percent": round((1 - compressed_size / original_size) * 100, 1) if original_size > 0 else 0
        }
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            saved_file.unlink()
        except:
            pass


# ============================================================================
# PDF to JPG API
# ============================================================================

@app.post("/api/pdf-to-jpg")
async def pdf_to_jpg(
    file: UploadFile = File(...),
    dpi: int = Form(150),
    pages: str = Form("all")  # "all" or "1,2,3" or "1-5"
):
    """Convert PDF pages to JPG images."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = save_upload_file(file)
    try:
        doc = fitz.open(str(saved_file))
        total_pages = doc.page_count
        output_files = []

        # Determine which pages to convert
        if pages == "all":
            page_list = list(range(total_pages))
        else:
            page_list = []
            for part in pages.split(","):
                part = part.strip()
                if "-" in part:
                    start, end = map(int, part.split("-"))
                    page_list.extend(range(start - 1, end))
                else:
                    page_list.append(int(part) - 1)
            page_list = [p for p in page_list if 0 <= p < total_pages]

        # Convert pages
        mat = fitz.Matrix(dpi / 72, dpi / 72)
        for idx, page_num in enumerate(page_list):
            page = doc[page_num]
            pix = page.get_pixmap(matrix=mat)
            out_path = OUTPUT_DIR / f"page_{page_num+1}_{uuid.uuid4().hex[:8]}.jpg"
            pix.save(str(out_path))
            output_files.append(out_path.name)
            jobs[job_id]["progress"] = int((idx + 1) / len(page_list) * 100)

        doc.close()
        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100
        jobs[job_id]["output_files"] = output_files

        # Create ZIP if multiple files
        if len(output_files) > 1:
            zip_path = create_output_path("images", ".zip")
            with zipfile.ZipFile(str(zip_path), "w", zipfile.ZIP_DEFLATED) as zf:
                for fname in output_files:
                    zf.write(OUTPUT_DIR / fname, fname)

            return {
                "job_id": job_id,
                "status": "done",
                "download_url": f"/api/download/{zip_path.name}",
                "filename": zip_path.name,
                "files": [{"filename": f, "download_url": f"/api/download/{f}"} for f in output_files],
                "total_images": len(output_files)
            }
        else:
            return {
                "job_id": job_id,
                "status": "done",
                "download_url": f"/api/download/{output_files[0]}",
                "filename": output_files[0],
                "total_images": 1
            }
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            saved_file.unlink()
        except:
            pass


# ============================================================================
# JPG to PDF API
# ============================================================================

@app.post("/api/jpg-to-pdf")
async def jpg_to_pdf(
    files: List[UploadFile] = File(...),
    order: str = Form(None),
    page_size: str = Form("a4")  # a4, letter, original
):
    """Convert JPG images to PDF."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_files = []
    try:
        # Save uploaded files
        for f in files:
            saved_files.append(save_upload_file(f))

        # Parse order if provided
        if order:
            indices = [int(i) for i in order.split(",")]
            saved_files = [saved_files[i] for i in indices if i < len(saved_files)]

        # Page size settings (in points)
        page_sizes = {
            "a4": (595, 842),
            "letter": (612, 792),
            "original": None
        }
        target_size = page_sizes.get(page_size, page_sizes["a4"])

        output_path = create_output_path("images_to_pdf")
        doc = fitz.open()

        total = len(saved_files)
        for i, img_path in enumerate(saved_files):
            img = fitz.open(str(img_path))

            if target_size:
                # Scale image to fit page
                rect = fitz.Rect(0, 0, target_size[0], target_size[1])
            else:
                # Use original image size
                pdfbytes = img.convert_to_pdf()
                img_pdf = fitz.open("pdf", pdfbytes)
                rect = img_pdf[0].rect

            pdfbytes = img.convert_to_pdf()
            img_pdf = fitz.open("pdf", pdfbytes)
            page = doc.new_page(width=rect.width, height=rect.height)
            page.show_pdf_page(rect, img_pdf, 0)

            img.close()
            img_pdf.close()
            jobs[job_id]["progress"] = int((i + 1) / total * 100)

        doc.save(str(output_path))
        doc.close()

        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_path.name}",
            "filename": output_path.name,
            "pages": len(saved_files)
        }
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        for f in saved_files:
            try:
                f.unlink()
            except:
                pass


# ============================================================================
# PDF to Word (DOCX) API
# ============================================================================

@app.post("/api/pdf-to-word")
async def pdf_to_word(file: UploadFile = File(...)):
    """Convert PDF to Word DOCX format."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = save_upload_file(file)
    try:
        doc = fitz.open(str(saved_file))

        # Try to use python-docx for proper DOCX output
        try:
            from docx import Document
            from docx.shared import Inches, Pt
            from docx.enum.text import WD_ALIGN_PARAGRAPH

            word_doc = Document()
            total_pages = doc.page_count

            for i, page in enumerate(doc):
                text = page.get_text()
                if text.strip():
                    paragraphs = text.split('\n')
                    for para in paragraphs:
                        if para.strip():
                            word_doc.add_paragraph(para)

                if i < total_pages - 1:
                    word_doc.add_page_break()

                jobs[job_id]["progress"] = int((i + 1) / total_pages * 100)

            output_path = create_output_path("converted", ".docx")
            word_doc.save(str(output_path))
            doc.close()

            jobs[job_id]["status"] = "done"
            jobs[job_id]["progress"] = 100

            return {
                "job_id": job_id,
                "status": "done",
                "download_url": f"/api/download/{output_path.name}",
                "filename": output_path.name,
                "pages": total_pages
            }
        except ImportError:
            # Fallback: Extract text to a simple text file
            output_path = create_output_path("converted", ".txt")
            with open(output_path, "w", encoding="utf-8") as f:
                for page in doc:
                    f.write(page.get_text())
                    f.write("\n\n--- Page Break ---\n\n")
            doc.close()

            jobs[job_id]["status"] = "done"
            jobs[job_id]["progress"] = 100

            return {
                "job_id": job_id,
                "status": "done",
                "download_url": f"/api/download/{output_path.name}",
                "filename": output_path.name,
                "note": "python-docx not installed, returned TXT instead"
            }
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            saved_file.unlink()
        except:
            pass


# ============================================================================
# Word (DOCX) to PDF API
# ============================================================================

@app.post("/api/word-to-pdf")
async def word_to_pdf(file: UploadFile = File(...)):
    """Convert Word DOCX to PDF format."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = save_upload_file(file)
    try:
        # Try using python-docx to read and fitz to create PDF
        try:
            from docx import Document

            word_doc = Document(str(saved_file))
            output_path = create_output_path("converted", ".pdf")

            # Create PDF with extracted text
            pdf_doc = fitz.open()
            page = pdf_doc.new_page()

            text_content = []
            for para in word_doc.paragraphs:
                if para.text.strip():
                    text_content.append(para.text)

            full_text = "\n\n".join(text_content)

            # Insert text into PDF
            rect = fitz.Rect(50, 50, page.rect.width - 50, page.rect.height - 50)

            # Handle multi-page content
            fontsize = 11
            lines_per_page = int((rect.height) / (fontsize * 1.5))
            lines = full_text.split("\n")

            current_page = page
            line_count = 0
            y_pos = rect.y0

            for line in lines:
                if line_count >= lines_per_page:
                    current_page = pdf_doc.new_page()
                    y_pos = rect.y0
                    line_count = 0

                current_page.insert_text(
                    (rect.x0, y_pos + fontsize),
                    line[:100],  # Truncate very long lines
                    fontsize=fontsize,
                    color=(0, 0, 0),
                    fontfile=UNICODE_FONT_PATH,
                    fontname="F1"
                )
                y_pos += fontsize * 1.5
                line_count += 1

            jobs[job_id]["progress"] = 80
            pdf_doc.save(str(output_path))
            pdf_doc.close()

            jobs[job_id]["status"] = "done"
            jobs[job_id]["progress"] = 100

            return {
                "job_id": job_id,
                "status": "done",
                "download_url": f"/api/download/{output_path.name}",
                "filename": output_path.name
            }
        except ImportError:
            raise HTTPException(status_code=500, detail="python-docx is required for DOCX conversion")
    except HTTPException:
        raise
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            saved_file.unlink()
        except:
            pass


# ============================================================================
# PDF to Excel API
# ============================================================================

@app.post("/api/pdf-to-excel")
async def pdf_to_excel(file: UploadFile = File(...)):
    """Convert PDF tables to Excel XLSX format using PyMuPDF table detection."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = save_upload_file(file)
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, Alignment, Border, Side

        doc = fitz.open(str(saved_file))
        wb = Workbook()
        ws = wb.active
        ws.title = "PDF Data"

        total_pages = doc.page_count
        current_row = 1
        tables_found = 0

        for i, page in enumerate(doc):
            # Try structured table detection first (PyMuPDF 1.23+)
            page_tables = page.find_tables()

            if page_tables and len(page_tables.tables) > 0:
                for table in page_tables.tables:
                    tables_found += 1
                    data = table.extract()

                    for row_idx, row in enumerate(data):
                        for col_idx, cell in enumerate(row):
                            cell_value = cell if cell else ""
                            cell_obj = ws.cell(row=current_row, column=col_idx + 1, value=cell_value)
                            # Bold header row (first row of each table)
                            if row_idx == 0:
                                cell_obj.font = Font(bold=True)
                            cell_obj.alignment = Alignment(wrap_text=True, vertical="top")
                        current_row += 1

                    current_row += 1  # Empty row between tables
            else:
                # Fallback: text-based extraction with smart column detection
                text = page.get_text()
                lines = text.strip().split('\n')

                for line in lines:
                    if line.strip():
                        # Try tab-separated first, then multi-space
                        cells = [cell.strip() for cell in line.split('\t') if cell.strip()]
                        if len(cells) == 1:
                            # Split by 2+ spaces (common in PDF table rendering)
                            cells = [cell.strip() for cell in re.split(r'\s{2,}', line) if cell.strip()]
                        if len(cells) == 1:
                            cells = [line.strip()]

                        for col_idx, cell in enumerate(cells, 1):
                            # Try to convert numeric strings
                            try:
                                if '.' in cell or ',' in cell:
                                    num_val = float(cell.replace(',', '.').replace(' ', ''))
                                    ws.cell(row=current_row, column=col_idx, value=num_val)
                                else:
                                    int_val = int(cell.replace(' ', ''))
                                    ws.cell(row=current_row, column=col_idx, value=int_val)
                            except (ValueError, AttributeError):
                                ws.cell(row=current_row, column=col_idx, value=cell)
                        current_row += 1

            if i < total_pages - 1:
                current_row += 1

            jobs[job_id]["progress"] = int((i + 1) / total_pages * 100)

        # Auto-fit column widths
        for column_cells in ws.columns:
            max_length = 0
            column = column_cells[0].column_letter
            for cell in column_cells:
                try:
                    if cell.value:
                        max_length = max(max_length, len(str(cell.value)))
                except:
                    pass
            ws.column_dimensions[column].width = min(max_length + 2, 50)

        doc.close()
        output_path = create_output_path("converted", ".xlsx")
        wb.save(str(output_path))

        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_path.name}",
            "filename": output_path.name,
            "pages": total_pages,
            "tables_detected": tables_found
        }
    except ImportError:
        raise HTTPException(status_code=500, detail="openpyxl is required for Excel conversion")
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            saved_file.unlink()
        except:
            pass


# ============================================================================
# Excel to PDF API
# ============================================================================

@app.post("/api/excel-to-pdf")
async def excel_to_pdf(file: UploadFile = File(...)):
    """Convert Excel XLSX to PDF format."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = save_upload_file(file)
    try:
        from openpyxl import load_workbook

        wb = load_workbook(str(saved_file))
        output_path = create_output_path("converted", ".pdf")

        pdf_doc = fitz.open()
        fontsize = 10
        margin = 50
        line_height = fontsize * 1.5

        for sheet_idx, sheet_name in enumerate(wb.sheetnames):
            ws = wb[sheet_name]
            page = pdf_doc.new_page()
            rect = page.rect

            y_pos = margin
            max_rows = int((rect.height - 2 * margin) / line_height)
            row_count = 0

            for row in ws.iter_rows(values_only=True):
                if row_count >= max_rows:
                    page = pdf_doc.new_page()
                    y_pos = margin
                    row_count = 0

                row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                page.insert_text(
                    (margin, y_pos + fontsize),
                    row_text[:200],
                    fontsize=fontsize,
                    color=(0, 0, 0),
                    fontfile=UNICODE_FONT_PATH,
                    fontname="F1"
                )
                y_pos += line_height
                row_count += 1

            jobs[job_id]["progress"] = int((sheet_idx + 1) / len(wb.sheetnames) * 100)

        pdf_doc.save(str(output_path))
        pdf_doc.close()

        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_path.name}",
            "filename": output_path.name,
            "sheets": len(wb.sheetnames)
        }
    except ImportError:
        raise HTTPException(status_code=500, detail="openpyxl is required for Excel conversion")
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            saved_file.unlink()
        except:
            pass


# ============================================================================
# PDF to PowerPoint API
# ============================================================================

@app.post("/api/pdf-to-powerpoint")
async def pdf_to_powerpoint(file: UploadFile = File(...)):
    """Convert PDF to PowerPoint PPTX format."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = save_upload_file(file)
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        doc = fitz.open(str(saved_file))
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        total_pages = doc.page_count

        for i, page in enumerate(doc):
            # Create slide with blank layout
            blank_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_layout)

            # Render page as image
            mat = fitz.Matrix(2, 2)  # 2x scale for better quality
            pix = page.get_pixmap(matrix=mat)
            img_path = UPLOAD_DIR / f"temp_slide_{i}_{uuid.uuid4().hex[:8]}.png"
            pix.save(str(img_path))

            # Add image to slide
            slide.shapes.add_picture(
                str(img_path),
                Inches(0),
                Inches(0),
                width=prs.slide_width,
                height=prs.slide_height
            )

            # Clean up temp image
            try:
                img_path.unlink()
            except:
                pass

            jobs[job_id]["progress"] = int((i + 1) / total_pages * 100)

        doc.close()
        output_path = create_output_path("converted", ".pptx")
        prs.save(str(output_path))

        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_path.name}",
            "filename": output_path.name,
            "slides": total_pages
        }
    except ImportError:
        raise HTTPException(status_code=500, detail="python-pptx is required for PowerPoint conversion")
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            saved_file.unlink()
        except:
            pass


# ============================================================================
# PowerPoint to PDF API
# ============================================================================

@app.post("/api/powerpoint-to-pdf")
async def powerpoint_to_pdf(file: UploadFile = File(...)):
    """Convert PowerPoint PPTX to PDF format with layout preservation."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = save_upload_file(file)
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from PIL import Image as PILImage, ImageDraw, ImageFont

        prs = Presentation(str(saved_file))
        output_path = create_output_path("converted", ".pdf")

        # Get slide dimensions
        slide_width = prs.slide_width or Inches(13.333)
        slide_height = prs.slide_height or Inches(7.5)
        pdf_width = 792  # Standard landscape
        pdf_height = int(pdf_width * (slide_height / slide_width))

        pdf_doc = fitz.open()
        total_slides = len(prs.slides)

        for i, slide in enumerate(prs.slides):
            # Render slide as image for accurate layout
            img_width, img_height = 1920, int(1920 * (slide_height / slide_width))
            img = PILImage.new("RGB", (img_width, img_height), "white")
            draw = ImageDraw.Draw(img)

            # Try to get a font for text rendering
            font_path = UNICODE_FONT_PATH or "arial.ttf"
            scale_x = img_width / (slide_width / 914400)  # EMU to inches to pixels
            scale_y = img_height / (slide_height / 914400)

            # Draw background if slide has a fill
            if hasattr(slide, 'background') and slide.background.fill:
                try:
                    bg = slide.background.fill
                    if bg.type is not None and hasattr(bg, 'fore_color') and bg.fore_color:
                        rgb = bg.fore_color.rgb
                        draw.rectangle([0, 0, img_width, img_height], fill=f"#{rgb}")
                except:
                    pass

            # Render shapes with positions
            for shape in sorted(slide.shapes, key=lambda s: (s.top or 0)):
                if not hasattr(shape, 'left') or shape.left is None:
                    continue

                # Calculate position in pixels
                x = int((shape.left / slide_width) * img_width)
                y = int((shape.top / slide_height) * img_height)
                w = int((shape.width / slide_width) * img_width) if shape.width else img_width - x
                h = int((shape.height / slide_height) * img_height) if shape.height else 50

                # Handle images
                if shape.shape_type == 13 and hasattr(shape, 'image'):  # Picture
                    try:
                        img_bytes = shape.image.blob
                        shape_img = PILImage.open(io.BytesIO(img_bytes))
                        shape_img = shape_img.convert("RGBA")
                        shape_img = shape_img.resize((max(w, 1), max(h, 1)), PILImage.LANCZOS)
                        img.paste(shape_img, (x, y), shape_img if shape_img.mode == "RGBA" else None)
                    except:
                        pass

                # Handle text
                elif hasattr(shape, "text_frame"):
                    try:
                        tf = shape.text_frame
                        text_y = y
                        for para in tf.paragraphs:
                            if not para.text.strip():
                                text_y += 20
                                continue
                            # Determine font size
                            fs = 18  # default
                            font_color = (0, 0, 0)
                            is_bold = False
                            for run in para.runs:
                                if run.font.size:
                                    fs = int(run.font.size / Pt(1))
                                if run.font.bold:
                                    is_bold = True
                                if run.font.color and run.font.color.rgb:
                                    rgb = str(run.font.color.rgb)
                                    font_color = (int(rgb[0:2], 16), int(rgb[2:4], 16), int(rgb[4:6], 16))

                            # Scale font size
                            scaled_fs = max(int(fs * img_width / 1920 * 1.5), 10)
                            try:
                                font = ImageFont.truetype(font_path, scaled_fs)
                            except:
                                font = ImageFont.load_default()

                            draw.text((x + 5, text_y), para.text, fill=font_color, font=font)
                            text_y += scaled_fs + 4
                    except:
                        pass

                # Handle tables
                elif shape.has_table:
                    try:
                        table = shape.table
                        rows = len(table.rows)
                        cols = len(table.columns)
                        cell_w = w // max(cols, 1)
                        cell_h = h // max(rows, 1)

                        for r_idx, row in enumerate(table.rows):
                            for c_idx, cell in enumerate(row.cells):
                                cx = x + c_idx * cell_w
                                cy = y + r_idx * cell_h
                                draw.rectangle([cx, cy, cx + cell_w, cy + cell_h], outline=(180, 180, 180))
                                try:
                                    font = ImageFont.truetype(font_path, max(int(12 * img_width / 1920 * 1.5), 8))
                                except:
                                    font = ImageFont.load_default()
                                draw.text((cx + 3, cy + 3), cell.text[:50], fill=(0, 0, 0), font=font)
                    except:
                        pass

            # Convert PIL image to PDF page
            img_bytes_io = io.BytesIO()
            img.save(img_bytes_io, format="PNG")
            img_bytes_io.seek(0)

            page = pdf_doc.new_page(width=pdf_width, height=pdf_height)
            page.insert_image(page.rect, stream=img_bytes_io.read())

            jobs[job_id]["progress"] = int((i + 1) / total_slides * 100)

        pdf_doc.save(str(output_path))
        pdf_doc.close()

        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_path.name}",
            "filename": output_path.name,
            "pages": total_slides
        }
    except ImportError as e:
        raise HTTPException(status_code=500, detail=f"python-pptx and Pillow are required: {e}")
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            saved_file.unlink()
        except:
            pass


# ============================================================================
# PNG to PDF API
# ============================================================================

@app.post("/api/png-to-pdf")
async def png_to_pdf(
    files: List[UploadFile] = File(...),
    order: str = Form(None),
    page_size: str = Form("a4")
):
    """Convert PNG images to PDF."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_files = []
    try:
        for f in files:
            saved_files.append(save_upload_file(f))

        if order:
            indices = [int(i) for i in order.split(",")]
            saved_files = [saved_files[i] for i in indices if i < len(saved_files)]

        page_sizes = {
            "a4": (595, 842),
            "letter": (612, 792),
            "original": None
        }
        target_size = page_sizes.get(page_size, page_sizes["a4"])

        output_path = create_output_path("images_to_pdf")
        doc = fitz.open()

        total = len(saved_files)
        for i, img_path in enumerate(saved_files):
            img = fitz.open(str(img_path))

            if target_size:
                rect = fitz.Rect(0, 0, target_size[0], target_size[1])
            else:
                pdfbytes = img.convert_to_pdf()
                img_pdf = fitz.open("pdf", pdfbytes)
                rect = img_pdf[0].rect

            pdfbytes = img.convert_to_pdf()
            img_pdf = fitz.open("pdf", pdfbytes)
            page = doc.new_page(width=rect.width, height=rect.height)
            page.show_pdf_page(rect, img_pdf, 0)

            img.close()
            img_pdf.close()
            jobs[job_id]["progress"] = int((i + 1) / total * 100)

        doc.save(str(output_path))
        doc.close()

        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_path.name}",
            "filename": output_path.name,
            "pages": len(saved_files)
        }
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        for f in saved_files:
            try:
                f.unlink()
            except:
                pass


# ============================================================================
# PDF to PNG API
# ============================================================================

@app.post("/api/pdf-to-png")
async def pdf_to_png(
    file: UploadFile = File(...),
    dpi: int = Form(150),
    pages: str = Form("all")
):
    """Convert PDF pages to PNG images."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = save_upload_file(file)
    try:
        doc = fitz.open(str(saved_file))
        total_pages = doc.page_count
        output_files = []

        if pages == "all":
            page_list = list(range(total_pages))
        else:
            page_list = []
            for part in pages.split(","):
                part = part.strip()
                if "-" in part:
                    start, end = map(int, part.split("-"))
                    page_list.extend(range(start - 1, end))
                else:
                    page_list.append(int(part) - 1)
            page_list = [p for p in page_list if 0 <= p < total_pages]

        mat = fitz.Matrix(dpi / 72, dpi / 72)
        for idx, page_num in enumerate(page_list):
            page = doc[page_num]
            pix = page.get_pixmap(matrix=mat)
            out_path = OUTPUT_DIR / f"page_{page_num+1}_{uuid.uuid4().hex[:8]}.png"
            pix.save(str(out_path))
            output_files.append(out_path.name)
            jobs[job_id]["progress"] = int((idx + 1) / len(page_list) * 100)

        doc.close()
        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100
        jobs[job_id]["output_files"] = output_files

        if len(output_files) > 1:
            zip_path = create_output_path("images", ".zip")
            with zipfile.ZipFile(str(zip_path), "w", zipfile.ZIP_DEFLATED) as zf:
                for fname in output_files:
                    zf.write(OUTPUT_DIR / fname, fname)

            return {
                "job_id": job_id,
                "status": "done",
                "download_url": f"/api/download/{zip_path.name}",
                "filename": zip_path.name,
                "files": [{"filename": f, "download_url": f"/api/download/{f}"} for f in output_files],
                "total_images": len(output_files)
            }
        else:
            return {
                "job_id": job_id,
                "status": "done",
                "download_url": f"/api/download/{output_files[0]}",
                "filename": output_files[0],
                "total_images": 1
            }
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            saved_file.unlink()
        except:
            pass


# ============================================================================
# PDF to Text API
# ============================================================================

@app.post("/api/pdf-to-text")
async def pdf_to_text(file: UploadFile = File(...)):
    """Extract text from PDF."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = save_upload_file(file)
    try:
        doc = fitz.open(str(saved_file))
        total_pages = doc.page_count

        output_path = create_output_path("extracted", ".txt")

        with open(output_path, "w", encoding="utf-8") as f:
            for i, page in enumerate(doc):
                text = page.get_text()
                f.write(f"--- Page {i + 1} ---\n\n")
                f.write(text)
                f.write("\n\n")
                jobs[job_id]["progress"] = int((i + 1) / total_pages * 100)

        doc.close()

        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_path.name}",
            "filename": output_path.name,
            "pages": total_pages
        }
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            saved_file.unlink()
        except:
            pass


# ============================================================================
# Text to PDF API
# ============================================================================

@app.post("/api/text-to-pdf")
async def text_to_pdf(
    file: UploadFile = File(None),
    text: str = Form(None),
    font_size: int = Form(12)
):
    """Create PDF from text input or text file."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = None
    try:
        content = ""

        if file and file.filename:
            saved_file = save_upload_file(file)
            with open(saved_file, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
        elif text:
            content = text
        else:
            raise HTTPException(status_code=400, detail="Either file or text must be provided")

        output_path = create_output_path("text_to_pdf", ".pdf")
        pdf_doc = fitz.open()

        # Page settings
        page_width = 595  # A4 width
        page_height = 842  # A4 height
        margin = 50
        line_height = font_size * 1.5
        max_lines = int((page_height - 2 * margin) / line_height)
        max_chars = int((page_width - 2 * margin) / (font_size * 0.5))

        # Split content into lines
        lines = content.split('\n')
        wrapped_lines = []
        for line in lines:
            if len(line) <= max_chars:
                wrapped_lines.append(line)
            else:
                # Word wrap
                words = line.split(' ')
                current_line = ""
                for word in words:
                    if len(current_line) + len(word) + 1 <= max_chars:
                        current_line += (" " if current_line else "") + word
                    else:
                        if current_line:
                            wrapped_lines.append(current_line)
                        current_line = word
                if current_line:
                    wrapped_lines.append(current_line)

        # Create pages
        total_lines = len(wrapped_lines)
        page = None
        y_pos = margin

        for idx, line in enumerate(wrapped_lines):
            if page is None or y_pos + line_height > page_height - margin:
                page = pdf_doc.new_page(width=page_width, height=page_height)
                y_pos = margin

            page.insert_text(
                (margin, y_pos + font_size),
                line,
                fontsize=font_size,
                color=(0, 0, 0),
                fontfile=UNICODE_FONT_PATH,
                fontname="F1"
            )
            y_pos += line_height
            jobs[job_id]["progress"] = int((idx + 1) / total_lines * 100) if total_lines > 0 else 100

        if pdf_doc.page_count == 0:
            pdf_doc.new_page()

        # Get page count before closing document
        page_count = pdf_doc.page_count

        pdf_doc.save(str(output_path))
        pdf_doc.close()

        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_path.name}",
            "filename": output_path.name,
            "pages": page_count
        }
    except HTTPException:
        raise
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        if saved_file:
            try:
                saved_file.unlink()
            except:
                pass


@app.post("/api/html-to-pdf")
async def html_to_pdf(
    file: UploadFile = File(None),
    html_content: str = Form(None)
):
    """Convert HTML to PDF."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = None
    try:
        if file and file.filename:
            saved_file = save_upload_file(file)
            with open(saved_file, "r", encoding="utf-8", errors="ignore") as f:
                html_text = f.read()
        elif html_content:
            html_text = html_content
        else:
            raise HTTPException(status_code=400, detail="No HTML content provided")

        output_path = create_output_path("html_converted", ".pdf")

        # Use fitz to create PDF from HTML
        doc = fitz.open()
        page = doc.new_page()
        # Insert HTML as text (basic conversion)
        rect = page.rect + fitz.Rect(36, 36, -36, -36)  # margins
        page.insert_htmlbox(rect, html_text)

        # Get page count before closing document
        page_count = doc.page_count

        doc.save(str(output_path))
        doc.close()

        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_path.name}",
            "filename": output_path.name,
            "pages": page_count
        }
    except HTTPException:
        raise
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        if saved_file:
            try:
                saved_file.unlink()
            except:
                pass


# ============================================================================
# Job Status Polling
# ============================================================================

@app.get("/api/job/{job_id}")
def get_job_status(job_id: str):
    """Get job status for progress polling."""
    if job_id not in jobs:
        raise HTTPException(status_code=404, detail="Job not found")

    job = jobs[job_id]
    return {
        "id": job["id"],
        "status": job["status"],
        "progress": job.get("progress", 0),
        "error": job.get("error"),
        "output_file": job.get("output_file"),
        "output_files": job.get("output_files", [])
    }


# ============================================================================
# Download endpoint for output files
# ============================================================================

@app.get("/api/download/{filename}")
def download_output(filename: str):
    """Download processed file from output directory."""
    file_path = OUTPUT_DIR / filename
    if not file_path.exists():
        # Also check in the old output directory
        old_path = Path("output") / filename
        if old_path.exists():
            file_path = old_path
        else:
            raise HTTPException(status_code=404, detail="File not found")

    # Determine media type
    suffix = file_path.suffix.lower()
    media_types = {
        ".pdf": "application/pdf",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".txt": "text/plain",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".png": "image/png",
        ".zip": "application/zip"
    }
    media_type = media_types.get(suffix, "application/octet-stream")

    return FileResponse(str(file_path), filename=filename, media_type=media_type)


# ============================================================================
# Rotate PDF API
# ============================================================================

@app.post("/api/rotate")
async def rotate_pdf(
    file: UploadFile = File(...),
    rotation: int = Form(90),
    pages: str = Form("all")
):
    """Rotate PDF pages by specified degrees (90, 180, 270)."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = save_upload_file(file)
    try:
        doc = fitz.open(str(saved_file))
        total_pages = doc.page_count

        if pages == "all":
            page_list = list(range(total_pages))
        else:
            page_list = []
            for part in pages.split(","):
                part = part.strip()
                if "-" in part:
                    start, end = map(int, part.split("-"))
                    page_list.extend(range(start - 1, end))
                else:
                    page_list.append(int(part) - 1)
            page_list = [p for p in page_list if 0 <= p < total_pages]

        for idx, page_num in enumerate(page_list):
            page = doc[page_num]
            page.set_rotation((page.rotation + rotation) % 360)
            jobs[job_id]["progress"] = int((idx + 1) / len(page_list) * 100)

        output_path = create_output_path("rotated")
        doc.save(str(output_path))
        doc.close()

        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_path.name}",
            "filename": output_path.name,
            "pages_rotated": len(page_list)
        }
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            saved_file.unlink()
        except:
            pass


# ============================================================================
# Delete Pages API
# ============================================================================

@app.post("/api/delete-pages")
async def delete_pages(
    file: UploadFile = File(...),
    pages: str = Form(...)
):
    """Delete specific pages from PDF."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = save_upload_file(file)
    try:
        doc = fitz.open(str(saved_file))
        total_pages = doc.page_count

        pages_to_delete = set()
        for part in pages.split(","):
            part = part.strip()
            if "-" in part:
                start, end = map(int, part.split("-"))
                pages_to_delete.update(range(start - 1, end))
            else:
                pages_to_delete.add(int(part) - 1)

        pages_to_delete = sorted([p for p in pages_to_delete if 0 <= p < total_pages], reverse=True)

        if len(pages_to_delete) >= total_pages:
            raise HTTPException(status_code=400, detail="Cannot delete all pages")

        for idx, page_num in enumerate(pages_to_delete):
            doc.delete_page(page_num)
            jobs[job_id]["progress"] = int((idx + 1) / len(pages_to_delete) * 100)

        output_path = create_output_path("pages_deleted")
        doc.save(str(output_path))
        doc.close()

        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_path.name}",
            "filename": output_path.name,
            "pages_deleted": len(pages_to_delete),
            "pages_remaining": total_pages - len(pages_to_delete)
        }
    except HTTPException:
        raise
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            saved_file.unlink()
        except:
            pass


# ============================================================================
# Extract Pages API
# ============================================================================

@app.post("/api/extract-pages")
async def extract_pages(
    file: UploadFile = File(...),
    pages: str = Form(...)
):
    """Extract specific pages as a new PDF."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = save_upload_file(file)
    try:
        doc = fitz.open(str(saved_file))
        total_pages = doc.page_count

        pages_to_extract = []
        for part in pages.split(","):
            part = part.strip()
            if "-" in part:
                start, end = map(int, part.split("-"))
                pages_to_extract.extend(range(start - 1, end))
            else:
                pages_to_extract.append(int(part) - 1)

        pages_to_extract = [p for p in pages_to_extract if 0 <= p < total_pages]

        if not pages_to_extract:
            raise HTTPException(status_code=400, detail="No valid pages to extract")

        new_doc = fitz.open()
        for idx, page_num in enumerate(pages_to_extract):
            new_doc.insert_pdf(doc, from_page=page_num, to_page=page_num)
            jobs[job_id]["progress"] = int((idx + 1) / len(pages_to_extract) * 100)

        output_path = create_output_path("extracted_pages")
        new_doc.save(str(output_path))
        new_doc.close()
        doc.close()

        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_path.name}",
            "filename": output_path.name,
            "pages_extracted": len(pages_to_extract)
        }
    except HTTPException:
        raise
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            saved_file.unlink()
        except:
            pass


# ============================================================================
# Crop PDF API
# ============================================================================

@app.post("/api/crop")
async def crop_pdf(
    file: UploadFile = File(...),
    top: float = Form(0),
    bottom: float = Form(0),
    left: float = Form(0),
    right: float = Form(0),
    pages: str = Form("all")
):
    """Crop PDF page margins (in points, 72 points = 1 inch)."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = save_upload_file(file)
    try:
        doc = fitz.open(str(saved_file))
        total_pages = doc.page_count

        if pages == "all":
            page_list = list(range(total_pages))
        else:
            page_list = []
            for part in pages.split(","):
                part = part.strip()
                if "-" in part:
                    start, end = map(int, part.split("-"))
                    page_list.extend(range(start - 1, end))
                else:
                    page_list.append(int(part) - 1)
            page_list = [p for p in page_list if 0 <= p < total_pages]

        for idx, page_num in enumerate(page_list):
            page = doc[page_num]
            rect = page.rect
            new_rect = fitz.Rect(
                rect.x0 + left,
                rect.y0 + top,
                rect.x1 - right,
                rect.y1 - bottom
            )
            page.set_cropbox(new_rect)
            jobs[job_id]["progress"] = int((idx + 1) / len(page_list) * 100)

        output_path = create_output_path("cropped")
        doc.save(str(output_path))
        doc.close()

        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_path.name}",
            "filename": output_path.name,
            "pages_cropped": len(page_list)
        }
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            saved_file.unlink()
        except:
            pass


# ============================================================================
# Watermark PDF API
# ============================================================================

@app.post("/api/watermark")
async def watermark_pdf(
    file: UploadFile = File(...),
    text: str = Form(None),
    watermark_image: UploadFile = File(None),
    position: str = Form("center"),
    opacity: float = Form(0.3),
    rotation: int = Form(45),
    font_size: int = Form(48),
    pages: str = Form("all")
):
    """Add text or image watermark to PDF."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = save_upload_file(file)
    watermark_file = None
    try:
        if not text and not (watermark_image and watermark_image.filename):
            raise HTTPException(status_code=400, detail="Either text or watermark image must be provided")

        doc = fitz.open(str(saved_file))
        total_pages = doc.page_count

        if pages == "all":
            page_list = list(range(total_pages))
        else:
            page_list = []
            for part in pages.split(","):
                part = part.strip()
                if "-" in part:
                    start, end = map(int, part.split("-"))
                    page_list.extend(range(start - 1, end))
                else:
                    page_list.append(int(part) - 1)
            page_list = [p for p in page_list if 0 <= p < total_pages]

        for idx, page_num in enumerate(page_list):
            page = doc[page_num]
            rect = page.rect

            if text:
                positions_map = {
                    "center": fitz.Point(rect.width / 2, rect.height / 2),
                    "top-left": fitz.Point(100, 100),
                    "top-right": fitz.Point(rect.width - 100, 100),
                    "bottom-left": fitz.Point(100, rect.height - 100),
                    "bottom-right": fitz.Point(rect.width - 100, rect.height - 100),
                }
                center = positions_map.get(position, positions_map["center"])

                # Create overlay page for transparency support
                overlay_doc = fitz.open()
                overlay_page = overlay_doc.new_page(width=rect.width, height=rect.height)

                # Calculate text position centered on the point
                text_length = fitz.get_text_length(text, fontsize=font_size)
                insert_point = fitz.Point(
                    center.x - text_length / 2,
                    center.y + font_size / 3
                )

                # Apply rotation via morph if needed
                if rotation != 0:
                    morph = (center, fitz.Matrix(rotation))
                    overlay_page.insert_text(
                        insert_point,
                        text,
                        fontsize=font_size,
                        color=(0.5, 0.5, 0.5),
                        fontfile=UNICODE_FONT_PATH,
                        fontname="F1",
                        morph=morph
                    )
                else:
                    overlay_page.insert_text(
                        insert_point,
                        text,
                        fontsize=font_size,
                        color=(0.5, 0.5, 0.5),
                        fontfile=UNICODE_FONT_PATH,
                        fontname="F1"
                    )

                # Merge overlay onto original page with opacity
                page.show_pdf_page(rect, overlay_doc, 0, overlay=(True, None, opacity))
                overlay_doc.close()

            elif watermark_image and watermark_image.filename:
                if not watermark_file:
                    watermark_file = save_upload_file(watermark_image)
                img_rect = fitz.Rect(0, 0, rect.width, rect.height)
                page.insert_image(img_rect, filename=str(watermark_file), overlay=True, keep_proportion=True)

            jobs[job_id]["progress"] = int((idx + 1) / len(page_list) * 100)

        output_path = create_output_path("watermarked")
        doc.save(str(output_path))
        doc.close()

        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_path.name}",
            "filename": output_path.name,
            "pages_watermarked": len(page_list)
        }
    except HTTPException:
        raise
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            saved_file.unlink()
        except:
            pass
        if watermark_file:
            try:
                watermark_file.unlink()
            except:
                pass


# ============================================================================
# Add Page Numbers API
# ============================================================================

@app.post("/api/add-page-numbers")
async def add_page_numbers(
    file: UploadFile = File(...),
    position: str = Form("bottom-center"),
    format_str: str = Form("Page {n} of {total}"),
    font_size: int = Form(12),
    start_number: int = Form(1),
    margin: int = Form(30)
):
    """Add page numbers to PDF."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = save_upload_file(file)
    try:
        doc = fitz.open(str(saved_file))
        total_pages = doc.page_count

        for i, page in enumerate(doc):
            page_num = start_number + i
            text = format_str.replace("{n}", str(page_num)).replace("{total}", str(total_pages))

            rect = page.rect
            text_length = fitz.get_text_length(text, fontsize=font_size)

            positions = {
                "top-left": (margin, margin + font_size),
                "top-center": ((rect.width - text_length) / 2, margin + font_size),
                "top-right": (rect.width - text_length - margin, margin + font_size),
                "bottom-left": (margin, rect.height - margin),
                "bottom-center": ((rect.width - text_length) / 2, rect.height - margin),
                "bottom-right": (rect.width - text_length - margin, rect.height - margin),
            }
            pos = positions.get(position, positions["bottom-center"])

            page.insert_text(
                fitz.Point(pos[0], pos[1]),
                text,
                fontsize=font_size,
                color=(0, 0, 0),
                fontfile=UNICODE_FONT_PATH,
                fontname="F1"
            )
            jobs[job_id]["progress"] = int((i + 1) / total_pages * 100)

        output_path = create_output_path("numbered")
        doc.save(str(output_path))
        doc.close()

        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_path.name}",
            "filename": output_path.name,
            "pages_numbered": total_pages
        }
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            saved_file.unlink()
        except:
            pass


# ============================================================================
# Protect PDF API
# ============================================================================

@app.post("/api/protect")
async def protect_pdf(
    file: UploadFile = File(...),
    user_password: str = Form(...),
    owner_password: str = Form(None),
    allow_print: str = Form("true"),
    allow_copy: str = Form("true"),
    allow_modify: str = Form("false")
):
    """Add password protection to PDF."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    # Parse bool strings from FormData
    _allow_print = allow_print.lower() in ("true", "1", "yes")
    _allow_copy = allow_copy.lower() in ("true", "1", "yes")
    _allow_modify = allow_modify.lower() in ("true", "1", "yes")

    saved_file = save_upload_file(file)
    try:
        doc = fitz.open(str(saved_file))

        owner_pwd = owner_password if owner_password else user_password

        permissions = fitz.PDF_PERM_ACCESSIBILITY
        if _allow_print:
            permissions |= fitz.PDF_PERM_PRINT | fitz.PDF_PERM_PRINT_HQ
        if _allow_copy:
            permissions |= fitz.PDF_PERM_COPY
        if _allow_modify:
            permissions |= fitz.PDF_PERM_MODIFY | fitz.PDF_PERM_ANNOTATE | fitz.PDF_PERM_FORM

        output_path = create_output_path("protected")
        doc.save(
            str(output_path),
            encryption=fitz.PDF_ENCRYPT_AES_256,
            user_pw=user_password,
            owner_pw=owner_pwd,
            permissions=permissions
        )
        doc.close()

        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_path.name}",
            "filename": output_path.name,
            "encrypted": True
        }
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            saved_file.unlink()
        except:
            pass


# ============================================================================
# Unlock PDF API
# ============================================================================

@app.post("/api/unlock")
async def unlock_pdf(
    file: UploadFile = File(...),
    password: str = Form(...)
):
    """Remove password protection from PDF."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = save_upload_file(file)
    try:
        doc = fitz.open(str(saved_file))

        if doc.is_encrypted:
            if not doc.authenticate(password):
                raise HTTPException(status_code=400, detail="Incorrect password")

        output_path = create_output_path("unlocked")
        doc.save(str(output_path))
        doc.close()

        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_path.name}",
            "filename": output_path.name,
            "unlocked": True
        }
    except HTTPException:
        raise
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            saved_file.unlink()
        except:
            pass


# ============================================================================
# Edit PDF API - Add text annotations
# ============================================================================

@app.post("/api/edit-pdf")
async def edit_pdf(
    file: UploadFile = File(...),
    text: str = Form(...),
    page_number: int = Form(1),
    position_x: float = Form(300),
    position_y: float = Form(420),
    font_size: int = Form(14),
    color: str = Form("0,0,0")
):
    """Add text annotation to a specific page of a PDF."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = save_upload_file(file)
    try:
        if not text.strip():
            raise HTTPException(status_code=400, detail="Text must not be empty")

        doc = fitz.open(str(saved_file))
        total_pages = doc.page_count

        # Convert 1-based page number to 0-based index
        page_idx = page_number - 1
        if page_idx < 0 or page_idx >= total_pages:
            raise HTTPException(
                status_code=400,
                detail=f"Page number {page_number} is out of range. PDF has {total_pages} page(s)."
            )

        # Parse color from "R,G,B" string (0-255 range)
        try:
            color_parts = [int(c.strip()) for c in color.split(",")]
            if len(color_parts) != 3:
                raise ValueError
            text_color = (color_parts[0] / 255.0, color_parts[1] / 255.0, color_parts[2] / 255.0)
        except (ValueError, IndexError):
            text_color = (0, 0, 0)

        jobs[job_id]["progress"] = 20

        page = doc[page_idx]
        point = fitz.Point(position_x, position_y)

        # Insert text using Unicode font for diacritics support
        page.insert_text(
            point,
            text,
            fontsize=font_size,
            fontname="helv",
            fontfile=UNICODE_FONT_PATH,
            color=text_color,
        )

        jobs[job_id]["progress"] = 70

        output_path = create_output_path("edited")
        doc.save(str(output_path), deflate=True, garbage=3)
        doc.close()

        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_path.name}",
            "filename": output_path.name,
            "page_edited": page_number,
        }
    except HTTPException:
        raise
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            saved_file.unlink()
        except:
            pass


# ============================================================================
# Annotate PDF API
# ============================================================================

@app.post("/api/annotate-pdf")
async def annotate_pdf(
    file: UploadFile = File(...),
    annotation_text: str = Form(...),
    pages: str = Form("all"),
    color: str = Form("yellow"),
):
    """Search for text in PDF and add highlight annotations."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = save_upload_file(file)
    try:
        doc = fitz.open(str(saved_file))
        total_pages = doc.page_count

        # Map color names to RGB tuples
        color_map = {
            "yellow": (1, 1, 0),
            "green": (0, 1, 0),
            "blue": (0, 0.5, 1),
            "pink": (1, 0.5, 0.75),
        }
        highlight_color = color_map.get(color.lower(), (1, 1, 0))

        # Determine which pages to process
        if pages.strip().lower() == "all":
            page_indices = list(range(total_pages))
        else:
            page_indices = []
            for part in pages.split(","):
                part = part.strip()
                if "-" in part:
                    start, end = part.split("-", 1)
                    for p in range(int(start) - 1, min(int(end), total_pages)):
                        if 0 <= p < total_pages:
                            page_indices.append(p)
                else:
                    p = int(part) - 1
                    if 0 <= p < total_pages:
                        page_indices.append(p)

        annotations_added = 0
        for i, page_idx in enumerate(page_indices):
            page = doc[page_idx]
            text_instances = page.search_for(annotation_text)
            for inst in text_instances:
                annot = page.add_highlight_annot(inst)
                annot.set_colors(stroke=highlight_color)
                annot.update()
                annotations_added += 1
            jobs[job_id]["progress"] = int((i + 1) / len(page_indices) * 80)

        output_path = create_output_path("annotated")
        doc.save(str(output_path), deflate=True, garbage=3)
        doc.close()

        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_path.name}",
            "filename": output_path.name,
            "annotations_added": annotations_added,
        }
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            saved_file.unlink()
        except:
            pass


# ============================================================================
# Redact PDF API
# ============================================================================

@app.post("/api/redact-pdf")
async def redact_pdf(
    file: UploadFile = File(...),
    redact_text: str = Form(...),
    pages: str = Form("all"),
):
    """Search for text in PDF and permanently redact (black out) it."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = save_upload_file(file)
    try:
        doc = fitz.open(str(saved_file))
        total_pages = doc.page_count

        # Determine which pages to process
        if pages.strip().lower() == "all":
            page_indices = list(range(total_pages))
        else:
            page_indices = []
            for part in pages.split(","):
                part = part.strip()
                if "-" in part:
                    start, end = part.split("-", 1)
                    for p in range(int(start) - 1, min(int(end), total_pages)):
                        if 0 <= p < total_pages:
                            page_indices.append(p)
                else:
                    p = int(part) - 1
                    if 0 <= p < total_pages:
                        page_indices.append(p)

        redactions_applied = 0
        for i, page_idx in enumerate(page_indices):
            page = doc[page_idx]
            text_instances = page.search_for(redact_text)
            for inst in text_instances:
                page.add_redact_annot(inst, fill=(0, 0, 0))
                redactions_applied += 1
            page.apply_redactions()
            jobs[job_id]["progress"] = int((i + 1) / len(page_indices) * 80)

        output_path = create_output_path("redacted")
        doc.save(str(output_path), deflate=True, garbage=3)
        doc.close()

        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_path.name}",
            "filename": output_path.name,
            "redactions_applied": redactions_applied,
        }
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            saved_file.unlink()
        except:
            pass


# ============================================================================
# Organize PDF API
# ============================================================================

@app.post("/api/organize-pdf")
async def organize_pdf(
    file: UploadFile = File(...),
    page_order: str = Form(""),
    reverse: bool = Form(False),
):
    """Reorder or reverse pages in a PDF document."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = save_upload_file(file)
    try:
        doc = fitz.open(str(saved_file))
        total_pages = doc.page_count

        if reverse:
            # Reverse all pages
            new_order = list(range(total_pages - 1, -1, -1))
        elif page_order.strip():
            # Parse custom page order (1-based input, convert to 0-based)
            new_order = []
            for part in page_order.split(","):
                part = part.strip()
                if part:
                    p = int(part) - 1
                    if p < 0 or p >= total_pages:
                        raise HTTPException(
                            status_code=400,
                            detail=f"Invalid page number {int(part)}. PDF has {total_pages} pages."
                        )
                    new_order.append(p)
            if not new_order:
                raise HTTPException(status_code=400, detail="No valid page numbers provided")
        else:
            raise HTTPException(status_code=400, detail="Please provide a page order or set reverse to true")

        jobs[job_id]["progress"] = 50

        doc.select(new_order)

        output_path = create_output_path("organized")
        doc.save(str(output_path), deflate=True, garbage=3)
        doc.close()

        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_path.name}",
            "filename": output_path.name,
            "total_pages": len(new_order),
        }
    except HTTPException:
        raise
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            saved_file.unlink()
        except:
            pass


# ============================================================================
# Flatten PDF API
# ============================================================================

@app.post("/api/flatten")
async def flatten_pdf(file: UploadFile = File(...)):
    """Flatten PDF form fields and annotations."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = save_upload_file(file)
    try:
        doc = fitz.open(str(saved_file))
        total_pages = doc.page_count

        for i, page in enumerate(doc):
            page.clean_contents()
            annots = list(page.annots()) if page.annots() else []
            for annot in annots:
                annot.set_flags(fitz.PDF_ANNOT_IS_HIDDEN)
            jobs[job_id]["progress"] = int((i + 1) / total_pages * 50)

        output_path = create_output_path("flattened")

        doc.save(str(output_path), deflate=True, garbage=3)
        doc.close()
        jobs[job_id]["progress"] = 100

        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_path.name}",
            "filename": output_path.name,
            "pages_flattened": total_pages
        }
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            saved_file.unlink()
        except:
            pass


# ============================================================================
# Repair PDF API
# ============================================================================

@app.post("/api/repair")
async def repair_pdf(file: UploadFile = File(...)):
    """Attempt to repair corrupted PDF."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = save_upload_file(file)
    try:
        doc = fitz.open(str(saved_file))
        jobs[job_id]["progress"] = 30

        output_path = create_output_path("repaired")
        doc.save(
            str(output_path),
            garbage=4,
            clean=True,
            deflate=True
        )
        jobs[job_id]["progress"] = 90

        repaired_doc = fitz.open(str(output_path))
        page_count = repaired_doc.page_count
        repaired_doc.close()
        doc.close()

        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_path.name}",
            "filename": output_path.name,
            "pages": page_count,
            "repaired": True
        }
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            saved_file.unlink()
        except:
            pass


# ============================================================================
# Split by Text Pattern API
# ============================================================================

@app.post("/api/split-by-text")
async def split_by_text(
    file: UploadFile = File(...),
    pattern: str = Form(...),
    include_match_in: str = Form("next"),  # "next", "previous", "separate"
    dpi: int = Form(150)
):
    """Split PDF when specific text pattern is found via OCR."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = save_upload_file(file)
    try:
        ocr = get_ocr()
        doc = fitz.open(str(saved_file))
        total_pages = doc.page_count

        split_points = []
        for i in range(total_pages):
            page = doc[i]
            # Try native text first (faster, preserves diacritics)
            full_text = page.get_text("text")

            # Fallback to OCR only if page has very little native text (scanned)
            if len(full_text.strip()) < 10:
                mat = fitz.Matrix(dpi / 72, dpi / 72)
                pix = page.get_pixmap(matrix=mat)
                img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, 3)

                result, _ = ocr(img)
                texts = [r[1] for r in result] if result else []
                full_text = ' '.join(texts)

            if re.search(pattern, full_text, re.IGNORECASE):
                split_points.append(i)

            jobs[job_id]["progress"] = int((i + 1) / total_pages * 70)

        if not split_points:
            raise HTTPException(status_code=400, detail="Pattern not found in any page")

        output_files = []
        if include_match_in == "separate":
            ranges = []
            prev = 0
            for sp in split_points:
                if sp > prev:
                    ranges.append((prev, sp - 1))
                ranges.append((sp, sp))
                prev = sp + 1
            if prev < total_pages:
                ranges.append((prev, total_pages - 1))
        elif include_match_in == "previous":
            ranges = []
            prev = 0
            for sp in split_points:
                ranges.append((prev, sp))
                prev = sp + 1
            if prev < total_pages:
                ranges.append((prev, total_pages - 1))
        else:
            ranges = []
            prev = 0
            for sp in split_points:
                if sp > prev:
                    ranges.append((prev, sp - 1))
                prev = sp
            ranges.append((prev, total_pages - 1))

        for idx, (start, end) in enumerate(ranges):
            out_path = create_output_path(f"split_text_{idx + 1}")
            new_doc = fitz.open()
            new_doc.insert_pdf(doc, from_page=start, to_page=end)
            new_doc.save(str(out_path))
            new_doc.close()
            output_files.append(out_path.name)
            jobs[job_id]["progress"] = 70 + int((idx + 1) / len(ranges) * 30)

        doc.close()
        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100

        return {
            "job_id": job_id,
            "status": "done",
            "files": [{"filename": f, "download_url": f"/api/download/{f}"} for f in output_files],
            "total_files": len(output_files),
            "split_points": [p + 1 for p in split_points]
        }
    except HTTPException:
        raise
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            saved_file.unlink()
        except:
            pass


# ============================================================================
# Split Invoices API (by invoice number detection)
# ============================================================================

@app.post("/api/split-invoices")
async def split_invoices(
    file: UploadFile = File(...),
    invoice_pattern: str = Form(r"(?:Factura|Invoice|Nr\.?\s*fact(?:ura)?|Invoice\s*(?:No|Number|#)?)[:\s]*([A-Z0-9\-\/]+)", flags=0),
    dpi: int = Form(150),
    filename_template: str = Form("{invoice}")
):
    """Split multi-invoice PDFs by detecting invoice numbers."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = save_upload_file(file)
    try:
        ocr = get_ocr()
        doc = fitz.open(str(saved_file))
        total_pages = doc.page_count

        page_invoices = []
        for i in range(total_pages):
            page = doc[i]
            # Try native text first (faster, preserves diacritics)
            full_text = page.get_text("text")

            # Fallback to OCR only if page has very little native text (scanned)
            if len(full_text.strip()) < 10:
                mat = fitz.Matrix(dpi / 72, dpi / 72)
                pix = page.get_pixmap(matrix=mat)
                img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, 3)

                result, _ = ocr(img)
                texts = [r[1] for r in result] if result else []
                full_text = ' '.join(texts)

            match = re.search(invoice_pattern, full_text, re.IGNORECASE)
            invoice_num = match.group(1) if match else None
            page_invoices.append((i, invoice_num))

            jobs[job_id]["progress"] = int((i + 1) / total_pages * 70)

        groups = {}
        current_invoice = None
        orphan_pages = []

        for page_idx, invoice_num in page_invoices:
            if invoice_num and invoice_num != current_invoice:
                current_invoice = invoice_num
                if current_invoice not in groups:
                    groups[current_invoice] = []
                groups[current_invoice].append(page_idx)
            elif current_invoice:
                groups[current_invoice].append(page_idx)
            else:
                orphan_pages.append(page_idx)

        output_files = []
        for idx, (invoice, pages) in enumerate(groups.items()):
            safe_invoice = re.sub(r'[^\w\-]', '_', invoice)
            fname = filename_template.replace("{invoice}", safe_invoice).replace("{pages}", str(len(pages))).replace("{index}", str(idx + 1))
            out_path = create_output_path(fname)

            new_doc = fitz.open()
            for page_idx in pages:
                new_doc.insert_pdf(doc, from_page=page_idx, to_page=page_idx)
            new_doc.save(str(out_path))
            new_doc.close()

            output_files.append({
                "filename": out_path.name,
                "download_url": f"/api/download/{out_path.name}",
                "invoice": invoice,
                "pages": len(pages)
            })
            jobs[job_id]["progress"] = 70 + int((idx + 1) / max(len(groups), 1) * 25)

        if orphan_pages:
            out_path = create_output_path("_no_invoice_detected")
            new_doc = fitz.open()
            for page_idx in orphan_pages:
                new_doc.insert_pdf(doc, from_page=page_idx, to_page=page_idx)
            new_doc.save(str(out_path))
            new_doc.close()
            output_files.append({
                "filename": out_path.name,
                "download_url": f"/api/download/{out_path.name}",
                "invoice": "N/A",
                "pages": len(orphan_pages)
            })

        doc.close()
        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100

        return {
            "job_id": job_id,
            "status": "done",
            "files": output_files,
            "total_invoices": len(groups),
            "total_files": len(output_files)
        }
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            saved_file.unlink()
        except:
            pass


# ============================================================================
# Auto-Rename PDF API
# ============================================================================

@app.post("/api/auto-rename")
async def auto_rename_pdf(
    file: UploadFile = File(...),
    template: str = Form("{date}_{type}_{number}"),
    dpi: int = Form(150)
):
    """OCR extracts key data and suggests/applies auto-rename."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = save_upload_file(file)
    try:
        ocr = get_ocr()
        doc = fitz.open(str(saved_file))

        page = doc[0]
        mat = fitz.Matrix(dpi / 72, dpi / 72)
        pix = page.get_pixmap(matrix=mat)
        img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, 3)

        result, _ = ocr(img)
        texts = [r[1] for r in result] if result else []
        full_text = ' '.join(texts)

        jobs[job_id]["progress"] = 50

        date_match = re.search(r'(\d{1,2}[\/\.\-]\d{1,2}[\/\.\-]\d{2,4})', full_text)
        date_str = date_match.group(1).replace('/', '-').replace('.', '-') if date_match else "nodate"

        invoice_match = re.search(r'(?:Factura|Invoice|Nr\.?\s*fact)[:\s]*([A-Z0-9\-\/]+)', full_text, re.IGNORECASE)
        number = invoice_match.group(1) if invoice_match else ""

        if not number:
            order_match = re.search(r'(?:Comanda|Order|Nr\.?\s*comanda)[:\s]*(\d+)', full_text, re.IGNORECASE)
            number = order_match.group(1) if order_match else "unknown"

        doc_type = "document"
        type_patterns = {
            "factura": r'\b(factura|invoice)\b',
            "contract": r'\b(contract|agreement)\b',
            "bon": r'\b(bon\s*fiscal|receipt|chitanta)\b',
            "comanda": r'\b(comanda|order)\b',
            "oferta": r'\b(oferta|quote|quotation)\b'
        }
        for t, p in type_patterns.items():
            if re.search(p, full_text, re.IGNORECASE):
                doc_type = t
                break

        safe_number = re.sub(r'[^\w\-]', '_', number)
        new_name = template.replace("{date}", date_str).replace("{type}", doc_type).replace("{number}", safe_number)
        new_name = re.sub(r'[^\w\-_]', '', new_name)
        if not new_name:
            new_name = f"document_{uuid.uuid4().hex[:8]}"

        out_path = OUTPUT_DIR / f"{new_name}.pdf"
        counter = 1
        while out_path.exists():
            out_path = OUTPUT_DIR / f"{new_name}_{counter}.pdf"
            counter += 1

        doc.save(str(out_path))
        doc.close()

        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{out_path.name}",
            "original_name": file.filename,
            "new_name": out_path.name,
            "extracted_data": {
                "date": date_str,
                "type": doc_type,
                "number": number
            }
        }
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            saved_file.unlink()
        except:
            pass


# ============================================================================
# Document Detector API
# ============================================================================

DOCUMENT_TYPES = {
    "invoice": {
        "keywords": ["factura", "invoice", "total", "subtotal", "tva", "vat", "suma", "amount", "pret", "price", "cantitate", "quantity"],
        "weight": 0
    },
    "contract": {
        "keywords": ["contract", "agreement", "party", "parties", "terms", "conditions", "clause", "semnat", "signed", "acord"],
        "weight": 0
    },
    "receipt": {
        "keywords": ["bon fiscal", "receipt", "chitanta", "casier", "cashier", "plata", "payment", "rest", "change"],
        "weight": 0
    },
    "form": {
        "keywords": ["formular", "form", "application", "cerere", "declaratie", "declaration", "completati", "fill"],
        "weight": 0
    },
    "letter": {
        "keywords": ["dear", "stimate", "sincerely", "regards", "cu stima", "referitor", "regarding", "subject"],
        "weight": 0
    },
    "report": {
        "keywords": ["raport", "report", "summary", "rezumat", "analysis", "analiza", "findings", "concluzii"],
        "weight": 0
    },
    "id_document": {
        "keywords": ["carte de identitate", "identity card", "pasaport", "passport", "cnp", "seria", "series"],
        "weight": 0
    },
    "bank_statement": {
        "keywords": ["extras de cont", "bank statement", "sold", "balance", "tranzactii", "transactions", "debit", "credit"],
        "weight": 0
    }
}


@app.post("/api/document-detector")
async def detect_document_type(
    file: UploadFile = File(...),
    dpi: int = Form(150),
    pages_to_scan: int = Form(2)
):
    """Classify document type using OCR and keyword matching."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = save_upload_file(file)
    try:
        ocr = get_ocr()
        doc = fitz.open(str(saved_file))
        total_pages = min(doc.page_count, pages_to_scan)

        full_text = ""
        for i in range(total_pages):
            page = doc[i]
            text = page.get_text()

            if len(text.strip()) < 50:
                mat = fitz.Matrix(dpi / 72, dpi / 72)
                pix = page.get_pixmap(matrix=mat)
                img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, 3)
                result, _ = ocr(img)
                texts = [r[1] for r in result] if result else []
                text = ' '.join(texts)

            full_text += text + " "
            jobs[job_id]["progress"] = int((i + 1) / total_pages * 80)

        full_text_lower = full_text.lower()

        scores = {}
        for doc_type, config in DOCUMENT_TYPES.items():
            score = 0
            matched_keywords = []
            for keyword in config["keywords"]:
                count = full_text_lower.count(keyword.lower())
                if count > 0:
                    score += count
                    matched_keywords.append(keyword)
            scores[doc_type] = {"score": score, "keywords": matched_keywords}

        detected_type = max(scores.keys(), key=lambda k: scores[k]["score"])
        confidence = "high" if scores[detected_type]["score"] > 5 else "medium" if scores[detected_type]["score"] > 2 else "low"

        if scores[detected_type]["score"] == 0:
            detected_type = "unknown"
            confidence = "low"

        doc.close()
        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100

        return {
            "job_id": job_id,
            "status": "done",
            "document_type": detected_type,
            "confidence": confidence,
            "all_scores": scores,
            "pages_scanned": total_pages,
            "filename": file.filename
        }
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            saved_file.unlink()
        except:
            pass


# ============================================================================
# OCR PDF (Add Text Layer) API
# ============================================================================

@app.post("/api/ocr-layer")
async def add_ocr_layer(
    file: UploadFile = File(...),
    language: str = Form("auto"),
    dpi: int = Form(300)
):
    """Add OCR text layer to scanned PDF, making it searchable."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = save_upload_file(file)
    try:
        ocr = get_ocr()
        doc = fitz.open(str(saved_file))
        total_pages = doc.page_count

        output_path = create_output_path("ocr_searchable")
        new_doc = fitz.open()

        for i in range(total_pages):
            page = doc[i]
            mat = fitz.Matrix(dpi / 72, dpi / 72)
            pix = page.get_pixmap(matrix=mat)

            img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, 3)
            result, _ = ocr(img)

            new_page = new_doc.new_page(width=page.rect.width, height=page.rect.height)
            new_page.show_pdf_page(new_page.rect, doc, i)

            if result:
                scale_x = page.rect.width / pix.width
                scale_y = page.rect.height / pix.height

                for item in result:
                    box = item[0]
                    text = item[1]

                    x0 = min(p[0] for p in box) * scale_x
                    y0 = min(p[1] for p in box) * scale_y
                    x1 = max(p[0] for p in box) * scale_x
                    y1 = max(p[1] for p in box) * scale_y

                    rect = fitz.Rect(x0, y0, x1, y1)
                    fontsize = max(6, min(12, (y1 - y0) * 0.8))

                    new_page.insert_text(
                        (x0, y1 - 2),
                        text,
                        fontsize=fontsize,
                        color=(1, 1, 1),
                        render_mode=3,
                        fontfile=UNICODE_FONT_PATH,
                        fontname="F1"
                    )

            jobs[job_id]["progress"] = int((i + 1) / total_pages * 100)

        new_doc.save(str(output_path))
        new_doc.close()
        doc.close()

        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_path.name}",
            "filename": output_path.name,
            "pages_processed": total_pages
        }
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            saved_file.unlink()
        except:
            pass


# ============================================================================
# Extract Text from PDF (OCR) API
# ============================================================================

@app.post("/api/extract-text-ocr")
async def extract_text_ocr(
    file: UploadFile = File(...),
    output_format: str = Form("txt"),  # txt, json, csv
    dpi: int = Form(150),
    pages: str = Form("all")
):
    """Extract text from PDF using OCR, export in various formats."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = save_upload_file(file)
    try:
        ocr = get_ocr()
        doc = fitz.open(str(saved_file))
        total_pages = doc.page_count

        if pages == "all":
            page_list = list(range(total_pages))
        else:
            page_list = []
            for part in pages.split(","):
                part = part.strip()
                if "-" in part:
                    start, end = map(int, part.split("-"))
                    page_list.extend(range(start - 1, end))
                else:
                    page_list.append(int(part) - 1)
            page_list = [p for p in page_list if 0 <= p < total_pages]

        extracted_data = []
        for idx, page_num in enumerate(page_list):
            page = doc[page_num]
            # Always try native text extraction first (preserves diacritics perfectly)
            text = page.get_text("text")

            # Only use OCR if native text is truly empty (scanned PDF)
            if len(text.strip()) < 10:
                mat = fitz.Matrix(dpi / 72, dpi / 72)
                pix = page.get_pixmap(matrix=mat)
                img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, 3)
                result, _ = ocr(img)
                texts = [r[1] for r in result] if result else []
                text = '\n'.join(texts)

            extracted_data.append({
                "page": page_num + 1,
                "text": text
            })
            jobs[job_id]["progress"] = int((idx + 1) / len(page_list) * 90)

        if output_format == "json":
            output_path = create_output_path("extracted_text", ".json")
            with open(output_path, "w", encoding="utf-8") as f:
                json.dump(extracted_data, f, ensure_ascii=False, indent=2)
        elif output_format == "csv":
            output_path = create_output_path("extracted_text", ".csv")
            with open(output_path, "w", encoding="utf-8") as f:
                f.write("page,text\n")
                for item in extracted_data:
                    escaped_text = item["text"].replace('"', '""').replace('\n', ' ')
                    f.write(f'{item["page"]},"{escaped_text}"\n')
        else:
            output_path = create_output_path("extracted_text", ".txt")
            with open(output_path, "w", encoding="utf-8") as f:
                for item in extracted_data:
                    f.write(f"--- Page {item['page']} ---\n\n")
                    f.write(item["text"])
                    f.write("\n\n")

        doc.close()
        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_path.name}",
            "filename": output_path.name,
            "pages_processed": len(page_list),
            "format": output_format
        }
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            saved_file.unlink()
        except:
            pass


# ============================================================================
# Sign PDF API
# ============================================================================

@app.post("/api/sign-pdf")
async def sign_pdf(
    file: UploadFile = File(...),
    signature_text: str = Form(...),
    position: str = Form("bottom-right"),
    page: str = Form("all"),
    font_size: int = Form(14),
    opacity: float = Form(1.0)
):
    """Add text signature to PDF pages."""
    job_id = uuid.uuid4().hex[:8]
    jobs[job_id] = {"id": job_id, "status": "processing", "progress": 0, "started_at": time.time()}

    saved_file = save_upload_file(file)
    try:
        if not signature_text.strip():
            raise HTTPException(status_code=400, detail="Signature text is required")

        doc = fitz.open(str(saved_file))
        total_pages = doc.page_count

        # Parse page selection
        if page == "all":
            page_list = list(range(total_pages))
        else:
            page_list = []
            for part in page.split(","):
                part = part.strip()
                if "-" in part:
                    start, end = map(int, part.split("-"))
                    page_list.extend(range(start - 1, end))
                else:
                    page_list.append(int(part) - 1)
            page_list = [p for p in page_list if 0 <= p < total_pages]

        position_map = {
            "top-left": (0.05, 0.07),
            "top-right": (0.60, 0.07),
            "bottom-left": (0.05, 0.95),
            "bottom-right": (0.60, 0.95),
            "center": (0.35, 0.50)
        }
        pos_x, pos_y = position_map.get(position, position_map["bottom-right"])

        for idx, page_num in enumerate(page_list):
            if 0 <= page_num < total_pages:
                pg = doc[page_num]
                rect = pg.rect
                x = rect.width * pos_x
                y = rect.height * pos_y

                # Use overlay for opacity support
                if opacity < 1.0:
                    overlay_doc = fitz.open()
                    overlay_page = overlay_doc.new_page(width=rect.width, height=rect.height)
                    overlay_page.insert_text(
                        (x, y), signature_text,
                        fontsize=font_size,
                        color=(0, 0, 0.4),
                        fontfile=UNICODE_FONT_PATH,
                        fontname="F1",
                        overlay=True
                    )
                    pg.show_pdf_page(rect, overlay_doc, 0, overlay=(True, None, opacity))
                    overlay_doc.close()
                else:
                    pg.insert_text(
                        (x, y), signature_text,
                        fontsize=font_size,
                        color=(0, 0, 0.4),
                        fontfile=UNICODE_FONT_PATH,
                        fontname="F1",
                        overlay=True
                    )

            jobs[job_id]["progress"] = int((idx + 1) / len(page_list) * 100)

        output_path = create_output_path("signed")
        doc.save(str(output_path))
        doc.close()

        jobs[job_id]["status"] = "done"
        jobs[job_id]["progress"] = 100

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_path.name}",
            "filename": output_path.name,
            "pages_signed": len(page_list)
        }
    except HTTPException:
        raise
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=f"Signature failed: {str(e)}")
    finally:
        try:
            saved_file.unlink()
        except:
            pass


# ============================================================================
# Batch Processing Configuration
# ============================================================================

from concurrent.futures import ThreadPoolExecutor, as_completed
import tempfile

# Configurable worker pool for batch operations
BATCH_WORKERS = int(os.environ.get("BATCH_WORKERS", 4))
batch_executor = ThreadPoolExecutor(max_workers=BATCH_WORKERS)

# Batch job tracking (separate from single-file jobs)
batch_jobs = {}


def require_batch_access(user: Optional[Dict[str, Any]]) -> None:
    """Require silver or gold tier for batch processing."""
    if not user:
        raise HTTPException(status_code=401, detail="Authentication required for batch processing")
    limits = get_user_plan_limits(user)
    if not limits.get("batch_processing", False):
        raise HTTPException(
            status_code=403,
            detail="Batch processing requires Silver or Gold plan. Upgrade at /pricing"
        )


# ============================================================================
# Archive Processor API - Process ZIP with multiple PDFs
# ============================================================================

DOCUMENT_TYPE_PATTERNS = {
    "invoice": [
        r"(?:factura|invoice|faktura|rechnung|facture)[:\s]*",
        r"(?:nr\.?\s*fact(?:ura)?|invoice\s*(?:no|number|#)?)[:\s]*",
        r"(?:total|subtotal|vat|tva)[:\s]*[\d\.,]+",
    ],
    "receipt": [
        r"(?:receipt|bon\s*fiscal|chitanta|quittung|reçu)[:\s]*",
        r"(?:bon\s*nr|receipt\s*#)[:\s]*",
        r"(?:cash|card|total)[:\s]*[\d\.,]+",
    ],
    "contract": [
        r"(?:contract|agreement|acord|vertrag|contrat)[:\s]*",
        r"(?:parties|parti|between|intre)[:\s]*",
        r"(?:terms\s*and\s*conditions|termeni|conditions)[:\s]*",
    ],
}


def detect_document_type(text: str) -> str:
    """Detect document type from OCR text."""
    text_lower = text.lower()
    scores = {doc_type: 0 for doc_type in DOCUMENT_TYPE_PATTERNS}

    for doc_type, patterns in DOCUMENT_TYPE_PATTERNS.items():
        for pattern in patterns:
            if re.search(pattern, text_lower, re.IGNORECASE):
                scores[doc_type] += 1

    max_score = max(scores.values())
    if max_score > 0:
        for doc_type, score in scores.items():
            if score == max_score:
                return doc_type
    return "other"


def extract_document_identifier(text: str, doc_type: str) -> Optional[str]:
    """Extract identifier from document (invoice number, receipt number, etc.)."""
    patterns = {
        "invoice": r"(?:Factura|Invoice|Nr\.?\s*fact(?:ura)?|Invoice\s*(?:No|Number|#)?)[:\s]*([A-Z0-9\-\/]+)",
        "receipt": r"(?:Receipt|Bon|Nr\.?\s*bon|Receipt\s*#)[:\s]*([A-Z0-9\-\/]+)",
        "contract": r"(?:Contract|Nr\.?\s*contract|Agreement\s*#)[:\s]*([A-Z0-9\-\/]+)",
    }
    pattern = patterns.get(doc_type)
    if pattern:
        match = re.search(pattern, text, re.IGNORECASE)
        if match:
            return match.group(1)
    return None


def process_single_pdf_for_archive(pdf_path: Path, ocr, dpi: int = 150) -> Dict[str, Any]:
    """Process a single PDF for archive organization."""
    try:
        doc = fitz.open(str(pdf_path))
        first_page = doc[0]
        mat = fitz.Matrix(dpi / 72, dpi / 72)
        pix = first_page.get_pixmap(matrix=mat)
        img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, 3)

        result, _ = ocr(img)
        texts = [r[1] for r in result] if result else []
        full_text = ' '.join(texts)

        doc_type = detect_document_type(full_text)
        identifier = extract_document_identifier(full_text, doc_type)

        doc.close()

        return {
            "path": pdf_path,
            "type": doc_type,
            "identifier": identifier,
            "pages": doc.page_count if hasattr(doc, 'page_count') else 1,
            "size": pdf_path.stat().st_size,
            "success": True
        }
    except Exception as e:
        return {
            "path": pdf_path,
            "type": "error",
            "identifier": None,
            "error": str(e),
            "success": False
        }


@app.post("/api/archive-processor")
async def process_archive(
    file: UploadFile = File(...),
    organize_by: str = Form("type"),  # "type" or "identifier"
    dpi: int = Form(150),
    credentials: HTTPAuthorizationCredentials = Depends(security)
):
    """Process ZIP archive with multiple PDFs, organize into folders."""
    user = get_current_user(credentials)
    require_batch_access(user)

    job_id = uuid.uuid4().hex[:8]
    batch_jobs[job_id] = {
        "id": job_id,
        "type": "archive-processor",
        "status": "processing",
        "progress": 0,
        "total_files": 0,
        "processed_files": 0,
        "files_by_type": {},
        "errors": [],
        "started_at": time.time()
    }

    saved_file = save_upload_file(file)
    temp_dir = None

    try:
        # Extract ZIP
        temp_dir = Path(tempfile.mkdtemp(prefix="archive_"))
        extract_dir = temp_dir / "extracted"
        extract_dir.mkdir()

        with zipfile.ZipFile(str(saved_file), 'r') as zf:
            zf.extractall(str(extract_dir))

        # Find all PDFs
        pdf_files = list(extract_dir.rglob("*.pdf"))
        batch_jobs[job_id]["total_files"] = len(pdf_files)

        if not pdf_files:
            raise HTTPException(status_code=400, detail="No PDF files found in archive")

        # Process PDFs concurrently
        ocr = get_ocr()
        results = []

        with ThreadPoolExecutor(max_workers=BATCH_WORKERS) as executor:
            futures = {
                executor.submit(process_single_pdf_for_archive, pdf_path, ocr, dpi): pdf_path
                for pdf_path in pdf_files
            }

            for future in as_completed(futures):
                result = future.result()
                results.append(result)
                batch_jobs[job_id]["processed_files"] += 1
                batch_jobs[job_id]["progress"] = int(
                    batch_jobs[job_id]["processed_files"] / len(pdf_files) * 70
                )

                if result["success"]:
                    doc_type = result["type"]
                    if doc_type not in batch_jobs[job_id]["files_by_type"]:
                        batch_jobs[job_id]["files_by_type"][doc_type] = 0
                    batch_jobs[job_id]["files_by_type"][doc_type] += 1
                else:
                    batch_jobs[job_id]["errors"].append({
                        "file": result["path"].name,
                        "error": result.get("error", "Unknown error")
                    })

        # Create organized ZIP
        output_dir = temp_dir / "organized"
        output_dir.mkdir()

        # Create folder structure
        folder_mapping = {
            "invoice": "Invoices",
            "receipt": "Receipts",
            "contract": "Contracts",
            "other": "Other",
            "error": "Errors"
        }

        for folder_name in folder_mapping.values():
            (output_dir / folder_name).mkdir(exist_ok=True)

        # Move files to appropriate folders
        for result in results:
            if result["success"]:
                doc_type = result["type"]
                folder = folder_mapping.get(doc_type, "Other")

                if organize_by == "identifier" and result["identifier"]:
                    new_name = f"{result['identifier']}.pdf"
                else:
                    new_name = result["path"].name

                dest_path = output_dir / folder / new_name
                counter = 1
                while dest_path.exists():
                    stem = Path(new_name).stem
                    dest_path = output_dir / folder / f"{stem}_{counter}.pdf"
                    counter += 1

                shutil.copy2(str(result["path"]), str(dest_path))

        batch_jobs[job_id]["progress"] = 90

        # Create output ZIP
        output_zip_path = create_output_path("organized_archive", ".zip")
        with zipfile.ZipFile(str(output_zip_path), 'w', zipfile.ZIP_DEFLATED) as zf:
            for folder in output_dir.iterdir():
                if folder.is_dir():
                    for file_path in folder.iterdir():
                        zf.write(str(file_path), f"{folder.name}/{file_path.name}")

        batch_jobs[job_id]["status"] = "done"
        batch_jobs[job_id]["progress"] = 100
        batch_jobs[job_id]["finished_at"] = time.time()

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_zip_path.name}",
            "filename": output_zip_path.name,
            "total_files": len(pdf_files),
            "files_by_type": batch_jobs[job_id]["files_by_type"],
            "errors": batch_jobs[job_id]["errors"]
        }

    except HTTPException:
        raise
    except Exception as e:
        batch_jobs[job_id]["status"] = "error"
        batch_jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        try:
            saved_file.unlink()
        except:
            pass
        if temp_dir and temp_dir.exists():
            try:
                shutil.rmtree(str(temp_dir))
            except:
                pass


# ============================================================================
# Batch Document Splitter API
# ============================================================================

def split_pdf_by_boundaries(pdf_path: Path, ocr, pattern: str, dpi: int = 150) -> List[Dict[str, Any]]:
    """Split a single PDF by detected boundaries."""
    try:
        doc = fitz.open(str(pdf_path))
        total_pages = doc.page_count

        page_markers = []
        for i in range(total_pages):
            page = doc[i]
            mat = fitz.Matrix(dpi / 72, dpi / 72)
            pix = page.get_pixmap(matrix=mat)
            img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, 3)

            result, _ = ocr(img)
            texts = [r[1] for r in result] if result else []
            full_text = ' '.join(texts)

            match = re.search(pattern, full_text, re.IGNORECASE)
            marker = match.group(1) if match else None
            page_markers.append((i, marker))

        # Group pages
        groups = {}
        current_marker = None
        orphan_pages = []

        for page_idx, marker in page_markers:
            if marker and marker != current_marker:
                current_marker = marker
                if current_marker not in groups:
                    groups[current_marker] = []
                groups[current_marker].append(page_idx)
            elif current_marker:
                groups[current_marker].append(page_idx)
            else:
                orphan_pages.append(page_idx)

        # Create split PDFs
        split_results = []
        for marker, pages in groups.items():
            new_doc = fitz.open()
            for page_idx in pages:
                new_doc.insert_pdf(doc, from_page=page_idx, to_page=page_idx)

            output_path = create_output_path(f"split_{marker}")
            new_doc.save(str(output_path))
            new_doc.close()

            split_results.append({
                "marker": marker,
                "pages": len(pages),
                "filename": output_path.name,
                "download_url": f"/api/download/{output_path.name}"
            })

        if orphan_pages:
            new_doc = fitz.open()
            for page_idx in orphan_pages:
                new_doc.insert_pdf(doc, from_page=page_idx, to_page=page_idx)
            output_path = create_output_path("split_no_marker")
            new_doc.save(str(output_path))
            new_doc.close()

            split_results.append({
                "marker": "no_marker",
                "pages": len(orphan_pages),
                "filename": output_path.name,
                "download_url": f"/api/download/{output_path.name}"
            })

        doc.close()

        return split_results
    except Exception as e:
        return [{"error": str(e), "source": pdf_path.name}]


@app.post("/api/batch-document-splitter")
async def batch_document_splitter(
    files: List[UploadFile] = File(...),
    pattern: str = Form(r"(?:Page|Document|Section)[:\s]*(\d+)"),
    dpi: int = Form(150),
    credentials: HTTPAuthorizationCredentials = Depends(security)
):
    """Batch split multiple PDFs by detected boundaries."""
    user = get_current_user(credentials)
    require_batch_access(user)

    if len(files) > 100:
        raise HTTPException(status_code=400, detail="Maximum 100 files per batch")

    job_id = uuid.uuid4().hex[:8]
    batch_jobs[job_id] = {
        "id": job_id,
        "type": "batch-splitter",
        "status": "processing",
        "progress": 0,
        "total_files": len(files),
        "processed_files": 0,
        "results": [],
        "errors": [],
        "started_at": time.time()
    }

    saved_files = []
    try:
        # Save all uploaded files
        for f in files:
            saved_files.append(save_upload_file(f))

        ocr = get_ocr()
        all_results = []

        for idx, pdf_path in enumerate(saved_files):
            results = split_pdf_by_boundaries(pdf_path, ocr, pattern, dpi)

            for r in results:
                if "error" in r:
                    batch_jobs[job_id]["errors"].append(r)
                else:
                    all_results.append(r)

            batch_jobs[job_id]["processed_files"] += 1
            batch_jobs[job_id]["progress"] = int((idx + 1) / len(files) * 90)

        batch_jobs[job_id]["results"] = all_results

        # Create ZIP with all split files
        output_zip_path = create_output_path("batch_split", ".zip")
        with zipfile.ZipFile(str(output_zip_path), 'w', zipfile.ZIP_DEFLATED) as zf:
            for result in all_results:
                file_path = OUTPUT_DIR / result["filename"]
                if file_path.exists():
                    zf.write(str(file_path), result["filename"])

        batch_jobs[job_id]["status"] = "done"
        batch_jobs[job_id]["progress"] = 100
        batch_jobs[job_id]["finished_at"] = time.time()

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_zip_path.name}",
            "zip_filename": output_zip_path.name,
            "total_source_files": len(files),
            "total_split_files": len(all_results),
            "files": all_results,
            "errors": batch_jobs[job_id]["errors"]
        }

    except HTTPException:
        raise
    except Exception as e:
        batch_jobs[job_id]["status"] = "error"
        batch_jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        for f in saved_files:
            try:
                f.unlink()
            except:
                pass


# ============================================================================
# Invoice Extractor API - Extract data to Excel/JSON
# ============================================================================

INVOICE_FIELD_PATTERNS = {
    "invoice_number": [
        r"(?:Factura|Invoice|Nr\.?\s*fact(?:ura)?|Invoice\s*(?:No|Number|#)?)[:\s]*([A-Z0-9\-\/]+)",
        r"(?:Invoice\s*#|Factura\s*nr\.?)[:\s]*([A-Z0-9\-\/]+)",
    ],
    "date": [
        r"(?:Data|Date|Datum)[:\s]*(\d{1,2}[\.\/\-]\d{1,2}[\.\/\-]\d{2,4})",
        r"(\d{1,2}[\.\/\-]\d{1,2}[\.\/\-]\d{2,4})",
    ],
    "company": [
        r"(?:Furnizor|Seller|From|Company)[:\s]*([A-Za-z0-9\s\.\,]+?)(?:\n|$)",
        r"(?:S\.?R\.?L\.?|S\.?A\.?|LLC|Inc\.?|Ltd\.?)\s*([A-Za-z0-9\s\.]+)",
    ],
    "total": [
        r"(?:Total|TOTAL|Grand\s*Total|Total\s*de\s*plata)[:\s]*[\$€£]?\s*([\d\.,]+)",
        r"(?:Amount\s*Due|Suma\s*de\s*plata)[:\s]*[\$€£]?\s*([\d\.,]+)",
    ],
    "vat": [
        r"(?:TVA|VAT|Tax)[:\s]*[\$€£]?\s*([\d\.,]+)",
    ],
    "subtotal": [
        r"(?:Subtotal|Sub-?total)[:\s]*[\$€£]?\s*([\d\.,]+)",
    ],
}


def extract_invoice_data(pdf_path: Path, ocr, dpi: int = 150) -> Dict[str, Any]:
    """Extract invoice data from a PDF."""
    try:
        doc = fitz.open(str(pdf_path))

        all_text = ""
        for page in doc:
            mat = fitz.Matrix(dpi / 72, dpi / 72)
            pix = page.get_pixmap(matrix=mat)
            img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, 3)

            result, _ = ocr(img)
            texts = [r[1] for r in result] if result else []
            all_text += ' '.join(texts) + "\n"

        doc.close()

        extracted = {"source_file": pdf_path.name, "pages": doc.page_count if hasattr(doc, 'page_count') else 1}

        for field, patterns in INVOICE_FIELD_PATTERNS.items():
            for pattern in patterns:
                match = re.search(pattern, all_text, re.IGNORECASE)
                if match:
                    extracted[field] = match.group(1).strip()
                    break
            if field not in extracted:
                extracted[field] = None

        return {"success": True, "data": extracted}

    except Exception as e:
        return {"success": False, "error": str(e), "source_file": pdf_path.name}


@app.post("/api/invoice-extractor")
async def extract_invoices(
    files: List[UploadFile] = File(...),
    output_format: str = Form("json"),  # "json" or "excel"
    dpi: int = Form(150),
    credentials: HTTPAuthorizationCredentials = Depends(security)
):
    """Extract invoice data from multiple PDFs to Excel or JSON."""
    user = get_current_user(credentials)
    require_batch_access(user)

    if len(files) > 100:
        raise HTTPException(status_code=400, detail="Maximum 100 files per batch")

    job_id = uuid.uuid4().hex[:8]
    batch_jobs[job_id] = {
        "id": job_id,
        "type": "invoice-extractor",
        "status": "processing",
        "progress": 0,
        "total_files": len(files),
        "processed_files": 0,
        "started_at": time.time()
    }

    saved_files = []
    try:
        for f in files:
            saved_files.append(save_upload_file(f))

        ocr = get_ocr()
        extracted_data = []
        errors = []

        with ThreadPoolExecutor(max_workers=BATCH_WORKERS) as executor:
            futures = {
                executor.submit(extract_invoice_data, pdf_path, ocr, dpi): pdf_path
                for pdf_path in saved_files
            }

            for future in as_completed(futures):
                result = future.result()
                batch_jobs[job_id]["processed_files"] += 1
                batch_jobs[job_id]["progress"] = int(
                    batch_jobs[job_id]["processed_files"] / len(files) * 80
                )

                if result["success"]:
                    extracted_data.append(result["data"])
                else:
                    errors.append({
                        "file": result.get("source_file", "unknown"),
                        "error": result.get("error", "Unknown error")
                    })

        # Generate output file
        if output_format == "excel":
            from openpyxl import Workbook

            wb = Workbook()
            ws = wb.active
            ws.title = "Invoice Data"

            # Header row
            headers = ["Source File", "Invoice Number", "Date", "Company", "Subtotal", "VAT", "Total"]
            for col, header in enumerate(headers, 1):
                ws.cell(row=1, column=col, value=header)

            # Data rows
            for row_idx, data in enumerate(extracted_data, 2):
                ws.cell(row=row_idx, column=1, value=data.get("source_file", ""))
                ws.cell(row=row_idx, column=2, value=data.get("invoice_number", ""))
                ws.cell(row=row_idx, column=3, value=data.get("date", ""))
                ws.cell(row=row_idx, column=4, value=data.get("company", ""))
                ws.cell(row=row_idx, column=5, value=data.get("subtotal", ""))
                ws.cell(row=row_idx, column=6, value=data.get("vat", ""))
                ws.cell(row=row_idx, column=7, value=data.get("total", ""))

            output_path = create_output_path("invoices_extracted", ".xlsx")
            wb.save(str(output_path))

            content_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        else:
            output_path = create_output_path("invoices_extracted", ".json")
            with open(output_path, "w", encoding="utf-8") as f:
                json.dump({
                    "invoices": extracted_data,
                    "errors": errors,
                    "total_processed": len(extracted_data),
                    "total_errors": len(errors)
                }, f, indent=2, ensure_ascii=False)

            content_type = "application/json"

        batch_jobs[job_id]["status"] = "done"
        batch_jobs[job_id]["progress"] = 100
        batch_jobs[job_id]["finished_at"] = time.time()

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_path.name}",
            "filename": output_path.name,
            "format": output_format,
            "total_processed": len(extracted_data),
            "total_errors": len(errors),
            "preview": extracted_data[:5],
            "errors": errors
        }

    except HTTPException:
        raise
    except Exception as e:
        batch_jobs[job_id]["status"] = "error"
        batch_jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        for f in saved_files:
            try:
                f.unlink()
            except:
                pass


# ============================================================================
# Receipt Extractor API
# ============================================================================

RECEIPT_FIELD_PATTERNS = {
    "merchant": [
        r"^([A-Z][A-Za-z0-9\s\.\,\&]+?)(?:\n|$)",
        r"(?:Magazin|Store|Shop)[:\s]*([A-Za-z0-9\s\.]+)",
    ],
    "date": [
        r"(?:Data|Date|Datum)[:\s]*(\d{1,2}[\.\/\-]\d{1,2}[\.\/\-]\d{2,4})",
        r"(\d{1,2}[\.\/\-]\d{1,2}[\.\/\-]\d{2,4})",
    ],
    "time": [
        r"(?:Ora|Time|Hour)[:\s]*(\d{1,2}:\d{2}(?::\d{2})?)",
        r"(\d{1,2}:\d{2}(?::\d{2})?(?:\s*[AP]M)?)",
    ],
    "total": [
        r"(?:Total|TOTAL)[:\s]*[\$€£RON]*\s*([\d\.,]+)",
        r"(?:Amount|Suma)[:\s]*[\$€£RON]*\s*([\d\.,]+)",
    ],
    "payment_method": [
        r"(?:Card|Cash|Numerar|Credit)[:\s]*",
        r"(VISA|MASTERCARD|AMEX|CASH|CARD)",
    ],
    "receipt_number": [
        r"(?:Bon\s*(?:fiscal)?|Receipt)[:\s#]*([A-Z0-9\-]+)",
        r"(?:Nr\.?\s*bon|Doc\s*#)[:\s]*([A-Z0-9\-]+)",
    ],
}


def extract_receipt_data(pdf_path: Path, ocr, dpi: int = 150) -> Dict[str, Any]:
    """Extract receipt data from a PDF."""
    try:
        doc = fitz.open(str(pdf_path))

        all_text = ""
        for page in doc:
            mat = fitz.Matrix(dpi / 72, dpi / 72)
            pix = page.get_pixmap(matrix=mat)
            img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, 3)

            result, _ = ocr(img)
            texts = [r[1] for r in result] if result else []
            all_text += ' '.join(texts) + "\n"

        doc.close()

        extracted = {"source_file": pdf_path.name}

        for field, patterns in RECEIPT_FIELD_PATTERNS.items():
            for pattern in patterns:
                match = re.search(pattern, all_text, re.IGNORECASE | re.MULTILINE)
                if match:
                    extracted[field] = match.group(1).strip() if match.lastindex else match.group(0).strip()
                    break
            if field not in extracted:
                extracted[field] = None

        return {"success": True, "data": extracted}

    except Exception as e:
        return {"success": False, "error": str(e), "source_file": pdf_path.name}


@app.post("/api/receipt-extractor")
async def extract_receipts(
    files: List[UploadFile] = File(...),
    output_format: str = Form("json"),  # "json" or "excel"
    dpi: int = Form(150),
    credentials: HTTPAuthorizationCredentials = Depends(security)
):
    """Extract receipt data from multiple PDFs to Excel or JSON."""
    user = get_current_user(credentials)
    require_batch_access(user)

    if len(files) > 100:
        raise HTTPException(status_code=400, detail="Maximum 100 files per batch")

    job_id = uuid.uuid4().hex[:8]
    batch_jobs[job_id] = {
        "id": job_id,
        "type": "receipt-extractor",
        "status": "processing",
        "progress": 0,
        "total_files": len(files),
        "processed_files": 0,
        "started_at": time.time()
    }

    saved_files = []
    try:
        for f in files:
            saved_files.append(save_upload_file(f))

        ocr = get_ocr()
        extracted_data = []
        errors = []

        with ThreadPoolExecutor(max_workers=BATCH_WORKERS) as executor:
            futures = {
                executor.submit(extract_receipt_data, pdf_path, ocr, dpi): pdf_path
                for pdf_path in saved_files
            }

            for future in as_completed(futures):
                result = future.result()
                batch_jobs[job_id]["processed_files"] += 1
                batch_jobs[job_id]["progress"] = int(
                    batch_jobs[job_id]["processed_files"] / len(files) * 80
                )

                if result["success"]:
                    extracted_data.append(result["data"])
                else:
                    errors.append({
                        "file": result.get("source_file", "unknown"),
                        "error": result.get("error", "Unknown error")
                    })

        # Generate output file
        if output_format == "excel":
            from openpyxl import Workbook

            wb = Workbook()
            ws = wb.active
            ws.title = "Receipt Data"

            headers = ["Source File", "Merchant", "Date", "Time", "Receipt #", "Total", "Payment Method"]
            for col, header in enumerate(headers, 1):
                ws.cell(row=1, column=col, value=header)

            for row_idx, data in enumerate(extracted_data, 2):
                ws.cell(row=row_idx, column=1, value=data.get("source_file", ""))
                ws.cell(row=row_idx, column=2, value=data.get("merchant", ""))
                ws.cell(row=row_idx, column=3, value=data.get("date", ""))
                ws.cell(row=row_idx, column=4, value=data.get("time", ""))
                ws.cell(row=row_idx, column=5, value=data.get("receipt_number", ""))
                ws.cell(row=row_idx, column=6, value=data.get("total", ""))
                ws.cell(row=row_idx, column=7, value=data.get("payment_method", ""))

            output_path = create_output_path("receipts_extracted", ".xlsx")
            wb.save(str(output_path))
        else:
            output_path = create_output_path("receipts_extracted", ".json")
            with open(output_path, "w", encoding="utf-8") as f:
                json.dump({
                    "receipts": extracted_data,
                    "errors": errors,
                    "total_processed": len(extracted_data),
                    "total_errors": len(errors)
                }, f, indent=2, ensure_ascii=False)

        batch_jobs[job_id]["status"] = "done"
        batch_jobs[job_id]["progress"] = 100
        batch_jobs[job_id]["finished_at"] = time.time()

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_path.name}",
            "filename": output_path.name,
            "format": output_format,
            "total_processed": len(extracted_data),
            "total_errors": len(errors),
            "preview": extracted_data[:5],
            "errors": errors
        }

    except HTTPException:
        raise
    except Exception as e:
        batch_jobs[job_id]["status"] = "error"
        batch_jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        for f in saved_files:
            try:
                f.unlink()
            except:
                pass


# ============================================================================
# Batch Job Status API
# ============================================================================

@app.get("/api/batch/status/{job_id}")
def get_batch_job_status(job_id: str):
    """Get status of a batch processing job."""
    if job_id not in batch_jobs:
        raise HTTPException(status_code=404, detail="Batch job not found")

    job = batch_jobs[job_id]
    return {
        "id": job["id"],
        "type": job.get("type"),
        "status": job["status"],
        "progress": job.get("progress", 0),
        "total_files": job.get("total_files", 0),
        "processed_files": job.get("processed_files", 0),
        "files_by_type": job.get("files_by_type", {}),
        "errors": job.get("errors", []),
        "error": job.get("error"),
        "started_at": job.get("started_at"),
        "finished_at": job.get("finished_at")
    }


# ============================================================================
# Batch Processing Dashboard - Multi-file operations
# ============================================================================

@app.post("/api/batch/process-multiple")
async def batch_process_multiple(
    files: List[UploadFile] = File(...),
    operation: str = Form(...),  # "compress", "convert-jpg", "ocr", etc.
    settings: str = Form("{}"),  # JSON string of operation-specific settings
    credentials: HTTPAuthorizationCredentials = Depends(security)
):
    """Process multiple files with the same operation."""
    user = get_current_user(credentials)
    require_batch_access(user)

    if len(files) > 100:
        raise HTTPException(status_code=400, detail="Maximum 100 files per batch")

    try:
        operation_settings = json.loads(settings)
    except json.JSONDecodeError:
        operation_settings = {}

    job_id = uuid.uuid4().hex[:8]
    batch_jobs[job_id] = {
        "id": job_id,
        "type": f"batch-{operation}",
        "status": "processing",
        "progress": 0,
        "total_files": len(files),
        "processed_files": 0,
        "results": [],
        "errors": [],
        "started_at": time.time()
    }

    saved_files = []
    output_files = []

    try:
        for f in files:
            saved_files.append((save_upload_file(f), f.filename))

        for idx, (pdf_path, original_name) in enumerate(saved_files):
            try:
                result = await process_single_file_operation(
                    pdf_path, original_name, operation, operation_settings
                )
                output_files.append(result)
                batch_jobs[job_id]["results"].append(result)
            except Exception as e:
                batch_jobs[job_id]["errors"].append({
                    "file": original_name,
                    "error": str(e)
                })

            batch_jobs[job_id]["processed_files"] += 1
            batch_jobs[job_id]["progress"] = int((idx + 1) / len(files) * 90)

        # Create ZIP with all output files
        output_zip_path = create_output_path(f"batch_{operation}", ".zip")
        with zipfile.ZipFile(str(output_zip_path), 'w', zipfile.ZIP_DEFLATED) as zf:
            for result in output_files:
                if "filename" in result:
                    file_path = OUTPUT_DIR / result["filename"]
                    if file_path.exists():
                        zf.write(str(file_path), result["filename"])

        batch_jobs[job_id]["status"] = "done"
        batch_jobs[job_id]["progress"] = 100
        batch_jobs[job_id]["finished_at"] = time.time()

        return {
            "job_id": job_id,
            "status": "done",
            "download_url": f"/api/download/{output_zip_path.name}",
            "zip_filename": output_zip_path.name,
            "total_files": len(files),
            "successful": len(output_files),
            "failed": len(batch_jobs[job_id]["errors"]),
            "results": output_files,
            "errors": batch_jobs[job_id]["errors"]
        }

    except HTTPException:
        raise
    except Exception as e:
        batch_jobs[job_id]["status"] = "error"
        batch_jobs[job_id]["error"] = str(e)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        for pdf_path, _ in saved_files:
            try:
                pdf_path.unlink()
            except:
                pass


async def process_single_file_operation(
    pdf_path: Path,
    original_name: str,
    operation: str,
    settings: Dict[str, Any]
) -> Dict[str, Any]:
    """Process a single file with the specified operation."""
    doc = fitz.open(str(pdf_path))

    if operation == "compress":
        quality = settings.get("quality", "medium")
        quality_settings = {
            "low": {"garbage": 1},
            "medium": {"garbage": 3},
            "high": {"garbage": 4}
        }
        q = quality_settings.get(quality, quality_settings["medium"])

        output_path = create_output_path(f"compressed_{Path(original_name).stem}")
        doc.save(str(output_path), garbage=q["garbage"], deflate=True, clean=True)
        doc.close()

        return {
            "original": original_name,
            "filename": output_path.name,
            "download_url": f"/api/download/{output_path.name}",
            "original_size": pdf_path.stat().st_size,
            "compressed_size": output_path.stat().st_size
        }

    elif operation == "convert-jpg":
        dpi = settings.get("dpi", 150)
        mat = fitz.Matrix(dpi / 72, dpi / 72)

        images = []
        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=mat)
            img_path = OUTPUT_DIR / f"{Path(original_name).stem}_page_{i+1}_{uuid.uuid4().hex[:4]}.jpg"
            pix.save(str(img_path))
            images.append(img_path.name)

        doc.close()

        return {
            "original": original_name,
            "images": images,
            "total_pages": len(images)
        }

    elif operation == "rotate":
        rotation = settings.get("rotation", 90)
        for page in doc:
            page.set_rotation((page.rotation + rotation) % 360)

        output_path = create_output_path(f"rotated_{Path(original_name).stem}")
        doc.save(str(output_path))
        doc.close()

        return {
            "original": original_name,
            "filename": output_path.name,
            "download_url": f"/api/download/{output_path.name}"
        }

    else:
        doc.close()
        raise ValueError(f"Unknown operation: {operation}")


# ============================================================================
# Phase 7: Batch Processing & Archive Tools
# ============================================================================

def detect_document_type(doc: fitz.Document, ocr) -> str:
    """Detect document type based on content."""
    if len(doc) == 0:
        return "Unknown"

    page = doc[0]
    pix = page.get_pixmap(matrix=fitz.Matrix(150/72, 150/72))
    img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, 3)

    result, _ = ocr(img)
    text = ' '.join([r[1] for r in result] if result else []).lower()

    if any(kw in text for kw in ["invoice", "bill", "payment", "amount due", "total:"]):
        return "Invoice"
    elif any(kw in text for kw in ["contract", "agreement", "terms and conditions", "parties"]):
        return "Contract"
    elif any(kw in text for kw in ["receipt", "paid", "transaction", "merchant"]):
        return "Receipt"
    elif any(kw in text for kw in ["statement", "balance", "account"]):
        return "Statement"
    elif any(kw in text for kw in ["report", "summary", "analysis"]):
        return "Report"
    else:
        return "Other"


def extract_invoice_data(doc: fitz.Document, ocr) -> Dict[str, Any]:
    """Extract invoice metadata using OCR."""
    if len(doc) == 0:
        return {}

    page = doc[0]
    pix = page.get_pixmap(matrix=fitz.Matrix(200/72, 200/72))
    img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, 3)

    result, _ = ocr(img)
    text = ' '.join([r[1] for r in result] if result else [])

    data = {
        "invoice_number": None,
        "date": None,
        "company": None,
        "total": None,
        "raw_text": text[:500]
    }

    inv_match = re.search(r"(?:invoice|inv|#)\s*:?\s*([A-Z0-9-]+)", text, re.IGNORECASE)
    if inv_match:
        data["invoice_number"] = inv_match.group(1)

    date_match = re.search(r"(\d{1,2}[/-]\d{1,2}[/-]\d{2,4})", text)
    if date_match:
        data["date"] = date_match.group(1)

    company_lines = text.split('\n')[:3]
    if company_lines:
        data["company"] = company_lines[0].strip()[:50]

    total_match = re.search(r"(?:total|amount due|grand total)\s*:?\s*\$?\s*([\d,]+\.?\d*)", text, re.IGNORECASE)
    if total_match:
        data["total"] = total_match.group(1)

    return data


def extract_receipt_data(doc: fitz.Document, ocr) -> Dict[str, Any]:
    """Extract receipt metadata using OCR."""
    if len(doc) == 0:
        return {}

    page = doc[0]
    pix = page.get_pixmap(matrix=fitz.Matrix(200/72, 200/72))
    img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, 3)

    result, _ = ocr(img)
    text = ' '.join([r[1] for r in result] if result else [])

    data = {
        "merchant": None,
        "date": None,
        "amount": None,
        "raw_text": text[:500]
    }

    merchant_lines = text.split('\n')[:2]
    if merchant_lines:
        data["merchant"] = merchant_lines[0].strip()[:50]

    date_match = re.search(r"(\d{1,2}[/-]\d{1,2}[/-]\d{2,4})", text)
    if date_match:
        data["date"] = date_match.group(1)

    amount_match = re.search(r"(?:total|amount|paid)\s*:?\s*\$?\s*([\d,]+\.?\d*)", text, re.IGNORECASE)
    if amount_match:
        data["amount"] = amount_match.group(1)

    return data


@app.post("/api/archive-processor/upload")
async def archive_processor_upload(
    file: UploadFile = File(...),
    user: Optional[Dict[str, Any]] = Depends(get_current_user)
):
    """Upload ZIP archive with PDFs for processing."""
    limits = get_user_plan_limits(user)

    if not limits.get("batch_processing"):
        raise HTTPException(status_code=403, detail="Batch processing requires Bronze tier or higher")

    if not file.filename.endswith('.zip'):
        raise HTTPException(status_code=400, detail="Only ZIP files are accepted")

    content = await file.read()
    if len(content) > limits["max_file_size_mb"] * 1024 * 1024:
        raise HTTPException(status_code=413, detail=f"File exceeds {limits['max_file_size_mb']}MB limit")

    job_id = uuid.uuid4().hex[:12]
    batch_processing_jobs[job_id] = {
        "id": job_id,
        "status": "processing",
        "progress": 0,
        "started_at": time.time(),
        "total_files": 0,
        "processed_files": 0,
        "files_by_type": {},
        "errors": []
    }

    def process_archive():
        try:
            ocr = get_ocr()
            temp_dir = tempfile.mkdtemp()
            zip_buffer = io.BytesIO(content)

            with zipfile.ZipFile(zip_buffer, 'r') as zip_ref:
                zip_ref.extractall(temp_dir)

            pdf_files = []
            for root, dirs, files in os.walk(temp_dir):
                for filename in files:
                    if filename.lower().endswith('.pdf'):
                        pdf_files.append(os.path.join(root, filename))

            batch_processing_jobs[job_id]["total_files"] = len(pdf_files)

            organized_files = {
                "Invoices": [],
                "Contracts": [],
                "Receipts": [],
                "Statements": [],
                "Reports": [],
                "Other": []
            }

            for idx, pdf_path in enumerate(pdf_files):
                try:
                    doc = fitz.open(pdf_path)
                    doc_type = detect_document_type(doc, ocr)

                    new_name = f"{doc_type}_{idx+1}_{uuid.uuid4().hex[:6]}.pdf"

                    folder = doc_type + "s" if doc_type in ["Invoice", "Contract", "Receipt", "Statement", "Report"] else "Other"
                    organized_files[folder].append((pdf_path, new_name))

                    if doc_type not in batch_processing_jobs[job_id]["files_by_type"]:
                        batch_processing_jobs[job_id]["files_by_type"][doc_type] = 0
                    batch_processing_jobs[job_id]["files_by_type"][doc_type] += 1

                    doc.close()
                except Exception as e:
                    batch_processing_jobs[job_id]["errors"].append({
                        "file": os.path.basename(pdf_path),
                        "error": str(e)
                    })

                batch_processing_jobs[job_id]["processed_files"] += 1
                batch_processing_jobs[job_id]["progress"] = int((idx + 1) / len(pdf_files) * 70)

            batch_processing_jobs[job_id]["progress"] = 80

            output_zip_path = OUTPUT_DIR / f"organized_{job_id}.zip"
            with zipfile.ZipFile(str(output_zip_path), 'w', zipfile.ZIP_DEFLATED) as zip_out:
                for folder, files_list in organized_files.items():
                    if files_list:
                        for pdf_path, new_name in files_list:
                            zip_out.write(pdf_path, f"{folder}/{new_name}")

            shutil.rmtree(temp_dir)

            batch_processing_jobs[job_id]["status"] = "done"
            batch_processing_jobs[job_id]["progress"] = 100
            batch_processing_jobs[job_id]["output_file"] = output_zip_path.name
            batch_processing_jobs[job_id]["finished_at"] = time.time()

        except Exception as e:
            batch_processing_jobs[job_id]["status"] = "error"
            batch_processing_jobs[job_id]["error"] = str(e)
            if 'temp_dir' in locals():
                shutil.rmtree(temp_dir, ignore_errors=True)

    threading.Thread(target=process_archive, daemon=True).start()

    return {"job_id": job_id, "status": "processing"}


@app.post("/api/batch-document-splitter/upload")
async def batch_document_splitter_upload(
    files: List[UploadFile] = File(...),
    split_pattern: str = Form("order"),
    user: Optional[Dict[str, Any]] = Depends(get_current_user)
):
    """Upload multiple PDFs for batch splitting."""
    limits = get_user_plan_limits(user)

    if not limits.get("batch_processing"):
        raise HTTPException(status_code=403, detail="Batch processing requires Bronze tier or higher")

    if len(files) > 100:
        raise HTTPException(status_code=400, detail="Maximum 100 files per batch")

    job_id = uuid.uuid4().hex[:12]
    batch_processing_jobs[job_id] = {
        "id": job_id,
        "status": "processing",
        "progress": 0,
        "started_at": time.time(),
        "total_files": len(files),
        "processed_files": 0,
        "results": [],
        "errors": []
    }

    saved_files = []
    for file in files:
        content = await file.read()
        temp_path = OUTPUT_DIR / f"temp_{uuid.uuid4().hex[:8]}_{file.filename}"
        temp_path.write_bytes(content)
        saved_files.append((temp_path, file.filename))

    def process_batch_split():
        try:
            all_results = []

            for idx, (pdf_path, original_name) in enumerate(saved_files):
                try:
                    result = split_single_document(str(pdf_path), original_name, split_pattern)
                    all_results.append(result)
                except Exception as e:
                    batch_processing_jobs[job_id]["errors"].append({
                        "file": original_name,
                        "error": str(e)
                    })

                batch_processing_jobs[job_id]["processed_files"] += 1
                batch_processing_jobs[job_id]["progress"] = int((idx + 1) / len(saved_files) * 90)

            batch_processing_jobs[job_id]["results"] = all_results

            output_zip_path = OUTPUT_DIR / f"batch_split_{job_id}.zip"
            with zipfile.ZipFile(str(output_zip_path), 'w', zipfile.ZIP_DEFLATED) as zip_out:
                for result in all_results:
                    for split_file in result.get("files", []):
                        file_path = OUTPUT_DIR / split_file["filename"]
                        if file_path.exists():
                            zip_out.write(str(file_path), split_file["filename"])

            batch_processing_jobs[job_id]["status"] = "done"
            batch_processing_jobs[job_id]["progress"] = 100
            batch_processing_jobs[job_id]["output_file"] = output_zip_path.name
            batch_processing_jobs[job_id]["finished_at"] = time.time()

        except Exception as e:
            batch_processing_jobs[job_id]["status"] = "error"
            batch_processing_jobs[job_id]["error"] = str(e)
        finally:
            for pdf_path, _ in saved_files:
                pdf_path.unlink(missing_ok=True)

    threading.Thread(target=process_batch_split, daemon=True).start()

    return {"job_id": job_id, "status": "processing"}


def split_single_document(pdf_path: str, original_name: str, split_pattern: str) -> Dict[str, Any]:
    """Split a single document by detected boundaries."""
    ocr = get_ocr()
    doc = fitz.open(pdf_path)

    if split_pattern == "order":
        pattern = r"(?:order|invoice|#)\s*:?\s*([A-Z0-9-]+)"
    elif split_pattern == "page":
        output_files = []
        for i in range(len(doc)):
            new_doc = fitz.open()
            new_doc.insert_pdf(doc, from_page=i, to_page=i)
            output_path = OUTPUT_DIR / f"{Path(original_name).stem}_page_{i+1}.pdf"
            new_doc.save(str(output_path))
            new_doc.close()
            output_files.append({
                "filename": output_path.name,
                "pages": 1
            })
        doc.close()
        return {
            "original": original_name,
            "files": output_files
        }
    else:
        pattern = r"(?:order|invoice|#)\s*:?\s*([A-Z0-9-]+)"

    page_orders = []
    for i in range(len(doc)):
        page = doc[i]
        order_num, _ = extract_order_from_page(page, ocr, pattern)
        page_orders.append((i, order_num))

    groups = {}
    current_order = None
    for page_idx, order_num in page_orders:
        if order_num and order_num != current_order:
            current_order = order_num
            if current_order not in groups:
                groups[current_order] = []
            groups[current_order].append(page_idx)
        elif current_order:
            groups[current_order].append(page_idx)

    output_files = []
    for order, pages in groups.items():
        new_doc = fitz.open()
        for page_idx in pages:
            new_doc.insert_pdf(doc, from_page=page_idx, to_page=page_idx)
        output_path = OUTPUT_DIR / f"{Path(original_name).stem}_{order}.pdf"
        new_doc.save(str(output_path))
        new_doc.close()
        output_files.append({
            "filename": output_path.name,
            "order": order,
            "pages": len(pages)
        })

    doc.close()

    return {
        "original": original_name,
        "files": output_files
    }


@app.post("/api/invoice-extractor/upload")
async def invoice_extractor_upload(
    files: List[UploadFile] = File(...),
    export_format: str = Form("json"),
    user: Optional[Dict[str, Any]] = Depends(get_current_user)
):
    """Extract invoice data from PDFs and export to Excel/JSON."""
    limits = get_user_plan_limits(user)

    if not limits.get("smart_tools"):
        raise HTTPException(status_code=403, detail="Invoice extractor requires Silver tier or higher")

    if len(files) > 100:
        raise HTTPException(status_code=400, detail="Maximum 100 files per batch")

    job_id = uuid.uuid4().hex[:12]
    batch_processing_jobs[job_id] = {
        "id": job_id,
        "status": "processing",
        "progress": 0,
        "started_at": time.time(),
        "total_files": len(files),
        "processed_files": 0,
        "results": [],
        "errors": []
    }

    saved_files = []
    for file in files:
        content = await file.read()
        temp_path = OUTPUT_DIR / f"temp_{uuid.uuid4().hex[:8]}_{file.filename}"
        temp_path.write_bytes(content)
        saved_files.append((temp_path, file.filename))

    def process_invoice_extraction():
        try:
            ocr = get_ocr()
            all_results = []

            with ThreadPoolExecutor(max_workers=MAX_BATCH_WORKERS) as executor:
                futures = {}
                for pdf_path, original_name in saved_files:
                    future = executor.submit(extract_invoice_data_from_file, pdf_path, original_name, ocr)
                    futures[future] = original_name

                for idx, future in enumerate(as_completed(futures)):
                    try:
                        result = future.result()
                        all_results.append(result)
                    except Exception as e:
                        batch_processing_jobs[job_id]["errors"].append({
                            "file": futures[future],
                            "error": str(e)
                        })

                    batch_processing_jobs[job_id]["processed_files"] += 1
                    batch_processing_jobs[job_id]["progress"] = int(
                        batch_processing_jobs[job_id]["processed_files"] / len(saved_files) * 80
                    )

            batch_processing_jobs[job_id]["results"] = all_results

            if export_format == "excel":
                try:
                    import openpyxl
                    wb = openpyxl.Workbook()
                    ws = wb.active
                    ws.title = "Invoices"
                    ws.append(["Filename", "Invoice Number", "Date", "Company", "Total"])

                    for r in all_results:
                        ws.append([
                            r["filename"],
                            r["invoice_number"],
                            r["date"],
                            r["company"],
                            r["total"]
                        ])

                    output_path = OUTPUT_DIR / f"invoices_{job_id}.xlsx"
                    wb.save(str(output_path))
                    batch_processing_jobs[job_id]["output_file"] = output_path.name
                    batch_processing_jobs[job_id]["export_format"] = "excel"
                except ImportError:
                    raise HTTPException(status_code=500, detail="openpyxl not installed")
            else:
                output_path = OUTPUT_DIR / f"invoices_{job_id}.json"
                output_path.write_text(json.dumps(all_results, indent=2))
                batch_processing_jobs[job_id]["output_file"] = output_path.name
                batch_processing_jobs[job_id]["export_format"] = "json"

            batch_processing_jobs[job_id]["status"] = "done"
            batch_processing_jobs[job_id]["progress"] = 100
            batch_processing_jobs[job_id]["finished_at"] = time.time()

        except Exception as e:
            batch_processing_jobs[job_id]["status"] = "error"
            batch_processing_jobs[job_id]["error"] = str(e)
        finally:
            for pdf_path, _ in saved_files:
                pdf_path.unlink(missing_ok=True)

    threading.Thread(target=process_invoice_extraction, daemon=True).start()

    return {"job_id": job_id, "status": "processing"}


def extract_invoice_data_from_file(pdf_path: Path, original_name: str, ocr) -> Dict[str, Any]:
    """Extract invoice data from a single file."""
    doc = fitz.open(str(pdf_path))
    data = extract_invoice_data(doc, ocr)
    doc.close()

    return {
        "filename": original_name,
        "invoice_number": data.get("invoice_number"),
        "date": data.get("date"),
        "company": data.get("company"),
        "total": data.get("total")
    }


@app.post("/api/receipt-extractor/upload")
async def receipt_extractor_upload(
    files: List[UploadFile] = File(...),
    export_format: str = Form("json"),
    user: Optional[Dict[str, Any]] = Depends(get_current_user)
):
    """Extract receipt data from PDFs and export to Excel/JSON."""
    limits = get_user_plan_limits(user)

    if not limits.get("smart_tools"):
        raise HTTPException(status_code=403, detail="Receipt extractor requires Silver tier or higher")

    if len(files) > 100:
        raise HTTPException(status_code=400, detail="Maximum 100 files per batch")

    job_id = uuid.uuid4().hex[:12]
    batch_processing_jobs[job_id] = {
        "id": job_id,
        "status": "processing",
        "progress": 0,
        "started_at": time.time(),
        "total_files": len(files),
        "processed_files": 0,
        "results": [],
        "errors": []
    }

    saved_files = []
    for file in files:
        content = await file.read()
        temp_path = OUTPUT_DIR / f"temp_{uuid.uuid4().hex[:8]}_{file.filename}"
        temp_path.write_bytes(content)
        saved_files.append((temp_path, file.filename))

    def process_receipt_extraction():
        try:
            ocr = get_ocr()
            all_results = []

            with ThreadPoolExecutor(max_workers=MAX_BATCH_WORKERS) as executor:
                futures = {}
                for pdf_path, original_name in saved_files:
                    future = executor.submit(extract_receipt_data_from_file, pdf_path, original_name, ocr)
                    futures[future] = original_name

                for idx, future in enumerate(as_completed(futures)):
                    try:
                        result = future.result()
                        all_results.append(result)
                    except Exception as e:
                        batch_processing_jobs[job_id]["errors"].append({
                            "file": futures[future],
                            "error": str(e)
                        })

                    batch_processing_jobs[job_id]["processed_files"] += 1
                    batch_processing_jobs[job_id]["progress"] = int(
                        batch_processing_jobs[job_id]["processed_files"] / len(saved_files) * 80
                    )

            batch_processing_jobs[job_id]["results"] = all_results

            if export_format == "excel":
                try:
                    import openpyxl
                    wb = openpyxl.Workbook()
                    ws = wb.active
                    ws.title = "Receipts"
                    ws.append(["Filename", "Merchant", "Date", "Amount"])

                    for r in all_results:
                        ws.append([
                            r["filename"],
                            r["merchant"],
                            r["date"],
                            r["amount"]
                        ])

                    output_path = OUTPUT_DIR / f"receipts_{job_id}.xlsx"
                    wb.save(str(output_path))
                    batch_processing_jobs[job_id]["output_file"] = output_path.name
                    batch_processing_jobs[job_id]["export_format"] = "excel"
                except ImportError:
                    raise HTTPException(status_code=500, detail="openpyxl not installed")
            else:
                output_path = OUTPUT_DIR / f"receipts_{job_id}.json"
                output_path.write_text(json.dumps(all_results, indent=2))
                batch_processing_jobs[job_id]["output_file"] = output_path.name
                batch_processing_jobs[job_id]["export_format"] = "json"

            batch_processing_jobs[job_id]["status"] = "done"
            batch_processing_jobs[job_id]["progress"] = 100
            batch_processing_jobs[job_id]["finished_at"] = time.time()

        except Exception as e:
            batch_processing_jobs[job_id]["status"] = "error"
            batch_processing_jobs[job_id]["error"] = str(e)
        finally:
            for pdf_path, _ in saved_files:
                pdf_path.unlink(missing_ok=True)

    threading.Thread(target=process_receipt_extraction, daemon=True).start()

    return {"job_id": job_id, "status": "processing"}


def extract_receipt_data_from_file(pdf_path: Path, original_name: str, ocr) -> Dict[str, Any]:
    """Extract receipt data from a single file."""
    doc = fitz.open(str(pdf_path))
    data = extract_receipt_data(doc, ocr)
    doc.close()

    return {
        "filename": original_name,
        "merchant": data.get("merchant"),
        "date": data.get("date"),
        "amount": data.get("amount")
    }


@app.get("/api/batch-processing/status/{job_id}")
def get_batch_processing_status(job_id: str):
    """Get status of batch processing job."""
    if job_id not in batch_processing_jobs:
        raise HTTPException(status_code=404, detail="Job not found")

    job = batch_processing_jobs[job_id]

    response = {
        "id": job["id"],
        "status": job["status"],
        "progress": job["progress"],
        "total_files": job.get("total_files", 0),
        "processed_files": job.get("processed_files", 0),
        "started_at": job["started_at"],
        "elapsed_seconds": int(time.time() - job["started_at"])
    }

    if job["status"] == "done":
        response["finished_at"] = job.get("finished_at")
        response["output_file"] = job.get("output_file")
        response["download_url"] = f"/api/download/{job.get('output_file')}"
        response["files_by_type"] = job.get("files_by_type", {})
        response["results"] = job.get("results", [])
        response["export_format"] = job.get("export_format")

    if job["status"] == "error":
        response["error"] = job.get("error")

    response["errors"] = job.get("errors", [])

    return response


if __name__ == "__main__":
    import uvicorn
    print("Starting 4uPDF API on http://localhost:3099")
    uvicorn.run(app, host="0.0.0.0", port=3099)
